# /// script
# requires-python = ">=3.11"
# dependencies = ["pytest>=7.0", "pyyaml>=6.0", "python-pptx>=1.0.0", "pydantic>=2.0", "lxml>=4.9", "Pillow>=10.0"]
# ///
"""S5 misc 非 renderer 修复回归 (fix-ppt-misc-medium-low).

- parse._parse_svg_dimensions: 百分比单位 width/height 不被误抽成 px (无 viewBox → None fallback)
- parse.main: total_files 单次扫描计数 (不重复 is_file)
- render framework: bare framework 列数自适应 point 数, 多 point 不再静默丢内容
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

PLUGIN_ROOT = Path(__file__).resolve().parent.parent
if str(PLUGIN_ROOT) not in sys.path:
    sys.path.insert(0, str(PLUGIN_ROOT))


# ============================================================
# parse._parse_svg_dimensions 百分比守护
# ============================================================

def test_svg_percent_dimensions_return_none(tmp_path):
    """width/height 为百分比且无 viewBox → (None, None), 不被误抽成 (100, 100)。"""
    from engine.parse import _parse_svg_dimensions
    svg = tmp_path / "pct.svg"
    svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="100%" height="100%"></svg>',
        encoding="utf-8",
    )
    assert _parse_svg_dimensions(svg) == (None, None)


def test_svg_viewbox_takes_priority(tmp_path):
    """viewBox 优先解析, 即便同时存在百分比 width。"""
    from engine.parse import _parse_svg_dimensions
    svg = tmp_path / "vb.svg"
    svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 800 600" width="100%"></svg>',
        encoding="utf-8",
    )
    assert _parse_svg_dimensions(svg) == (800, 600)


def test_svg_px_dimensions_parsed(tmp_path):
    """绝对 px 单位仍正常解析。"""
    from engine.parse import _parse_svg_dimensions
    svg = tmp_path / "px.svg"
    svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="320px" height="240px"></svg>',
        encoding="utf-8",
    )
    assert _parse_svg_dimensions(svg) == (320, 240)


# ============================================================
# parse.main total_files 单次扫描
# ============================================================

def test_main_total_files_counts_all_files(tmp_path):
    """total_files 计所有 is_file (含非支持类型), text_sources 仅 md。"""
    from engine.parse import main
    (tmp_path / "a.md").write_text("# Hi\n\n正文内容足够长以供解析使用。\n", encoding="utf-8")
    (tmp_path / "note.txt").write_text("unsupported but counted", encoding="utf-8")
    out = tmp_path / "parsed.json"
    main(str(tmp_path), str(out))
    data = json.loads(out.read_text(encoding="utf-8"))
    assert data["assets"]["total_files"] == 2, data["assets"]
    assert data["assets"]["text_sources"] == 1, data["assets"]


# ============================================================
# render framework 多 point 自适应列数
# ============================================================

def test_framework_renders_all_points(tmp_path):
    """bare framework 3 个 point → 列数自适应 3, 三列内容都渲染 (不丢 point[1:])。"""
    from engine.render import render_presentation, load_theme
    from schemas.slide_plan import SlidePlan
    from pptx import Presentation

    plan_data = {
        "meta": {"title": "t", "preset": "research-report", "theme": "huawei",
                 "generated_at": "2026-06-15T00:00:00"},
        "narrative": {"thesis": "x", "arc": "a", "total_slides": 1},
        "slides": [
            {"id": 1, "role": "content", "visual_type": "framework",
             "content": {
                 "title": "三列框架",
                 "key_points": [
                     {"heading": "A", "body": "第一列内容唯一标识 alphaword"},
                     {"heading": "B", "body": "第二列内容唯一标识 bravoword"},
                     {"heading": "C", "body": "第三列内容唯一标识 charlieword"},
                 ],
             }},
        ],
    }
    plan = SlidePlan.model_validate(plan_data)
    theme = load_theme("huawei")
    out = tmp_path / "fw.pptx"
    render_presentation(plan, theme, str(out))

    prs = Presentation(str(out))
    blob = " ".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "alphaword" in blob and "bravoword" in blob and "charlieword" in blob, (
        f"framework 应渲染全部 3 个 point (列数自适应), 实际文本: {blob}"
    )
