# /// script
# requires-python = ">=3.11"
# dependencies = ["pytest>=7.0", "pyyaml>=6.0", "python-pptx>=1.0.0", "pydantic>=2.0"]
# ///
"""域 A — orchestrator 状态机测试 (S3-P0, test-specs TC-A1..A8 + TC-B1/B2 单测层).

被测对象: engine/orchestrator.py (run / resume / replay 三子命令).

测试手段: subprocess 黑盒调用 orchestrator CLI (与真实使用同路径), 验证:
- exit code 协议 (0/1/2/3, PRD §6.3)
- run 目录落盘结构 (manifest/state/error + 各阶段产物)
- state.json / manifest.json / error.json 字段 (plan §JSON Schema 冻结)
- replay 确定性 (两次跑结构等价)

所有依赖 orchestrator.py / us-military fixture 的用例 skipif 守护:
orchestrator 未就位 → skip; fixture 未就位 → skip。Dev 交付后自然激活。

参考范本:
- tests/integration/test_huawei_full_deck.py:105 _extract_text 文本提取
- tests/integration/test_huawei_full_deck.py:81 idempotency 结构比对范式
"""
from __future__ import annotations

import importlib.util
import json
import re
import shutil
import subprocess
import sys
from pathlib import Path

import pytest

PLUGIN_ROOT = Path(__file__).resolve().parent.parent
ORCHESTRATOR = PLUGIN_ROOT / "engine" / "orchestrator.py"
FIXTURE_DIR = PLUGIN_ROOT / "tests" / "fixtures" / "us-military-rotation"
FIXTURE_INPUT = FIXTURE_DIR / "input"
FIXTURE_RECORDED = FIXTURE_DIR / "recorded"

# run 目录名格式: <YYYYMMDD-HHMMSS> (plan 任务2 GOTCHA: 可带毫秒/pid 后缀防撞)
RUN_DIR_RE = re.compile(r"^\d{8}-\d{6}")

# ---- skipif 守护谓词 (test-specs §命名与位置约定) ----
ORCH_READY = ORCHESTRATOR.exists()
FIXTURE_READY = FIXTURE_INPUT.exists() and FIXTURE_RECORDED.exists()

pytestmark = pytest.mark.skipif(
    not ORCH_READY,
    reason=f"engine/orchestrator.py 尚未就绪 (Dev 任务 1-6 前): {ORCHESTRATOR}",
)


# ============================================================
# 公共辅助
# ============================================================

def _run_orchestrator(*args: str, cwd: Path | None = None) -> subprocess.CompletedProcess:
    """subprocess 调 orchestrator CLI, 返回 CompletedProcess (含 returncode/stdout/stderr).

    用 uv run --script 走 PEP 723 (与真实使用同路径)。捕获文本输出, 不抛异常。
    """
    cmd = [
        "uv", "run", "--script", str(ORCHESTRATOR), *args,
    ]
    return subprocess.run(
        cmd,
        cwd=str(cwd or PLUGIN_ROOT),
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=300,
    )


def _latest_run_dir(output_dir: Path) -> Path:
    """定位 output_dir/.ppt-workdir/runs/ 下最新的 run 目录。"""
    runs = output_dir / ".ppt-workdir" / "runs"
    assert runs.is_dir(), f"runs 目录未创建: {runs}"
    candidates = sorted(p for p in runs.iterdir() if p.is_dir())
    assert candidates, f"runs 目录下无 run 子目录: {runs}"
    return candidates[-1]


def test_plugin_version_tolerates_packaged_metadata_absence(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """CodeAgent 打包产物不含 .claude-plugin/plugin.json 时不应启动失败。"""
    spec = importlib.util.spec_from_file_location("ppt_orchestrator_under_test", ORCHESTRATOR)
    assert spec is not None
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)

    monkeypatch.setattr(module, "PLUGIN_ROOT", tmp_path)

    assert module._plugin_version() is None


def _read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def _extract_pptx_texts(pptx_path: Path) -> list[list[str]]:
    """读 PPTX, 返回每页的 shape 文本列表 (range范本: test_huawei_full_deck.py:105)。"""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    pages: list[list[str]] = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append("\n".join(p.text for p in shape.text_frame.paragraphs))
            else:
                texts.append("")
        pages.append(texts)
    return pages


def _slide_plan_structure(yaml_path: Path) -> dict:
    """从 slide-plan.yaml 提取结构签名 (排除非确定字段如 meta.generated_at)。

    key_points 在真实 schema 中位于 slide.content 下 (SlideContent.key_points),
    不是 slide 顶层 (schemas/slide_plan.py:86-98 verify)。
    """
    import yaml

    data = yaml.safe_load(yaml_path.read_text(encoding="utf-8"))
    slides = data.get("slides", []) if isinstance(data, dict) else []

    def _point_count(s: dict) -> int:
        content = s.get("content", {}) or {}
        return len(content.get("key_points", []) or [])

    return {
        "count": len(slides),
        "visual_types": [s.get("visual_type") for s in slides],
        "roles": [s.get("role") for s in slides],
        "point_counts": [_point_count(s) for s in slides],
    }


# ============================================================
# TC-A7: 参数 / 用法错误 → exit 2 (不依赖 fixture, 可独立跑)
# ============================================================

class TestUsageErrors:
    """TC-A7: 缺子命令 / 未知子命令 / run-dir 不存在 → exit 2。"""

    def test_no_subcommand_exits_2(self):
        result = _run_orchestrator()
        assert result.returncode == 2, (
            f"无子命令应 exit 2, 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )

    def test_unknown_subcommand_exits_2(self):
        result = _run_orchestrator("foobar")
        assert result.returncode == 2, (
            f"未知子命令应 exit 2, 实际 {result.returncode}"
        )

    def test_resume_nonexistent_run_dir_exits_2(self, tmp_path):
        result = _run_orchestrator("resume", str(tmp_path / "does-not-exist"))
        assert result.returncode == 2, (
            f"resume 不存在的 run-dir 应 exit 2, 实际 {result.returncode}\n"
            f"stderr={result.stderr}"
        )


# ============================================================
# 以下用例依赖 us-military fixture (run/resume/replay 需要真实输入)
# ============================================================

_fixture_skip = pytest.mark.skipif(
    not FIXTURE_READY,
    reason=f"us-military fixture 尚未就绪 (Dev 任务 9 前): {FIXTURE_DIR}",
)


# ============================================================
# TC-A1 / A2: run 新建 run 目录, 停在 architect gate (exit 3), manifest 字段
# ============================================================

@_fixture_skip
class TestRunSubcommand:
    """TC-A1: run → exit 3 + run 目录结构; TC-A2: manifest 字段完整。"""

    @pytest.fixture
    def run_result(self, tmp_path):
        out = tmp_path / "deck.pptx"
        result = _run_orchestrator(
            "run", str(FIXTURE_INPUT),
            "--theme", "huawei",
            "--output", str(out),
        )
        return result, tmp_path

    def test_run_exits_3_need_llm(self, run_result):
        """TC-A1: run 到 architect gate → exit 3 (NEED_LLM)。"""
        result, _ = run_result
        assert result.returncode == 3, (
            f"run 应 exit 3 (NEED_LLM), 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )

    def test_run_creates_run_dir(self, run_result):
        """TC-A1: .ppt-workdir/runs/<YYYYMMDD-HHMMSS>/ 被创建。"""
        _, out_dir = run_result
        run_dir = _latest_run_dir(out_dir)
        assert RUN_DIR_RE.match(run_dir.name), (
            f"run 目录名 {run_dir.name!r} 不匹配 <YYYYMMDD-HHMMSS>"
        )

    def test_run_writes_parse_and_architect_prompt(self, run_result):
        """TC-A1: 01-parse/parsed-content.json + 02-architect/prompt-context.md 落盘。"""
        _, out_dir = run_result
        run_dir = _latest_run_dir(out_dir)
        assert (run_dir / "01-parse" / "parsed-content.json").exists()
        assert (run_dir / "02-architect" / "prompt-context.md").exists()

    def test_run_state_json_awaiting_architect_output(self, run_result):
        """TC-A1: state.json 指向 architect gate, mode=bridge。"""
        _, out_dir = run_result
        run_dir = _latest_run_dir(out_dir)
        state = _read_json(run_dir / "state.json")
        assert state["current_stage"] == "architect"
        assert state["awaiting"] == "02-architect/output.yaml"
        assert state["mode"] == "bridge"

    def test_run_no_architect_output_yet(self, run_result):
        """TC-A1: 02-architect/output.yaml 不应存在 (等会话 LLM 写)。"""
        _, out_dir = run_result
        run_dir = _latest_run_dir(out_dir)
        assert not (run_dir / "02-architect" / "output.yaml").exists()

    def test_manifest_fields_complete(self, run_result):
        """TC-A2: manifest.json 顶层平铺字段完整 (plan §JSON Schema 冻结)。"""
        _, out_dir = run_result
        run_dir = _latest_run_dir(out_dir)
        manifest = _read_json(run_dir / "manifest.json")
        # 顶层键名冻结
        assert "input" in manifest
        assert re.fullmatch(r"[0-9a-f]{64}", manifest["input_sha256"]), (
            f"input_sha256 应为 64 hex 小写, 实际 {manifest.get('input_sha256')!r}"
        )
        assert manifest["input_kind"] in ("dir", "file")
        assert re.fullmatch(r"\d+\.\d+\.\d+", manifest["plugin_version"]), (
            f"plugin_version 应为 semver 格式, 实际 {manifest.get('plugin_version')!r}"
        )
        assert manifest["taste_requested"] is False
        # stages: parse 已 completed
        assert manifest["stages"]["parse"]["status"] == "completed"


# ============================================================
# TC-A3 / A4 / A5 / A6: resume 校验与推进
# ============================================================

@_fixture_skip
class TestResumeArchitectGate:
    """TC-A3: resume 通过推进到 plan gate; TC-A4/A5: 协议层 fail-fast。

    resume architect 阶段两层语义 (team-lead 2026-06-11 裁决):
    1. 协议层硬门: output.yaml 缺失/空/非法 YAML → error.json (stage=architect,
       rule=protocol.*) + exit 1。
    2. validate_architecture 内容质量 issues 降级 advisory: 写 validation.json +
       不 fail, 继续 exit 3 推进 (依据: v3.4.0 SKILL.md Stage 2 从不调该校验,
       硬 fail 构成行为回归; 与 replay 路径语义对称)。
    """

    @pytest.fixture
    def run_dir_at_architect(self, tmp_path):
        """跑 run 停在 architect gate, 返回 run 目录。"""
        out = tmp_path / "deck.pptx"
        _run_orchestrator(
            "run", str(FIXTURE_INPUT), "--theme", "huawei", "--output", str(out),
        )
        return _latest_run_dir(tmp_path)

    def test_resume_with_valid_architect_advances_to_plan(self, run_dir_at_architect):
        """TC-A3: 写入合法 architect output → resume → exit 3 + 推进到 plan。"""
        run_dir = run_dir_at_architect
        # 模拟会话 LLM 产出: 拷 fixture recorded/architect-output.yaml
        src = FIXTURE_RECORDED / "architect-output.yaml"
        shutil.copyfile(src, run_dir / "02-architect" / "output.yaml")

        result = _run_orchestrator("resume", str(run_dir))
        assert result.returncode == 3, (
            f"resume (architect 通过) 应 exit 3 推进到 plan gate, 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )
        assert (run_dir / "02-architect" / "validation.json").exists()
        state = _read_json(run_dir / "state.json")
        assert state["current_stage"] == "plan"
        assert state["awaiting"] == "03-plan/output.yaml"
        assert (run_dir / "03-plan" / "prompt-context.md").exists()
        assert (run_dir / "03-plan" / "theme-prompt.md").exists()

    def test_resume_architect_content_issues_advisory_not_fail(self, run_dir_at_architect):
        """TC-A3b: validate_architecture 内容 issues → advisory 推进 (exit 3, 不 fail)。

        裁决第 2 层: 内容质量 issues (如 chapter title < 10 字, architect.py:75)
        写 validation.json 但不 fail, 继续 exit 3 推进到 plan。
        构造: 在 fixture architecture 基础上把某 chapter title 改短 (< 10 字) 触发 issue,
        但 schema 仍合法 (协议层通过)。
        """
        import yaml as _yaml

        run_dir = run_dir_at_architect
        src = FIXTURE_RECORDED / "architect-output.yaml"
        arch = _yaml.safe_load(src.read_text(encoding="utf-8"))
        # 把首个 chapter title 改短触发 action_title issue (但结构合法)
        chapters = arch.get("chapters") or []
        if chapters:
            chapters[0]["title"] = "短"  # < 10 字 → validate_architecture 报 issue
        (run_dir / "02-architect" / "output.yaml").write_text(
            _yaml.safe_dump(arch, allow_unicode=True), encoding="utf-8",
        )

        result = _run_orchestrator("resume", str(run_dir))
        # advisory: 不 fail, 继续推进到 plan
        assert result.returncode == 3, (
            f"内容 issues 应 advisory 推进 (exit 3), 不 fail, 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )
        # validation.json 含 issues (advisory 记录)
        vjson = run_dir / "02-architect" / "validation.json"
        assert vjson.exists(), "validation.json 应记录 advisory issues"
        # 未写 error.json (协议层未失败) — advisory 不触发 error
        state = _read_json(run_dir / "state.json")
        assert state["current_stage"] == "plan", "advisory 后应推进到 plan"

    def test_resume_missing_output_fails(self, run_dir_at_architect):
        """TC-A5: output.yaml 不存在 → 协议层 exit 1 + error.json 指明期望路径。"""
        run_dir = run_dir_at_architect
        # 不写 output.yaml
        result = _run_orchestrator("resume", str(run_dir))
        assert result.returncode == 1, (
            f"resume 缺 output.yaml 应 exit 1 (协议层硬门), 实际 {result.returncode}"
        )
        err = _read_json(run_dir / "error.json")
        assert err["stage"] == "architect"
        # 协议层失败 rule=protocol.* (裁决第 1 层)
        assert "protocol" in (err.get("rule") or ""), (
            f"协议层失败 rule 应含 protocol, 实际 {err.get('rule')!r}"
        )
        assert "02-architect/output.yaml" in err["message"], (
            f"error.json.message 应指明期望路径, 实际 {err['message']!r}"
        )

    def test_resume_empty_output_fails(self, run_dir_at_architect):
        """TC-A4: 空 output.yaml → 协议层 exit 1 + error.json.stage=architect。"""
        run_dir = run_dir_at_architect
        (run_dir / "02-architect" / "output.yaml").write_text("", encoding="utf-8")
        result = _run_orchestrator("resume", str(run_dir))
        assert result.returncode == 1
        err = _read_json(run_dir / "error.json")
        assert err["stage"] == "architect"
        assert "protocol" in (err.get("rule") or ""), (
            f"协议层失败 rule 应含 protocol, 实际 {err.get('rule')!r}"
        )

    def test_resume_invalid_yaml_fails(self, run_dir_at_architect):
        """TC-A4: 非法 YAML → 协议层 exit 1 + error.json.stage=architect。"""
        run_dir = run_dir_at_architect
        (run_dir / "02-architect" / "output.yaml").write_text(
            "key: [unclosed\n  bad: :::", encoding="utf-8",
        )
        result = _run_orchestrator("resume", str(run_dir))
        assert result.returncode == 1
        err = _read_json(run_dir / "error.json")
        assert err["stage"] == "architect"
        assert "protocol" in (err.get("rule") or ""), (
            f"协议层失败 rule 应含 protocol, 实际 {err.get('rule')!r}"
        )


@_fixture_skip
class TestResumePlanContentVolume:
    """TC-A6: plan 内容量失败 → 顶层 exit 1 + error.json 记 rule (不透传 validate_plan 内部 1/2)。"""

    @pytest.fixture
    def run_dir_at_plan(self, tmp_path):
        """跑到 plan gate, 返回 run 目录。"""
        out = tmp_path / "deck.pptx"
        _run_orchestrator(
            "run", str(FIXTURE_INPUT), "--theme", "huawei", "--output", str(out),
        )
        run_dir = _latest_run_dir(tmp_path)
        shutil.copyfile(
            FIXTURE_RECORDED / "architect-output.yaml",
            run_dir / "02-architect" / "output.yaml",
        )
        _run_orchestrator("resume", str(run_dir))  # 推进到 plan
        return run_dir

    def test_plan_content_volume_failure_top_exit_1(self, run_dir_at_plan):
        """TC-A6: 内容量不达标的 plan → orchestrator 顶层 exit 1 (非 validate_plan 内部 1/2 透传)。

        bad_plan 已用真实 schema 自验 (2026-06-11):
        - 通过 SlidePlan pydantic 校验 (走真实字段名 id/role/content.key_points + Narrative dict)
        - validate_plan 判内容量 FAIL ("内容区 6 字 (min 200)"), ok=False
        bullets 类型走整页总量校验 (TOTAL_MIN=200), key_points 总字数远不足 → 必 FAIL。
        """
        run_dir = run_dir_at_plan
        # 内容量不达标的 plan: bullets + key_points 总字数 6 << 200 (整页总量 FAIL)
        bad_plan = (
            "meta:\n"
            "  title: \"测试 deck\"\n"
            "  preset: research-report\n"
            "  theme: huawei\n"
            "  generated_at: \"2026-04-24T23:46:00\"\n"
            "narrative:\n"
            "  thesis: \"测试论点\"\n"
            "  arc: \"opening -> closing\"\n"
            "  total_slides: 1\n"
            "slides:\n"
            "  - id: 1\n"
            "    role: content\n"
            "    visual_type: bullets\n"
            "    content:\n"
            "      title: \"行动标题至少十个字测试用\"\n"
            "      key_points:\n"
            "        - \"短一\"\n"
            "        - \"短二\"\n"
            "        - \"短三\"\n"
        )
        (run_dir / "03-plan" / "output.yaml").write_text(bad_plan, encoding="utf-8")
        result = _run_orchestrator("resume", str(run_dir))
        # 顶层收敛为 1 (即便 validate_plan 内部内容量是 exit 1 / 应用率是 exit 2)
        assert result.returncode == 1, (
            f"plan 内容量失败应顶层 exit 1, 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )
        err = _read_json(run_dir / "error.json")
        assert err["stage"] == "plan", f"error.json.stage 应为 plan, 实际 {err.get('stage')!r}"
        # rule 应指向内容量类规则 (不空)
        assert err.get("rule"), f"error.json.rule 应非空 (定位具体规则), 实际 {err.get('rule')!r}"


# ============================================================
# TC-A8: --with-taste 在 run 阶段 (S3-P1 档位 B 接入后: 记 flag + 推迟到 render 后)
# ============================================================

@_fixture_skip
class TestWithTasteRun:
    """--with-taste 在 run 阶段: 不报错 + manifest.taste_requested=true + 仍停在 architect gate。

    S3-P1 档位 B 接入后, taste 是 render 完成后的可选末阶段; run 阶段只记 flag,
    流程不变 (仍停 architect gate, exit 3), taste gate 要等 plan resume render 完才触发。
    (替换 P0 期 'stdout 含 S3-P1 占位提示' 旧断言 — 占位已移除。)
    """

    @pytest.fixture
    def run_result(self, tmp_path):
        out = tmp_path / "deck.pptx"
        result = _run_orchestrator(
            "run", str(FIXTURE_INPUT),
            "--theme", "huawei", "--output", str(out), "--with-taste",
        )
        return result, tmp_path

    def test_with_taste_does_not_error_on_arg(self, run_result):
        """--with-taste 不应导致 exit 2 (参数被接受)。"""
        result, _ = run_result
        assert result.returncode != 2, (
            f"--with-taste 不应触发用法错误 exit 2, 实际 {result.returncode}\n"
            f"stderr={result.stderr}"
        )

    def test_with_taste_run_still_pauses_at_architect(self, run_result):
        """run --with-taste 仍停在 architect gate (exit 3); taste 推迟到 render 后, 不在 run 阶段触发。"""
        result, out_dir = run_result
        assert result.returncode == 3, (
            f"run --with-taste 应仍 exit 3 (停 architect gate), 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )
        # 旧 S3-P1 占位提示已移除
        combined = result.stdout + result.stderr
        assert "占位" not in combined, "S3-P1 占位提示应已移除 (taste 已真接入)"
        # taste 产物此刻不应存在 (要等 render 后)
        run_dir = _latest_run_dir(out_dir)
        assert not (run_dir / "05-taste" / "taste-report.json").exists()
        state = _read_json(run_dir / "state.json")
        assert state["current_stage"] == "architect"

    def test_with_taste_manifest_flag_true(self, run_result):
        """manifest.taste_requested == true。"""
        _, out_dir = run_result
        run_dir = _latest_run_dir(out_dir)
        manifest = _read_json(run_dir / "manifest.json")
        assert manifest["taste_requested"] is True


# ============================================================
# S3-P1 档位 B: taste 可选末阶段接入 (bridge taste gate + replay taste)
# ============================================================

# 真实 deck taste 基线 fixture (v3.8.0 交付, 合法 TasteReport), 复用作 taste 产出样本
TASTE_FIXTURE_JSON = PLUGIN_ROOT / "tests" / "fixtures" / "taste" / "renderer-deck.taste-report.json"


@_fixture_skip
@pytest.mark.skipif(
    not TASTE_FIXTURE_JSON.exists(),
    reason=f"taste baseline fixture 未就位: {TASTE_FIXTURE_JSON}",
)
class TestTasteGateBridge:
    """bridge mode: run/resume 流程 render 完成后转 taste gate (exit 3), resume 校验 JSON 契约。

    LLM 视觉评分留会话层 (orchestrator 零 LLM): orchestrator 只在 render 后暂停 NEED_LLM,
    会话产出 05-taste/taste-report.json 后 resume 校验契约 -> done。
    """

    @pytest.fixture
    def run_dir_at_taste_gate(self, tmp_path):
        """跑完整 bridge 流程到 taste gate (render 已完成, exit 3 等 taste JSON), 返回 (run_dir, resume_result)。"""
        out = tmp_path / "deck.pptx"
        _run_orchestrator(
            "run", str(FIXTURE_INPUT), "--theme", "huawei", "--output", str(out), "--with-taste",
        )
        run_dir = _latest_run_dir(tmp_path)
        # architect gate
        shutil.copyfile(
            FIXTURE_RECORDED / "architect-output.yaml",
            run_dir / "02-architect" / "output.yaml",
        )
        _run_orchestrator("resume", str(run_dir))  # -> plan gate
        # plan gate -> render -> taste gate
        shutil.copyfile(
            FIXTURE_RECORDED / "plan-output.yaml",
            run_dir / "03-plan" / "output.yaml",
        )
        result = _run_orchestrator("resume", str(run_dir))
        return run_dir, result

    def test_render_done_pauses_at_taste_gate(self, run_dir_at_taste_gate):
        """render 完成后转 taste gate: exit 3 + state=taste + prompt-context + deck 已渲染。"""
        run_dir, result = run_dir_at_taste_gate
        assert result.returncode == 3, (
            f"render 完成 + taste_requested 应转 taste gate (exit 3), 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )
        state = _read_json(run_dir / "state.json")
        assert state["current_stage"] == "taste"
        assert state["awaiting"] == "05-taste/taste-report.json"
        assert (run_dir / "05-taste" / "prompt-context.md").exists()
        # deck 已渲染 (taste 在 render 之后)
        assert (run_dir / "04-render" / "deck.pptx").exists()

    def test_taste_resume_valid_json_done(self, run_dir_at_taste_gate):
        """taste gate 写合法 taste-report.json -> resume -> exit 0 + state=done + manifest.taste_result。"""
        run_dir, _ = run_dir_at_taste_gate
        shutil.copyfile(TASTE_FIXTURE_JSON, run_dir / "05-taste" / "taste-report.json")
        result = _run_orchestrator("resume", str(run_dir))
        assert result.returncode == 0, (
            f"合法 taste JSON 应 exit 0 (done), 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )
        state = _read_json(run_dir / "state.json")
        assert state["current_stage"] == "done"
        manifest = _read_json(run_dir / "manifest.json")
        assert "taste_result" in manifest, "manifest 应记 taste_result"
        tr = manifest["taste_result"]
        assert tr["report_path"] == "05-taste/taste-report.json"
        assert isinstance(tr["layout_avg"], (int, float))
        assert isinstance(tr["palette_avg"], (int, float))
        assert manifest["stages"]["taste"]["status"] == "completed"

    def test_taste_resume_missing_json_fails(self, run_dir_at_taste_gate):
        """taste gate 未产出 JSON -> resume 协议层 exit 1 + error.json(stage=taste, rule=protocol.*)。"""
        run_dir, _ = run_dir_at_taste_gate
        # 不写 taste-report.json
        result = _run_orchestrator("resume", str(run_dir))
        assert result.returncode == 1, (
            f"缺 taste JSON 应 exit 1 (协议层), 实际 {result.returncode}"
        )
        err = _read_json(run_dir / "error.json")
        assert err["stage"] == "taste"
        assert "protocol" in (err.get("rule") or ""), (
            f"协议层失败 rule 应含 protocol, 实际 {err.get('rule')!r}"
        )

    def test_taste_resume_invalid_schema_fails(self, run_dir_at_taste_gate):
        """taste gate 写不合契约的 JSON (avg 与 slides 不自洽) -> exit 1 + error.json(rule=schema.taste_report)。"""
        run_dir, _ = run_dir_at_taste_gate
        # 取合法 fixture, 破坏内部一致性 (layout_avg 改成明显错的值) -> 触发 model_validator
        data = json.loads(TASTE_FIXTURE_JSON.read_text(encoding="utf-8"))
        data["summary"]["layout_avg"] = 1.01  # 与 slides 实际平均不符 (容差 0.01 外)
        (run_dir / "05-taste" / "taste-report.json").write_text(
            json.dumps(data, ensure_ascii=False), encoding="utf-8",
        )
        result = _run_orchestrator("resume", str(run_dir))
        assert result.returncode == 1, (
            f"不合契约 taste JSON 应 exit 1 (schema 层), 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )
        err = _read_json(run_dir / "error.json")
        assert err["stage"] == "taste"
        assert err.get("rule") == "schema.taste_report", (
            f"schema 层失败 rule 应为 schema.taste_report, 实际 {err.get('rule')!r}"
        )


@_fixture_skip
@pytest.mark.skipif(
    not TASTE_FIXTURE_JSON.exists(),
    reason=f"taste baseline fixture 未就位: {TASTE_FIXTURE_JSON}",
)
class TestReplayTaste:
    """replay mode: --with-taste 确定性回放 (有 fixture 校验 / 无 fixture 优雅跳过)。"""

    def test_replay_with_taste_fixture_validates(self, tmp_path):
        """fixture 含 recorded/taste-report.json -> replay --with-taste 拷入 + 校验 + 落 manifest, exit 0。"""
        # copytree us-military fixture 到临时目录, 补一份 recorded/taste-report.json
        fx = tmp_path / "fx"
        shutil.copytree(FIXTURE_DIR, fx)
        shutil.copyfile(TASTE_FIXTURE_JSON, fx / "recorded" / "taste-report.json")
        out = tmp_path / "deck.pptx"
        result = _run_orchestrator(
            "replay", str(fx / "input"), "--fixture", str(fx),
            "--theme", "huawei", "--output", str(out), "--with-taste",
        )
        assert result.returncode == 0, (
            f"replay --with-taste (有 fixture) 应 exit 0, 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )
        run_dir = _latest_run_dir(tmp_path)
        assert (run_dir / "05-taste" / "taste-report.json").exists(), "05-taste 应有回放的 taste JSON"
        manifest = _read_json(run_dir / "manifest.json")
        assert "taste_result" in manifest, "manifest 应记 taste_result"
        assert manifest["stages"]["taste"]["status"] == "completed"

    def test_replay_with_taste_no_fixture_skips_gracefully(self, tmp_path):
        """fixture 无 recorded/taste-report.json -> replay --with-taste 优雅跳过 (exit 0, 无 taste 产物)。"""
        out = tmp_path / "deck.pptx"
        # 真实 us-military fixture 不含 recorded/taste-report.json
        result = _run_orchestrator(
            "replay", str(FIXTURE_INPUT), "--fixture", str(FIXTURE_DIR),
            "--theme", "huawei", "--output", str(out), "--with-taste",
        )
        assert result.returncode == 0, (
            f"replay --with-taste (无 fixture) 应优雅跳过 exit 0, 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )
        run_dir = _latest_run_dir(tmp_path)
        assert not (run_dir / "05-taste" / "taste-report.json").exists(), (
            "无 fixture 时不应产出 taste JSON (优雅跳过)"
        )
        manifest = _read_json(run_dir / "manifest.json")
        assert "taste_result" not in manifest, "优雅跳过时 manifest 不应有 taste_result"


# ============================================================
# TC-B1 / B2: replay 确定性 (两次跑结构等价) — 单测层 (E2E 也覆盖, 这里聚焦确定性)
# ============================================================

@_fixture_skip
class TestReplayDeterminism:
    """TC-B1: 两次 replay slide-plan 结构等价; TC-B2: PPTX 结构等价。"""

    @pytest.fixture
    def two_runs(self, tmp_path):
        """连续两次独立 replay, 返回 (run_dir_a, pptx_a, run_dir_b, pptx_b)。"""
        out_a = tmp_path / "a" / "deck.pptx"
        out_b = tmp_path / "b" / "deck.pptx"
        out_a.parent.mkdir(parents=True)
        out_b.parent.mkdir(parents=True)
        r_a = _run_orchestrator(
            "replay", str(FIXTURE_INPUT), "--fixture", str(FIXTURE_DIR),
            "--theme", "huawei", "--output", str(out_a),
        )
        r_b = _run_orchestrator(
            "replay", str(FIXTURE_INPUT), "--fixture", str(FIXTURE_DIR),
            "--theme", "huawei", "--output", str(out_b),
        )
        assert r_a.returncode == 0, f"replay A 失败: {r_a.stdout}\n{r_a.stderr}"
        assert r_b.returncode == 0, f"replay B 失败: {r_b.stdout}\n{r_b.stderr}"
        return (
            _latest_run_dir(tmp_path / "a"), out_a,
            _latest_run_dir(tmp_path / "b"), out_b,
        )

    def test_slide_plan_structure_equivalent(self, two_runs):
        """TC-B1: 两次 03-plan/output.yaml 结构签名相等 (排除时间戳)。"""
        run_a, _, run_b, _ = two_runs
        struct_a = _slide_plan_structure(run_a / "03-plan" / "output.yaml")
        struct_b = _slide_plan_structure(run_b / "03-plan" / "output.yaml")
        assert struct_a == struct_b, (
            f"两次 replay slide-plan 结构不等价 (非确定源!):\nA={struct_a}\nB={struct_b}"
        )

    def test_pptx_structure_equivalent(self, two_runs):
        """TC-B2: 两次 PPTX 页数/shape 数/文本集合相等。"""
        _, pptx_a, _, pptx_b = two_runs
        pages_a = _extract_pptx_texts(pptx_a)
        pages_b = _extract_pptx_texts(pptx_b)
        assert len(pages_a) == len(pages_b), "两次 PPTX 页数不一致 (非确定性)"
        for idx, (pa, pb) in enumerate(zip(pages_a, pages_b), start=1):
            assert len(pa) == len(pb), f"Slide {idx} shape 数不一致"
            assert pa == pb, f"Slide {idx} 文本不一致:\nA={pa}\nB={pb}"


# ============================================================
# code-review 遗留修复回归 (M1 / M2 / L2)
# ============================================================

@_fixture_skip
class TestErrorJsonArchivedOnSuccess:
    """M1: gate 失败写 error.json, 修复后成功推进时归档为 error-resolved-N.json。

    避免残留 error.json 与 state.json(done/推进) 矛盾误导人工排查。
    """

    def test_error_archived_after_recovery(self, tmp_path):
        out = tmp_path / "deck.pptx"
        _run_orchestrator(
            "run", str(FIXTURE_INPUT), "--theme", "huawei", "--output", str(out),
        )
        run_dir = _latest_run_dir(tmp_path)
        # 第一次 resume 缺 output → 协议层失败写 error.json
        r1 = _run_orchestrator("resume", str(run_dir))
        assert r1.returncode == 1, f"缺 output 应 exit 1, 实际 {r1.returncode}"
        assert (run_dir / "error.json").exists(), "失败应写 error.json"
        # 修复: 写合法 architect output → resume 成功推进, error.json 归档
        shutil.copyfile(
            FIXTURE_RECORDED / "architect-output.yaml",
            run_dir / "02-architect" / "output.yaml",
        )
        r2 = _run_orchestrator("resume", str(run_dir))
        assert r2.returncode == 3, (
            f"修复后 resume 应 exit 3 推进, 实际 {r2.returncode}\n"
            f"stdout={r2.stdout}\nstderr={r2.stderr}"
        )
        assert not (run_dir / "error.json").exists(), "成功推进后 error.json 应被归档"
        assert (run_dir / "error-resolved-1.json").exists(), (
            "error.json 应改名归档为 error-resolved-1.json (保留失败-恢复追溯)"
        )


@_fixture_skip
class TestReplayFixtureInvalidYaml:
    """L2: replay 路径 fixture architect-output.yaml 非法 YAML → exit 1 + error.json。

    与 resume 协议层 invalid_yaml 对称, 防非法 fixture 直接崩进程 (无 try/except 时
    yaml.safe_load 抛 ScannerError 未捕获 → traceback)。
    """

    def test_replay_invalid_architect_fixture_exits_1(self, tmp_path):
        # 拷 fixture 到临时目录, 破坏 recorded/architect-output.yaml
        fx = tmp_path / "fx"
        shutil.copytree(FIXTURE_DIR, fx)
        (fx / "recorded" / "architect-output.yaml").write_text(
            "key: [unclosed\n  bad: :::", encoding="utf-8",
        )
        out = tmp_path / "deck.pptx"
        result = _run_orchestrator(
            "replay", str(fx / "input"), "--fixture", str(fx),
            "--theme", "huawei", "--output", str(out),
        )
        assert result.returncode == 1, (
            f"非法 fixture architect yaml 应 exit 1 (不崩 traceback), 实际 {result.returncode}\n"
            f"stdout={result.stdout}\nstderr={result.stderr}"
        )
        run_dir = _latest_run_dir(tmp_path)
        err = _read_json(run_dir / "error.json")
        assert err["stage"] == "architect"
        assert "fixture" in (err.get("rule") or ""), (
            f"rule 应标识 fixture 来源, 实际 {err.get('rule')!r}"
        )


class TestSafeCopyOutputLockFallback:
    """M2: 最终输出 copyfile 被锁 (PermissionError) → 自动落 _v2 (Windows 防锁规范)。

    import 层单测 (subprocess 难可靠模拟文件锁): monkeypatch shutil.copyfile 让首次
    落 deck.pptx 抛 PermissionError, 验证 fallback 到 deck_v2.pptx 且内容完整。
    """

    def test_locked_output_falls_back_to_v2(self, tmp_path, monkeypatch):
        from engine import orchestrator as orch

        src = tmp_path / "src.pptx"
        src.write_bytes(b"deck-bytes-payload")
        dst = tmp_path / "out" / "deck.pptx"

        real_copy = shutil.copyfile

        def fake_copy(s, d):
            # 模拟 deck.pptx 被 Office 锁定: 首选落点抛 PermissionError, _vN 落点放行
            if Path(d).name == "deck.pptx":
                raise PermissionError("locked by Office")
            return real_copy(s, d)

        monkeypatch.setattr(orch.shutil, "copyfile", fake_copy)
        actual = orch._safe_copy_output(src, dst)
        assert actual.name == "deck_v2.pptx", f"应 fallback 到 deck_v2.pptx, 实际 {actual.name}"
        assert actual.exists()
        assert actual.read_bytes() == b"deck-bytes-payload", "fallback 落点内容应完整"
