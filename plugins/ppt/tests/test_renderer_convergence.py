# /// script
# requires-python = ">=3.11"
# dependencies = ["pytest>=7.0", "pyyaml>=6.0", "python-pptx>=1.0.0", "pydantic>=2.0", "lxml>=4.9", "Pillow>=10.0"]
# ///
"""自适应母版收敛回归 (S4 收敛重构 T16, 对应 AC-004).

验证 Phase 3 收敛主体:
- 无后缀 cards/process/comparison 自适应母版: 布局由 len(points) 派生, 全 point 不丢。
- 老后缀 deprecated alias (cards-2..5 等) 向后兼容: 两条路径都不丢 point。
  - path Y: content.key_points (legacy 兜底, 现可达)
  - path X: slide.variant.points (CardsContent, 等 T8 让 get_points 读 variant.points)

设计张力 (已与 team-lead 对账, AC-004 口径采两路全覆盖):
T7 给无后缀 + 老后缀 cards-2..5 建了 typed variant CardsContent(points 字段),
但 T8 前老 renderer 的 _render_cards fallback 只认 variant.cards/content.cards, 不认 variant.points
→ path X 会静默丢 point。故 path X + 无后缀渲染断言用 ADAPTIVE_READY 守卫, T8 注册 renderer
+ get_points 读 variant.points 后自动激活。

render->extract 模式沿用 tests/test_misc_s5_fixes.py:79 (唯一标识 token 拼 blob 断言)。
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

PLUGIN_ROOT = Path(__file__).resolve().parent.parent
if str(PLUGIN_ROOT) not in sys.path:
    sys.path.insert(0, str(PLUGIN_ROOT))


def _adaptive_ready() -> bool:
    """T8 是否已注册无后缀 cards/process/comparison renderer。"""
    try:
        import engine.render  # noqa: F401  触发全部注册
        from engine.renderer_kit import _RENDERERS
        return all(vt in _RENDERERS for vt in ("cards", "process", "comparison"))
    except Exception:
        return False


ADAPTIVE_READY = _adaptive_ready()
adaptive_only = pytest.mark.skipif(
    not ADAPTIVE_READY,
    reason="无后缀 cards/process/comparison renderer 未注册 (等 T8); path X + 派生布局断言 blocked",
)


# ============================================================
# 共享 render->extract 工具
# ============================================================

def _render_blob(plan_data: dict, out_path: Path) -> str:
    """渲染 plan 并返回所有 shape 文本拼接 blob。"""
    from pptx import Presentation

    from engine.render import load_theme, render_presentation
    from schemas.slide_plan import SlidePlan
    plan = SlidePlan.model_validate(plan_data)
    theme = load_theme("huawei")
    render_presentation(plan, theme, str(out_path))
    prs = Presentation(str(out_path))
    return " ".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )


def _plan_path_y(visual_type: str, n: int) -> dict:
    """老后缀走 path Y (content.key_points)。每 point body 含唯一 token。"""
    pts = [{"heading": f"H{i}", "body": f"ptY-{visual_type}-{i}-tok"} for i in range(n)]
    return {
        "meta": {"title": "t", "preset": "research-report", "theme": "huawei",
                 "generated_at": "2026-06-18T00:00:00"},
        "narrative": {"thesis": "x", "arc": "a", "total_slides": 1},
        "slides": [{"id": 1, "role": "content", "visual_type": visual_type,
                    "content": {"title": f"{visual_type} 标题", "key_points": pts}}],
    }


def _plan_path_x_cards(visual_type: str, n: int) -> dict:
    """无后缀/老后缀 cards 走 path X (slide.variant=CardsContent.points)。"""
    pts = [{"heading": f"H{i}", "body": f"ptX-{visual_type}-{i}-tok"} for i in range(n)]
    return {
        "meta": {"title": "t", "preset": "research-report", "theme": "huawei",
                 "generated_at": "2026-06-18T00:00:00"},
        "narrative": {"thesis": "x", "arc": "a", "total_slides": 1},
        "slides": [{"id": 1, "role": "content", "visual_type": visual_type,
                    "content": {"title": f"{visual_type} path-X"},
                    "variant": {"visual_type": visual_type, "title": f"{visual_type} path-X",
                                "points": pts}}],
    }


def _tokens_present(blob: str, prefix: str, n: int) -> list[str]:
    """返回缺失的 token 列表 (空 = 全部渲染)。"""
    return [f"{prefix}-{i}-tok" for i in range(n) if f"{prefix}-{i}-tok" not in blob]


def _card_grid(plan_data: dict, out_path: Path, n: int) -> tuple[int, int]:
    """渲染后返回卡片网格 (行数, 列数)。

    卡片标题 shape 文本为 H0..H{n-1} (见 _plan_path_x_cards 构造)。distinct top
    坐标数 = 行数, distinct left = 列数; EMU 取整到 0.01 inch 去噪。
    """
    from pptx import Presentation

    from engine.render import load_theme, render_presentation
    from schemas.slide_plan import SlidePlan
    plan = SlidePlan.model_validate(plan_data)
    render_presentation(plan, load_theme("huawei"), str(out_path))
    prs = Presentation(str(out_path))
    headings = {f"H{i}" for i in range(n)}
    tops, lefts = set(), set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if (shape.has_text_frame and shape.text_frame.text in headings
                    and shape.top is not None and shape.left is not None):
                tops.add(round(shape.top / 914400, 2))
                lefts.add(round(shape.left / 914400, 2))
    return (len(tops), len(lefts))


# ============================================================
# 老后缀 alias 向后兼容 — path Y (现可达, 不依赖 T8)
# ============================================================

class TestLegacyAliasPathY:
    """老后缀 cards-N 走 content.key_points, legacy 兜底渲染, 全 point 不丢。"""

    @pytest.mark.parametrize("vt,n", [("cards-2", 2), ("cards-4", 4), ("cards-5", 5)])
    def test_legacy_cards_path_y_renders_all_points(self, vt, n, tmp_path):
        blob = _render_blob(_plan_path_y(vt, n), tmp_path / f"{vt}_y.pptx")
        missing = _tokens_present(blob, f"ptY-{vt}", n)
        assert not missing, f"{vt} (path Y) 丢 point: {missing}"

    def test_cards6_path_y_renders_all_points(self, tmp_path):
        """cards-6 (typed variant Cards6Content) path Y 兜底也不丢 (6 点)。"""
        blob = _render_blob(_plan_path_y("cards-6", 6), tmp_path / "c6_y.pptx")
        missing = _tokens_present(blob, "ptY-cards-6", 6)
        assert not missing, f"cards-6 (path Y) 丢 point: {missing}"


# ============================================================
# 老后缀 alias 向后兼容 — path X (等 T8: get_points 读 variant.points)
# ============================================================

@adaptive_only
class TestLegacyAliasPathX:
    """老后缀 cards-N 走 slide.variant.points (CardsContent), T8 后全 point 不丢。"""

    @pytest.mark.parametrize("vt,n", [("cards-2", 2), ("cards-4", 4)])
    def test_legacy_cards_path_x_renders_all_points(self, vt, n, tmp_path):
        blob = _render_blob(_plan_path_x_cards(vt, n), tmp_path / f"{vt}_x.pptx")
        missing = _tokens_present(blob, f"ptX-{vt}", n)
        assert not missing, f"{vt} (path X variant.points) 丢 point: {missing}"


# ============================================================
# 无后缀自适应母版 — 派生布局 + 全 point 不丢 (等 T8)
# ============================================================

@adaptive_only
class TestAdaptiveCards:
    """无后缀 cards 由 len(points) 派生布局, 2/4/6 全渲染。"""

    @pytest.mark.parametrize("n", [2, 4, 6])
    def test_bare_cards_renders_all_points(self, n, tmp_path):
        blob = _render_blob(_plan_path_x_cards("cards", n), tmp_path / f"cards_{n}.pptx")
        missing = _tokens_present(blob, f"ptX-cards", n)
        assert not missing, f"无后缀 cards n={n} 丢 point: {missing}"

    def test_layout_differs_by_point_count(self, tmp_path):
        """2/4/6 point 应派生不同布局 — 用卡片网格 (行数, 列数) 结构化特征作证据。

        卡片标题 shape 文本为 H0..H{n-1}; 其 distinct top 坐标数 = 行数, distinct
        left 坐标数 = 列数 (比 shape 总数计数更 robust, 直接反映布局行列结构, 不受
        装饰/title/footer shape 数量波动干扰)。

        实测派生矩阵 (huawei 主题, T8 后):
          n=2 → (1行, 2列)  单行
          n=4 → (2行, 2列)  双行 2x2
          n=6 → (2行, 3列)  3x2 网格
        三者网格指纹各不相同 — 这是 len(points) 派生布局的确定性证据。
        """
        grids = {n: _card_grid(_plan_path_x_cards("cards", n), tmp_path / f"g_{n}.pptx", n)
                 for n in (2, 4, 6)}
        # 派生形态: n=2 单行, n=4/6 双行
        assert grids[2][0] == 1, f"n=2 应单行, 实际 {grids[2]}"
        assert grids[4][0] == 2, f"n=4 应双行, 实际 {grids[4]}"
        assert grids[6][0] == 2, f"n=6 应双行(网格), 实际 {grids[6]}"
        # 三个 (行,列) 指纹两两不同 — 布局确实随 point 数派生而变
        assert len({grids[2], grids[4], grids[6]}) == 3, (
            f"2/4/6 应派生 3 种不同网格布局, 实际 {grids}"
        )


@adaptive_only
class TestAdaptiveProcess:
    """无后缀 process 自适应 2-5 步。"""

    @pytest.mark.parametrize("n", [2, 3, 5])
    def test_bare_process_renders_all_points(self, n, tmp_path):
        plan = {
            "meta": {"title": "t", "preset": "research-report", "theme": "huawei",
                     "generated_at": "2026-06-18T00:00:00"},
            "narrative": {"thesis": "x", "arc": "a", "total_slides": 1},
            "slides": [{"id": 1, "role": "content", "visual_type": "process",
                        "content": {"title": "流程 path-X"},
                        "variant": {"visual_type": "process", "title": "流程 path-X",
                                    "points": [{"heading": f"S{i}", "body": f"proc-{i}-tok"}
                                               for i in range(n)]}}],
        }
        blob = _render_blob(plan, tmp_path / f"proc_{n}.pptx")
        missing = [f"proc-{i}-tok" for i in range(n) if f"proc-{i}-tok" not in blob]
        assert not missing, f"无后缀 process n={n} 丢 point: {missing}"


@adaptive_only
class TestAdaptiveComparison:
    """无后缀 comparison 自适应列数。"""

    @pytest.mark.parametrize("n", [2, 3])
    def test_bare_comparison_renders_all_points(self, n, tmp_path):
        plan = {
            "meta": {"title": "t", "preset": "research-report", "theme": "huawei",
                     "generated_at": "2026-06-18T00:00:00"},
            "narrative": {"thesis": "x", "arc": "a", "total_slides": 1},
            "slides": [{"id": 1, "role": "content", "visual_type": "comparison",
                        "content": {"title": "对比 path-X"},
                        "variant": {"visual_type": "comparison", "title": "对比 path-X",
                                    "points": [{"heading": f"C{i}", "body": f"cmp-{i}-tok"}
                                               for i in range(n)]}}],
        }
        blob = _render_blob(plan, tmp_path / f"cmp_{n}.pptx")
        missing = [f"cmp-{i}-tok" for i in range(n) if f"cmp-{i}-tok" not in blob]
        assert not missing, f"无后缀 comparison n={n} 丢 point: {missing}"
