"""华为主题 renderer: governance (S4 Phase 1 从 renderers_huawei.py 拆出)。

helper 单一来源在 engine.renderer_kit; 本模块仅含单个 @register_renderer。
"""
from __future__ import annotations

from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from lib.font_fallback import resolve_font_for_pptx
from lib.margins import SafeArea
from schemas.slide_plan import SlideSpec, StructuredPoint

from engine.renderer_kit import (
    register_renderer, _RENDERERS, RendererContext,
    hex_to_rgb, add_textbox, add_textbox_rich, add_rounded_rect, set_slide_background,
    render_title_zone, render_footer, get_content_zone, get_points, get_point_bodies,
    _render_cards, _resolve_accent, _resolve_title_color, _resolve_body_color, _ve,
    compute_card_height, estimate_content_height, normalize_point,
    render_card_header, render_card_number_badge, _get_card_accent,
    _truncate_for_single_line, _add_textbox_no_wrap,
    _px_pt, _px_in, _tokens, _layouts, _color, _lerp_hex, _is_dark, _type_px, _layout,
    _extra, _extra_list_merge, _estimate_text_height_in, _dget, _add_rect,
    _canvas_dims_in, _canvas_padding_in, _slide_dims_in, _font_family, TREND_COLOR_TOKEN,
)


@register_renderer("governance")
def render_governance(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """治理结构：顶部决策委员会 rect + 下方 N 个工作组 rect + 连接线。"""
    ctx = RendererContext.build(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "governance")
    top_h_px = cfg.get("top_box_height_px", 120)
    unit_cols_max = cfg.get("unit_columns_max", 4)

    font = ctx.font
    area_left, area_top, area_w, area_h = ctx.area_left, ctx.area_top, ctx.area_w, ctx.area_h

    # schemas: spec.variant.top (GovTop {name, tag}); legacy: spec.content.top_box {title, desc}
    top_box = _extra(spec, "top_box", "top", default=None) or {}
    units = _extra(spec, "units", default=None) or []
    if not units:
        for p in get_points(spec):
            units.append({"title": p.heading or "", "desc": p.body or ""})

    primary = _color(theme, "primary", "#C7000B")
    paper = _color(theme, "paper", "#FFFFFF")
    paper_2 = _color(theme, "paper_2", "#F4F4F4")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    rule = _color(theme, "rule", "#D9D9D9")
    body_pt = _px_pt(_type_px(theme, "body", 24))
    small_pt = _px_pt(_type_px(theme, "small", 20))
    h4_pt = _px_pt(_type_px(theme, "h4", 28))

    # 顶部框（决策委员会）
    top_h = _px_in(top_h_px)
    top_w = area_w * 0.6
    top_x = area_left + (area_w - top_w) / 2
    top_y = area_top
    add_rounded_rect(slide, top_x, top_y, top_w, top_h, primary, corner_radius=0.1)

    # schemas: GovTop.{name, tag}; legacy: {title, desc}
    if isinstance(top_box, (dict,)) or hasattr(top_box, "__class__") and not isinstance(top_box, str):
        top_title = _dget(top_box, "title", "name", default="")
        top_desc = _dget(top_box, "desc", "tag", default="")
    else:
        top_title = str(top_box) if top_box else ""
        top_desc = ""
    if top_title:
        add_textbox(
            slide, top_x, top_y + 0.15, top_w, h4_pt / 72 * 1.4 + 0.15,
            top_title, font, h4_pt,
            paper, alignment=PP_ALIGN.CENTER, bold=True,
        )
    if top_desc:
        add_textbox(
            slide, top_x, top_y + 0.15 + h4_pt / 72 * 1.4 + 0.1,
            top_w, small_pt / 72 * 1.5 + 0.15,
            top_desc, font, small_pt,
            paper, alignment=PP_ALIGN.CENTER,
        )

    # 底部 N 个工作组
    n = min(len(units), unit_cols_max) or 1
    gap = 0.25
    unit_w = (area_w - gap * (n - 1)) / n
    unit_top = top_y + top_h + 0.6
    # unit_h 按内容估算 (标题 + desc/items), 避免框过高内部大片留空
    _umax = 0.0
    for _it in units[:n]:
        if isinstance(_it, str):
            _ud2 = ""
        else:
            _ud2 = _dget(_it, "desc", "description", default="")
            if not _ud2:
                _items2 = _dget(_it, "items", default=None)
                if isinstance(_items2, (list, tuple)) and _items2:
                    _ud2 = "\n".join(f"• {x}" for x in _items2)
        _dh2 = _estimate_text_height_in(str(_ud2), unit_w - 0.3, small_pt, line_spacing=1.4, padding_in=0.1) if _ud2 else 0.0
        _umax = max(_umax, body_pt / 72 * 1.4 + 0.25 + _dh2)
    unit_h = max(1.3, min(area_h - top_h - 0.6 - 0.1, _umax + 0.3))

    # 连接线：顶框底边中点 → 水平横杆 → 下垂到每个 unit 顶部
    hub_y = top_y + top_h
    trunk_y = hub_y + 0.25
    hub_x_center = top_x + top_w / 2
    # 垂直干线
    slide.shapes.add_connector(
        1, Inches(hub_x_center), Inches(hub_y),
        Inches(hub_x_center), Inches(trunk_y),
    ).line.color.rgb = hex_to_rgb(rule)
    # 水平横杆（覆盖所有 unit 中心）
    if n > 1:
        first_cx = area_left + unit_w / 2
        last_cx = area_left + (n - 1) * (unit_w + gap) + unit_w / 2
        slide.shapes.add_connector(
            1, Inches(first_cx), Inches(trunk_y),
            Inches(last_cx), Inches(trunk_y),
        ).line.color.rgb = hex_to_rgb(rule)

    for i in range(n):
        item = units[i] if i < len(units) else {}
        if isinstance(item, str):
            item = {"title": item, "desc": ""}
        ux = area_left + i * (unit_w + gap)
        # 垂直支线
        cx = ux + unit_w / 2
        slide.shapes.add_connector(
            1, Inches(cx), Inches(trunk_y),
            Inches(cx), Inches(unit_top),
        ).line.color.rgb = hex_to_rgb(rule)

        add_rounded_rect(slide, ux, unit_top, unit_w, unit_h, paper_2, corner_radius=0.08)
        # schemas: GovUnit.{code, name, items[]}; legacy: {title, desc, description}
        ut = _dget(item, "title", "name", default="")
        ud = _dget(item, "desc", "description", default="")
        if not ud:
            items_list = _dget(item, "items", default=None)
            if isinstance(items_list, (list, tuple)) and items_list:
                ud = "\n".join(f"• {it}" for it in items_list)
        # 同时显示 code (如 "L1") 作为 ut 前缀, 若 schemas 给了 code
        unit_code = _dget(item, "code", default="")
        if unit_code and unit_code not in ut:
            ut = f"{unit_code} {ut}".strip()
        add_textbox(
            slide, ux + 0.15, unit_top + 0.15, unit_w - 0.3,
            body_pt / 72 * 1.4 + 0.15,
            ut, font, body_pt,
            ink, bold=True,
        )
        if ud:
            add_textbox(
                slide, ux + 0.15,
                unit_top + 0.15 + body_pt / 72 * 1.4 + 0.1,
                unit_w - 0.3,
                unit_h - 0.4 - body_pt / 72 * 1.4,
                ud, font, small_pt,
                ink_soft,
            )
