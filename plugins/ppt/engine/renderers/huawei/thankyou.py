"""华为主题 renderer: thankyou (S4 Phase 1 从 renderers_huawei.py 拆出)。

helper 单一来源在 engine.renderer_kit; 本模块仅含单个 @register_renderer。
"""
from __future__ import annotations

from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from lib.font_fallback import resolve_font_for_pptx
from lib.margins import SafeArea
from schemas.slide_plan import SlideSpec, StructuredPoint

from engine.renderer_kit import (
    register_renderer, _RENDERERS,
    hex_to_rgb, add_textbox, add_textbox_rich, add_rounded_rect, set_slide_background,
    render_title_zone, render_footer, get_content_zone, get_points, get_point_bodies,
    _render_cards, _resolve_accent, _resolve_title_color, _resolve_body_color, _ve,
    compute_card_height, estimate_content_height, normalize_point,
    render_card_header, render_card_number_badge, _get_card_accent,
    _truncate_for_single_line, _add_textbox_no_wrap,
    _px_pt, _px_in, _tokens, _layouts, _color, _lerp_hex, _is_dark, _type_px, _layout,
    _extra, _extra_list_merge, _estimate_text_height_in, _dget, _add_rect,
    _canvas_dims_in, _canvas_padding_in, _slide_dims_in, _font_family, TREND_COLOR_TOKEN,
)


@register_renderer("thankyou")
def render_thankyou(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """致谢页（浅底华为式, 2026-05-31 taste 选型 A）: 左侧克制大字 + 短红线, 右侧竖排联系方式。

    历史: 旧版红底白字 360px 巨字 (AP2 红满底 + AP3 巨字), 用户评 2/1。现照 golden banking-p043
    重排为白底: 标题降到 ~96px(~58pt), 红仅作短线点缀, 联系方式右侧竖排 label/value。
    """
    sw, sh = _slide_dims_in()
    cfg = _layout(theme, "thankyou")
    set_slide_background(slide, _color(theme, cfg.get("background_token", "paper"), "#FFFFFF"))

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)
    ink = _color(theme, cfg.get("title_color_token", "ink"), "#1F1F1F")
    red = _color(theme, "primary", "#C7000B")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    ink_mute = _color(theme, "ink_mute", "#9A9A9A")

    # 2026-06-12 issue: 标题从 160px(~96pt) 收敛到 ~52pt (对照 golden banking-p043, 旧值命中 AP3 巨字)
    title_pt = _px_pt(cfg.get("title_size_px", 88))     # ~52.8pt
    eyebrow_pt = _px_pt(cfg.get("eyebrow_size_px", 27)) # ~16pt 红色小标签
    emph_pt = _px_pt(cfg.get("emphasis_size_px", 54))   # ~32pt 红色 punchline (第二行强调)
    sub_pt = _px_pt(cfg.get("subtitle_size_px", 30))    # ~18pt
    label_pt = _px_pt(cfg.get("label_size_px", 20))     # ~12pt
    value_pt = _px_pt(cfg.get("value_size_px", 26))     # ~15.6pt
    left_w = cfg.get("left_column_in", 6.5)
    right_x = cfg.get("right_column_x_in", 8.0)

    # 消费 variant: eyebrow / title / emphasis (历史只读 content.title, eyebrow+emphasis 丢失)
    title_text = _extra(spec, "title", default=None) or spec.content.title or "谢谢"
    eyebrow_text = _extra(spec, "eyebrow", default=None)
    emphasis_text = _extra(spec, "emphasis", default=None)

    # 左侧栏: eyebrow → 标题 → 短红线 → emphasis (顶部纵向堆叠, 下划线随标题框避让, 不再固定 4.35 压字)
    lx = left_pad + 0.3
    y = 2.0
    if eyebrow_text:
        eb_h = eyebrow_pt / 72 * 1.5 + 0.05
        add_textbox(slide, lx, y, left_w, eb_h, str(eyebrow_text), font, eyebrow_pt, red, bold=True)
        y += eb_h + 0.12
    # 标题 (克制 ~52pt); 长中文标题在 left_w 内会折行, 用估算高度 (含折行) 让下划线/emphasis 避让,
    # 否则按单行高度放下划线会压住折行后的第二行 (2026-06-12 视觉验证抓到的回归)。
    title_h = max(
        title_pt / 72 * 1.3 + 0.1,
        _estimate_text_height_in(title_text, left_w, title_pt, line_spacing=1.15, padding_in=0.1),
    )
    tbx = slide.shapes.add_textbox(Inches(lx), Inches(y), Inches(left_w), Inches(title_h))
    tf = tbx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = font
    p.font.size = Pt(title_pt)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(ink)
    y += title_h + 0.12
    # 标题下短红线 (位于标题实际底部之下, 随折行高度避让, 不再压字)
    _add_rect(slide, lx + 0.03, y, 1.6, 0.06, red)
    y += 0.06 + 0.2
    # emphasis 红色 punchline (仅当 variant 提供)
    if emphasis_text:
        emph_h = emph_pt / 72 * 1.3 + 0.15
        add_textbox(slide, lx, y, left_w, emph_h, str(emphasis_text), font, emph_pt, red, bold=True)

    # 右侧文字栏: 副标题 (从 variant 取, 历史只读 content.subtitle 在 path X 下丢失) + 联系方式
    right_w = sw - right_x - right_pad
    subtitle_text = _extra(spec, "subtitle", default=None)
    if subtitle_text:
        add_textbox(
            slide, right_x, 1.9, right_w, 1.0,
            str(subtitle_text), font, sub_pt, ink_soft,
        )
    max_contacts = cfg.get("contact_columns", 4)  # 上限防 6+ 联系方式溢出底边 (保留旧 contacts[:n] 语义)
    contacts = (_extra(spec, "contacts", default=None) or [])[:max_contacts]
    cy = 3.1
    for c in contacts:
        label_v = _dget(c, "label", default="")
        value_v = _dget(c, "value", default="")
        if label_v:
            add_textbox(slide, right_x, cy, right_w, 0.3, str(label_v), font, label_pt, ink_mute, bold=True)
        if value_v:
            add_textbox(slide, right_x, cy + 0.28, right_w, 0.34, str(value_v), font, value_pt, ink_soft)
        cy += 0.78


# ===========================================================================
# P2 四版式（+ V17 变体集合由 design.layout_variant 分派）
# ===========================================================================

