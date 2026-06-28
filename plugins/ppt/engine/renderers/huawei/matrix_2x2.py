"""华为主题 renderer: matrix-2x2 (S4 Phase 1 从 renderers_huawei.py 拆出)。

helper 单一来源在 engine.renderer_kit; 本模块仅含单个 @register_renderer。
"""
from __future__ import annotations

from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from lib.font_fallback import resolve_font_for_pptx
from lib.margins import SafeArea
from schemas.slide_plan import SlideSpec, StructuredPoint

from engine.renderer_kit import (
    register_renderer, _RENDERERS, RendererContext,
    hex_to_rgb, add_textbox, add_textbox_rich, add_rounded_rect, set_slide_background,
    render_title_zone, render_footer, get_content_zone, get_points, get_point_bodies,
    _render_cards, _resolve_accent, _resolve_title_color, _resolve_body_color, _ve,
    compute_card_height, estimate_content_height, normalize_point,
    render_card_header, render_card_number_badge, _get_card_accent,
    _truncate_for_single_line, _add_textbox_no_wrap,
    _px_pt, _px_in, _resolve_gap_in, _tokens, _layouts, _color, _lerp_hex, _is_dark, _type_px, _layout,
    _extra, _extra_list_merge, _estimate_text_height_in, _dget, _add_rect,
    _canvas_dims_in, _canvas_padding_in, _slide_dims_in, _font_family, TREND_COLOR_TOKEN,
)
from engine.layout_engine import Region, split_grid


@register_renderer("matrix-2x2")
def render_matrix_2x2(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """2x2 矩阵：左 Y 轴 + 底 X 轴 + 4 象限卡片，可按 highlight 标记高亮象限。"""
    ctx = RendererContext.build(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "matrix_2x2")
    y_col_px = cfg.get("axis_y_column_px", 140)
    x_row_px = cfg.get("axis_x_row_px", 140)
    gap_px = cfg.get("grid_gap_px", 20)
    border_px = cfg.get("border_width_px", 2)
    border_token = cfg.get("border_token", "ink")
    heading_px = cfg.get("quadrant_heading_size_px", 22)
    y_label_px = cfg.get("y_axis_label_size_px", 20)
    small_px = cfg.get("small_label_size_px", 14)
    desc_px = cfg.get("desc_size_px", 16)

    font = ctx.font
    area_left, area_top, area_w, area_h = ctx.area_left, ctx.area_top, ctx.area_w, ctx.area_h

    y_col_in = _px_in(y_col_px)
    x_row_in = _px_in(x_row_px)
    gap_in = _px_in(gap_px)

    grid_left = area_left + y_col_in
    grid_top = area_top
    grid_w = area_w - y_col_in
    grid_h = area_h - x_row_in
    cell_w = (grid_w - gap_in) / 2
    cell_h = (grid_h - gap_in) / 2

    border_color = _color(theme, border_token, "#1F1F1F")
    primary = _color(theme, "primary", "#C7000B")
    primary_soft = _color(theme, "primary_soft", "#FDECEE")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    ink_mute = _color(theme, "ink_mute", "#9A9A9A")
    paper_2 = _color(theme, "paper_2", "#F4F4F4")

    # path X: variant.axis (MatrixAxis{x: Axis, y: Axis}, 各含 label)。
    # 历史 bug: renderer 只读顶层 y_axis/x_axis/axis_y/axis_x, schema 的嵌套 axis.x/y.label 静默丢失。
    axis_obj = _extra(spec, "axis", default=None)
    _axis_y_label = None
    _axis_x_label = None
    if axis_obj is not None:
        _ay = _dget(axis_obj, "y", default=None)
        _ax = _dget(axis_obj, "x", default=None)
        if _ay is not None:
            _axis_y_label = _dget(_ay, "label", default=None)
        if _ax is not None:
            _axis_x_label = _dget(_ax, "label", default=None)

    # Y 轴（左侧 — 2026-05-31 taste B 干净版: 主标签(括号前)真竖排 eaVert 字正立,
    # 括号补充作横排小字副标放底部 (避免 KD/括号等西文符号在 eaVert 竖排下侧躺)）
    y_axis = _axis_y_label or _extra(spec, "y_axis", default=None) or _extra(spec, "axis_y", default="")
    if y_axis:
        y_full = str(y_axis)
        # 拆: 括号前为竖排主标签, 括号内容为横排副标 (全角/半角括号都认)
        idx = y_full.find("（")
        if idx < 0:
            idx = y_full.find("(")
        if idx > 0:
            y_main = y_full[:idx].strip()
            y_sub = y_full[idx:].strip().strip("（）()").strip()
        else:
            y_main, y_sub = y_full, ""
        avail_h = grid_h - (0.5 if y_sub else 0.0)  # 有副标时底部留空
        n_ch = max(1, len(y_main))
        default_pt = _px_pt(y_label_px)
        # 竖排每字占高 ≈ pt/72 × 1.15 行距; 主标签须 ≤ avail_h, 否则缩字号 (下限 9pt 保可读)
        fit_pt = avail_h * 72 / n_ch / 1.15
        y_pt = max(9.0, min(default_pt, fit_pt))
        vbox = slide.shapes.add_textbox(
            Inches(area_left), Inches(grid_top),
            Inches(y_col_in), Inches(avail_h),
        )
        vtf = vbox.text_frame
        vtf.word_wrap = True
        vtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # eaVert = 东亚竖排, 字符正立自上而下 (区别于 rotation=270 整体旋转导致字侧躺)
        bodyPr = vtf._txBody.find(qn("a:bodyPr"))
        if bodyPr is None:
            bodyPr = vtf._txBody.makeelement(qn("a:bodyPr"), {})
            vtf._txBody.insert(0, bodyPr)
        bodyPr.set("vert", "eaVert")
        vp = vtf.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        vp.text = y_main
        vp.font.name = font
        vp.font.size = Pt(y_pt)
        vp.font.bold = True
        vp.font.color.rgb = hex_to_rgb(ink_soft)
        # 括号副标: 横排小字放 axis 列底部 (横排不侧躺)
        if y_sub:
            add_textbox(
                slide, area_left - 0.15, grid_top + grid_h - 0.45,
                y_col_in + 0.3, 0.42,
                y_sub, font, max(8.0, _px_pt(y_label_px) * 0.7),
                ink_mute, alignment=PP_ALIGN.CENTER,
            )

    # X 轴（底部）
    x_axis = _axis_x_label or _extra(spec, "x_axis", default=None) or _extra(spec, "axis_x", default="")
    if x_axis:
        add_textbox(
            slide, grid_left, grid_top + grid_h + 0.15, grid_w, x_row_in - 0.2,
            str(x_axis), font, _px_pt(y_label_px),
            ink_soft, alignment=PP_ALIGN.CENTER, bold=True,
        )

    # 4 象限（quadrants）
    quadrants = _extra(spec, "quadrants", default=None) or []
    if not quadrants:
        pts = get_points(spec)
        for p in pts[:4]:
            quadrants.append({
                "heading": p.heading or "",
                "desc": p.body or "",
                "highlight": False,
            })
    # 顺序约定：左上 / 右上 / 左下 / 右下
    # S4 阶段3: 2×2 网格切分交给 LayoutEngine, 在 Y 轴列/X 轴行 (fixed-leading) 之后的 grid 子
    # region 内分。split_grid row-major(idx 0=左上,1=右上,2=左下,3=右下) 与原固定坐标逐 EMU 恒等。
    # cell_w/cell_h 保留供下方 rect 宽高 + textbox 宽 (== cells[i].width/height)。
    layout_spec = cfg.get("layout")
    if layout_spec is not None:
        grid_region = Region(grid_left, grid_top, grid_w, grid_h)
        _cells = split_grid(grid_region, layout_spec["rows"], layout_spec["cols"],
                            _resolve_gap_in(layout_spec, cfg))
        positions = [(c.left, c.top) for c in _cells]
    else:
        positions = [
            (grid_left, grid_top),
            (grid_left + cell_w + gap_in, grid_top),
            (grid_left, grid_top + cell_h + gap_in),
            (grid_left + cell_w + gap_in, grid_top + cell_h + gap_in),
        ]
    heading_pt = _px_pt(heading_px)
    desc_pt = _px_pt(desc_px)
    small_pt = _px_pt(small_px)

    for i, (x, y) in enumerate(positions):
        item = quadrants[i] if i < len(quadrants) else {}
        if isinstance(item, str):
            item = {"heading": "", "desc": item}
        highlight = bool(_dget(item, "highlight", default=False))
        fill = primary_soft if highlight else paper_2
        _add_rect(slide, x, y, cell_w, cell_h, fill,
                  line_hex=border_color, line_width_pt=max(border_px * 0.5, 0.5))

        # 角标（tag / corner_tag 单数兼容 path Y; tags 列表为 path X schema QuadrantItem 字段）
        tag = _dget(item, "tag", default="") or _dget(item, "corner_tag", default="")
        if not tag:
            _tags = _dget(item, "tags", default=None)
            if _tags:
                if isinstance(_tags, (list, tuple)):
                    tag = " · ".join(str(t) for t in _tags if t)
                else:
                    tag = str(_tags)
        if tag:
            add_textbox(
                slide, x + 0.15, y + 0.1, cell_w - 0.3, small_pt / 72 * 1.5 + 0.1,
                str(tag), font, small_pt,
                primary if highlight else ink_mute, bold=True,
            )
        heading = _dget(item, "heading", default="") or _dget(item, "title", default="")
        add_textbox(
            slide, x + 0.2, y + 0.4, cell_w - 0.4, heading_pt / 72 * 1.4 + 0.2,
            str(heading), font, heading_pt,
            ink, bold=True,
        )
        desc = _dget(item, "desc", default="") or _dget(item, "description", default="") or _dget(item, "body", default="")
        if desc:
            add_textbox(
                slide, x + 0.2, y + 0.4 + heading_pt / 72 * 1.4 + 0.1,
                cell_w - 0.4, cell_h - 0.4 - heading_pt / 72 * 1.4 - 0.3,
                str(desc), font, desc_pt,
                ink_soft,
            )


