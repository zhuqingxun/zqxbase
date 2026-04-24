"""pptx_compat - python-pptx internal API compatibility wrapper.

Centralizes all private/undocumented python-pptx API access into a single
module. Public python-pptx API (Presentation, Slide, shapes, text, etc.)
is used directly by consuming modules -- NOT wrapped here.

Only internal APIs (underscore-prefixed attributes, oxml subpackages,
lxml element manipulation) are wrapped. This provides:
  - Single point of maintenance when python-pptx internals change
  - Import-time version validation
  - Clear documentation of version stability per function

Supported versions: python-pptx >=1.0.0, <2.0
"""

import importlib.metadata

# --- Version validation (import-time) ---
PPTX_VERSION = importlib.metadata.version("python-pptx")
_version_parts = PPTX_VERSION.split(".")
_MAJOR = int(_version_parts[0])
_MINOR = int(_version_parts[1]) if len(_version_parts) > 1 else 0
_PATCH = int(_version_parts[2]) if len(_version_parts) > 2 else 0

if _MAJOR < 1:
    raise ImportError(
        f"python-pptx {PPTX_VERSION} not supported. Requires >=1.0.0,<2.0. "
        f"Run: pip install 'python-pptx>=1.0.0,<2.0'"
    )
if _MAJOR >= 2:
    raise ImportError(
        f"python-pptx {PPTX_VERSION} not supported. Requires >=1.0.0,<2.0. "
        f"Major version bump may break internal APIs."
    )

# --- Core dependencies (after validation) ---
from pptx.oxml.ns import qn  # noqa: E402
from pptx.oxml import parse_xml  # noqa: E402
from pptx.oxml.xmlchemy import OxmlElement  # noqa: E402
from lxml import etree  # noqa: E402


# =============================================================================
# Category 1: Run Properties
# =============================================================================

def get_run_properties(run):
    """Get or create run properties element (a:rPr).

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: run._r.get_or_add_rPr()
    """
    return run._r.get_or_add_rPr()


# =============================================================================
# Category 2: Paragraph Properties
# =============================================================================

def get_paragraph_properties(paragraph):
    """Get or create paragraph properties element (a:pPr).

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: paragraph._p.get_or_add_pPr()
    """
    return paragraph._p.get_or_add_pPr()


def get_paragraph_element(paragraph):
    """Get the backing CT_TextParagraph element for direct access.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: paragraph._p
    """
    return paragraph._p


def find_paragraph_pPr(paragraph):
    """Find existing paragraph properties element (read-only, no create).

    Returns None if no pPr element exists.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: paragraph._p.find(qn("a:pPr"))
    """
    return paragraph._p.find(qn("a:pPr"))


# =============================================================================
# Category 3: Shape Element
# =============================================================================

def get_shape_element(shape):
    """Get the backing lxml element for a shape.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: shape._element
    """
    return shape._element


def remove_shape(shape):
    """Remove a shape from its parent slide by removing its XML element.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: shape._element.getparent().remove()
    """
    sp = shape._element
    sp.getparent().remove(sp)


def get_table_from_shape(table_shape):
    """Get the underlying tbl XML element from a table shape.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: table_shape._element.graphic.graphicData.tbl
    """
    return table_shape._element.graphic.graphicData.tbl


# =============================================================================
# Category 4: Slide Collection
# =============================================================================

def get_slide_id_list(prs):
    """Get the internal slide ID list for reordering/deletion.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: prs.slides._sldIdLst (CT_SlideIdList element)
    """
    return prs.slides._sldIdLst


def delete_slide(prs, index):
    """Delete a slide by index.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: prs.slides._sldIdLst, prs.part.drop_rel()
    """
    rId = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


# =============================================================================
# Category 5: Shape Tree
# =============================================================================

def insert_shape_element(slide, element):
    """Insert a shape element into a slide's shape tree.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: slide.shapes._spTree.insert_element_before()
    """
    slide.shapes._spTree.insert_element_before(element, 'p:extLst')


# =============================================================================
# Category 6: Relationship Access
# =============================================================================

def get_part_rels(part):
    """Get relationships collection from a part.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: part.rels (semi-public but undocumented)
    """
    return part.rels


def get_rel_target(rel):
    """Get the target of a relationship.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: rel._target
    """
    return rel._target


def drop_rel(part, rId):
    """Drop a relationship from a part by relationship ID.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: part.drop_rel()
    """
    part.drop_rel(rId)


def get_or_add_rel(rels, reltype, target_part):
    """Get existing or add new relationship to a part.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: rels.get_or_add(reltype, target_part)
    """
    return rels.get_or_add(reltype, target_part)


# =============================================================================
# Category 7: Theme/Part Access
# =============================================================================

def get_theme_part(prs):
    """Get the theme part from the first slide master.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: prs.slide_masters[0].part.rels, rel.target_part
    """
    for rel in prs.slide_masters[0].part.rels.values():
        if 'theme' in rel.reltype:
            return rel.target_part
    return None


def get_part_blob(part):
    """Get raw XML bytes from a part.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: part.blob
    """
    return part.blob


def get_presentation_from_slide(slide):
    """Get the Presentation element from a slide via internal part chain.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: slide.part.package.presentation_part.presentation
    """
    return slide.part.package.presentation_part.presentation


def get_slide_layout_master(shape):
    """Get the slide master from a shape's slide layout.

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: shape.part.slide_layout.slide_master
    """
    return shape.part.slide_layout.slide_master


# =============================================================================
# Category 8: Cell Properties
# =============================================================================

def get_cell_properties(cell):
    """Get or create table cell properties element (a:tcPr).

    Works: python-pptx >=0.6.21 through 1.0.2
    Internal: cell._tc.get_or_add_tcPr()
    """
    return cell._tc.get_or_add_tcPr()


# =============================================================================
# Category 9: XML Building (re-exports for centralized access)
# =============================================================================

def make_sub_element(parent, nsprefix_tag):
    """Create a sub-element with namespace-qualified tag.

    Wraps: etree.SubElement(parent, qn(nsprefix_tag))
    """
    return etree.SubElement(parent, qn(nsprefix_tag))


def find_element(parent, nsprefix_tag):
    """Find child element by namespace-qualified tag.

    Wraps: parent.find(qn(nsprefix_tag))
    """
    return parent.find(qn(nsprefix_tag))


def find_all_elements(parent, nsprefix_tag):
    """Find all child elements by namespace-qualified tag.

    Wraps: parent.findall(qn(nsprefix_tag))
    """
    return parent.findall(qn(nsprefix_tag))


# Re-exports: consumers import these from pptx_compat instead of pptx.oxml
# qn, parse_xml, OxmlElement, etree are already imported above and available.


# =============================================================================
# Version Info
# =============================================================================

def get_version():
    """Return the detected python-pptx version string."""
    return PPTX_VERSION


def supports_version(min_version="1.0.0"):
    """Check if installed version meets minimum requirement.

    Args:
        min_version: Dotted version string (e.g., "1.0.0")

    Returns:
        True if installed version >= min_version
    """
    min_parts = tuple(int(x) for x in min_version.split("."))
    cur_parts = tuple(int(x) for x in PPTX_VERSION.split(".")[:len(min_parts)])
    return cur_parts >= min_parts
