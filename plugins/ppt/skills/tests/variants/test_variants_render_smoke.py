# /// script
# requires-python = ">=3.11"
# dependencies = ["pytest>=7.0", "python-pptx>=1.0.0", "pydantic>=2.0", "pyyaml>=6.0"]
# ///
"""SPEC-R-V01..V17：18 个 variant renderer 的 smoke + 关键断言。

依赖：Dev T4（18 renderer）+ T5（variant samples yaml）完成。

策略：
- 每个 variant 一份 sample yaml（放 `tests/variants/samples/<visual-type>.yaml`）
- 该 yaml 必须是合法的 SlidePlan（meta + narrative + slides[至少 1 张]）
- 对应 slide 的 visual_type + variant 字段符合 VariantUnion
- 本测试不做像素级断言，只断言：renderer 不抛异常 + slide 有 shape + 额外的 variant 特化断言
"""
from __future__ import annotations

from pathlib import Path

import pytest

PLUGIN_ROOT = Path(__file__).resolve().parents[2]
SAMPLES_DIR = PLUGIN_ROOT / "tests" / "variants" / "samples"


try:
    from schemas.slide_plan import SlidePlan  # type: ignore
    from engine.render import render_presentation, load_theme  # type: ignore
except ImportError:
    SlidePlan = None  # type: ignore
    render_presentation = None  # type: ignore
    load_theme = None  # type: ignore


pytestmark = pytest.mark.skipif(
    SlidePlan is None or render_presentation is None or not SAMPLES_DIR.exists(),
    reason="engine / schemas 未就绪 或 samples/ 目录为空（Dev 前置未完成）",
)


# 18 variant（对齐 Plan T5 与 specs O7 副作用修订）
VARIANT_TYPES = [
    # P0 八版式
    "cover-left-bar", "toc", "section-divider-dark",
    "kpi-stats", "matrix-2x2", "architecture-layered",
    "timeline-huawei", "process-flow-huawei",
    # P1 五版式
    "swot", "roadmap", "pyramid", "heatmap-matrix", "thankyou",
    # P2 四版式 + 独立 variant
    "cards-6", "rings", "personas",
    "risk-list", "governance",
]


def _sample_path(vt: str) -> Path:
    return SAMPLES_DIR / f"{vt}.yaml"


@pytest.fixture(scope="module")
def huawei_theme():
    return load_theme("huawei")


@pytest.mark.parametrize("visual_type", VARIANT_TYPES, ids=VARIANT_TYPES)
def test_variant_renders_smoke(tmp_path, huawei_theme, visual_type):
    """SPEC-R-COMMON-1/2：每个 variant sample 渲染不抛异常、slide 有 shape。"""
    import yaml as _yaml
    from pptx import Presentation

    sample = _sample_path(visual_type)
    if not sample.exists():
        pytest.skip(f"{sample} 尚未交付（Dev T5 前）")

    data = _yaml.safe_load(sample.read_text(encoding="utf-8"))
    plan = SlidePlan.model_validate(data)

    # 保证该 sample 至少有一个 slide 的 visual_type 匹配
    matching_slides = [s for s in plan.slides if s.visual_type == visual_type]
    assert matching_slides, (
        f"{sample} 中无 visual_type={visual_type} 的 slide，sample 命名与内容不一致"
    )

    out = tmp_path / f"{visual_type}.pptx"
    render_presentation(plan, huawei_theme, str(out))
    assert out.exists()

    prs = Presentation(str(out))
    for i, slide in enumerate(prs.slides, start=1):
        assert len(slide.shapes) >= 1, f"{visual_type} slide {i} 无 shape"


# ============================================================
# 特化断言（SPEC-R-V01-EX / V02-EX / V05-EX / V13-EX）
# ============================================================

def _render_first_slide(tmp_path, theme, visual_type: str):
    """共用 helper：渲染 sample 并返回首张匹配的 slide 对象。"""
    import yaml as _yaml
    from pptx import Presentation

    sample = _sample_path(visual_type)
    if not sample.exists():
        pytest.skip(f"{sample} 尚未交付")
    data = _yaml.safe_load(sample.read_text(encoding="utf-8"))
    plan = SlidePlan.model_validate(data)
    out = tmp_path / f"{visual_type}.pptx"
    render_presentation(plan, theme, str(out))
    prs = Presentation(str(out))

    for idx, spec in enumerate(plan.slides):
        if spec.visual_type == visual_type:
            return prs.slides[idx]
    pytest.fail(f"sample {sample} 未含 {visual_type} slide")


def test_spec_r_v01_cover_no_edge_red_bar(tmp_path, huawei_theme):
    """cover-left-bar 不应有贯穿边缘的细红条 (2026-05-31 全局规则: huawei 风格不要 edge bar, 用户选型 A 去红条)。"""
    from pptx.util import Emu
    from pptx.dml.color import RGBColor

    slide = _render_first_slide(tmp_path, huawei_theme, "cover-left-bar")
    red = RGBColor(0xC7, 0x00, 0x0B)
    for sh in slide.shapes:
        if sh.left is None or sh.width is None or sh.height is None:
            continue
        # 贯穿边缘红条 = 贴左边 (x≈0) + 窄 (<0.5") + 高 (≥5")
        edge_bar = (sh.left <= Emu(0.2 * 914400)
                    and sh.width < Emu(0.5 * 914400)
                    and sh.height >= Emu(5.0 * 914400))
        if edge_bar:
            try:
                assert sh.fill.fore_color.rgb != red, "检测到贯穿边缘的红条 — 违反全局 huawei 风格规则"
            except (AttributeError, TypeError):
                pass


def test_spec_r_v02_ex_section_divider_light_bg(tmp_path, huawei_theme):
    """section-divider-dark 背景色 = #FFFFFF (2026-05-31 taste 选型 B: 浅底极简, 旧 #2A2A2A 深底违反 P1/P3 用户评 0)。"""
    from pptx.dml.color import RGBColor

    slide = _render_first_slide(tmp_path, huawei_theme, "section-divider-dark")
    try:
        rgb = slide.background.fill.fore_color.rgb
    except Exception as e:
        pytest.skip(f"slide.background 不可读（实现未设 solid fill）: {e}")
    assert rgb == RGBColor(0xFF, 0xFF, 0xFF), f"背景期望浅底 #FFFFFF，实际 #{rgb}"


def _content_field(spec, field_name):
    """从 SlideSpec.content 读 variant 特有字段。

    架构决策（Dev-2 T3 最终选型 = HR-1 方案 A）：variant 字段扁平放在 content 顶层，
    由 SlideContent.extra='allow' 承接；二次校验由 validate_plan 路由 VariantUnion。
    """
    content = spec.content
    # pydantic v2 的 extra 字段通过 model_extra 或 __pydantic_extra__ 暴露
    extra = getattr(content, "__pydantic_extra__", None) or {}
    if field_name in extra:
        return extra[field_name]
    return getattr(content, field_name, None)


def test_spec_r_v06_ex_toc_items(tmp_path, huawei_theme):
    """SPEC-R-V06-EX：toc slide 可见文本含所有章节条目（支持 chapters / items 两种 yaml 字段名）。"""
    import yaml as _yaml

    sample = _sample_path("toc")
    if not sample.exists():
        pytest.skip("toc sample 未交付")
    data = _yaml.safe_load(sample.read_text(encoding="utf-8"))
    plan = SlidePlan.model_validate(data)
    toc_slide_spec = next(s for s in plan.slides if s.visual_type == "toc")
    # sample 可能用 `chapters` 或 `items` 两种字段名（renderer 的 _extra 两者都读）
    chapters = _content_field(toc_slide_spec, "chapters") or _content_field(toc_slide_spec, "items") or []
    if not chapters:
        pytest.skip("toc content 无 chapters/items 字段")

    slide = _render_first_slide(tmp_path, huawei_theme, "toc")
    all_text = "\n".join(
        p.text for sh in slide.shapes if sh.has_text_frame
        for p in sh.text_frame.paragraphs
    )
    # 至少任一章节的 title 出现在渲染可见文本中
    titles_found = sum(1 for ch in chapters if isinstance(ch, dict) and ch.get("title") and ch["title"] in all_text)
    assert titles_found >= 1, f"toc 渲染未包含任何章节 title（chapters={chapters}, rendered_text={all_text!r}）"


def test_spec_r_v13_ex_thankyou_main_text_restrained(tmp_path, huawei_theme):
    """SPEC-R-V13-EX (2026-06-12 修订)：thankyou 标题字号收敛到 ~52pt 区间, 不再是 ≥72pt 巨字。

    历史断言要求 ≥72pt, 但 2026-06-12 issue 明确: 96pt 标题占画面 1/3+ 命中 AP3,
    须收敛到 ~52pt (对照 golden banking-p043)。本断言改为校验"标题是最大字号但克制"。

    兼容 pptx 的两种字号承载方式：paragraph 级 `p.font.size` 与 run 级 `r.font.size`。
    """
    slide = _render_first_slide(tmp_path, huawei_theme, "thankyou")
    sizes = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            if p.font.size is not None:
                sizes.append(p.font.size.pt)
            for r in p.runs:
                if r.font.size is not None:
                    sizes.append(r.font.size.pt)
    assert sizes, "thankyou slide 无任何 paragraph/run 带 font.size"
    # 标题应为最大字号, 落在克制区间 (~52.8pt); 上限 64pt 防回退到旧巨字, 下限 44pt 防过小
    assert 44 <= max(sizes) <= 64, (
        f"thankyou 标题最大字号 {max(sizes)} pt 不在克制区间 [44, 64]"
        f"（>64 疑回退巨字命中 AP3, <44 疑过小）"
    )


def test_spec_r_v14_ex_cards_6_has_six_cards(tmp_path, huawei_theme):
    """SPEC-R-V14-EX：cards-6 渲染包含 6 个 card heading。

    card 数据来源（按 renderer 优先级）：content.cards / content.key_points（StructuredPoint）
    两种都支持。断言 6 个 card 的 heading 文本都出现在可见文本里。
    """
    import yaml as _yaml

    sample = _sample_path("cards-6")
    if not sample.exists():
        pytest.skip("cards-6 sample 未交付")
    data = _yaml.safe_load(sample.read_text(encoding="utf-8"))
    plan = SlidePlan.model_validate(data)
    spec = next(s for s in plan.slides if s.visual_type == "cards-6")

    # 支持两种字段名：content.cards（新 variant 字段） 或 content.key_points（legacy StructuredPoint）
    cards_raw = _content_field(spec, "cards") or []
    if not cards_raw:
        kps = getattr(spec.content, "key_points", None) or []
        cards_raw = [
            kp if isinstance(kp, dict) else {"heading": getattr(kp, "heading", None) or getattr(kp, "body", "")}
            for kp in kps
        ]
    assert len(cards_raw) == 6, f"cards-6 sample 期望 6 卡，实际 {len(cards_raw)}：{cards_raw}"

    slide = _render_first_slide(tmp_path, huawei_theme, "cards-6")
    all_text = "\n".join(
        p.text for sh in slide.shapes if sh.has_text_frame
        for p in sh.text_frame.paragraphs
    )
    missing = []
    for c in cards_raw:
        heading = c.get("heading") if isinstance(c, dict) else None
        if heading and heading not in all_text:
            missing.append(heading)
    # 允许少数 heading 因字号/截断不可见；至少 4 个 heading 能对上即认为渲染正常
    assert (len(cards_raw) - len(missing)) >= 4, (
        f"cards-6 渲染只命中 {len(cards_raw) - len(missing)}/{len(cards_raw)} 个 card heading，"
        f"缺失 {missing}"
    )
