# /// script
# requires-python = ">=3.11"
# dependencies = ["pytest>=7.0", "python-pptx>=1.0.0", "pydantic>=2.0", "pyyaml>=6.0"]
# ///
"""path X (variant typed union) 渲染回归测试 — 全 18 版式。

背景 (2026-06-12 issue-ppt-renderer-variant-consumption-gaps):
现有 `test_variants_render_smoke.py` 的 sample 把 variant 字段**扁平放在 content 顶层** (path Y),
渲染层经由 `_extra` 检查 content 兜底, 因此能渲染出内容。但 codex-fusion live E2E deck 走的是
**path X** (数据放在 `slide.variant`, 被 schema 解析成 typed VariantUnion 子模型)。renderer 历史上
用错字段名 / 错路径取数 (如 roadmap 读 lanes/bars 而 schema 是 rows/cells), path X 下静默丢内容。

本测试为每个 visual_type 构造一份 **path X** plan (data 全部在 `variant` 块, content 仅留必填 title),
渲染后断言"该版式专属的关键文本"出现在 PPTX 可见文本里。覆盖文本框 + 表格单元格 (heatmap 用 table)。
"""
from __future__ import annotations

from pathlib import Path

import pytest

PLUGIN_ROOT = Path(__file__).resolve().parents[2]

try:
    from schemas.slide_plan import SlidePlan  # type: ignore
    from engine.render import render_presentation, load_theme  # type: ignore
except ImportError:
    SlidePlan = None  # type: ignore
    render_presentation = None  # type: ignore
    load_theme = None  # type: ignore


pytestmark = pytest.mark.skipif(
    SlidePlan is None or render_presentation is None,
    reason="engine / schemas 未就绪",
)


# ---------------------------------------------------------------------------
# 每个 visual_type: (variant 块, 必须出现的关键文本列表)
# variant 块严格符合 schemas/variants.py 的对应子模型 (path X)。
# ---------------------------------------------------------------------------
PATH_X_CASES: dict[str, tuple[dict, list[str]]] = {
    "cover-left-bar": (
        {
            "visual_type": "cover-left-bar",
            "title": "封面主标题",
            "eyebrow": "EB标识XY",
            "emphasis": "强调短语XY",
            "subtitle": "封面副标题XY",
            "meta": [{"label": "出品方", "value": "研究院ZZ"}],
        },
        ["EB标识XY", "强调短语XY", "封面副标题XY", "研究院ZZ"],
    ),
    "toc": (
        {
            "visual_type": "toc",
            "title": "目录",
            "chapters": [
                {"number": "01", "title": "第一章AB"},
                {"number": "02", "title": "第二章CD"},
                {"number": "03", "title": "第三章EF"},
            ],
        },
        ["第一章AB", "第二章CD", "第三章EF"],
    ),
    "section-divider-dark": (
        {
            "visual_type": "section-divider-dark",
            "big_number": "02",
            "eyebrow": "章节EB",
            "title": "分隔标题",
            "description": "分隔描述文本XYZ",
        },
        ["章节EB", "分隔描述文本XYZ"],
    ),
    "kpi-stats": (
        # 定性值 (无数字) → 走 qualitative 竖排分支, 该分支历史从 content.subtitle 取数 (path X 下丢)
        {
            "visual_type": "kpi-stats",
            "title": "KPI面板",
            "subtitle": "KPI副标题XY",
            "kpis": [
                {"label": "维度甲", "value": "优"},
                {"label": "维度乙", "value": "良"},
                {"label": "维度丙", "value": "中"},
            ],
        },
        ["KPI副标题XY", "维度甲", "维度乙"],
    ),
    "matrix-2x2": (
        {
            "visual_type": "matrix-2x2",
            "title": "二维矩阵",
            "axis": {"x": {"label": "X轴标签"}, "y": {"label": "Y轴标签"}},
            "quadrants": [
                {"position": "Q1", "heading": "象限一", "tags": ["角标甲", "角标乙"]},
                {"position": "Q2", "heading": "象限二"},
                {"position": "Q3", "heading": "象限三"},
                {"position": "Q4", "heading": "象限四"},
            ],
        },
        # "角标甲" 验证 QuadrantItem.tags 渲染 (历史 bug: renderer 只读单数 tag/corner_tag, schema 的 tags 列表静默丢)
        ["X轴标签", "Y轴标签", "象限一", "象限四", "角标甲"],
    ),
    "architecture-layered": (
        {
            "visual_type": "architecture-layered",
            "title": "分层架构",
            "layers": [
                {"header": {"code": "L1", "name": "接入层"}, "cells": [{"title": "网关甲", "highlight": True}]},
                {"header": {"code": "L2", "name": "服务层"}, "cells": [{"title": "服务乙"}]},
            ],
        },
        # cells[].highlight=True 验证 ArchCell forbid 承载 (历史 bug: ArchCell 无 highlight 字段, path X 高亮被 forbid 拒)
        ["接入层", "网关甲", "服务层", "服务乙"],
    ),
    "timeline-huawei": (
        {
            "visual_type": "timeline-huawei",
            "title": "演进时间线",
            "phases": [
                {"year": "2024", "label": "起步甲"},
                {"year": "2025", "label": "发展乙"},
                {"year": "2026", "label": "成熟丙"},
            ],
        },
        ["2024", "起步甲", "2025", "成熟丙"],
    ),
    "process-flow-huawei": (
        {
            "visual_type": "process-flow-huawei",
            "title": "处理流程",
            "steps": [
                {"label": "步骤甲"},
                {"label": "步骤乙"},
                {"label": "步骤丙"},
            ],
        },
        ["步骤甲", "步骤乙", "步骤丙"],
    ),
    "swot": (
        {
            "visual_type": "swot",
            "title": "SWOT分析",
            "quadrants": {
                "s": {"heading": "优势", "items": ["优势项甲", "优势项乙"]},
                "w": {"heading": "劣势", "items": ["劣势项丙"]},
                "o": {"heading": "机会", "items": ["机会项丁"]},
                "t": {"heading": "威胁", "items": ["威胁项戊"]},
            },
        },
        ["优势项甲", "劣势项丙", "机会项丁", "威胁项戊"],
    ),
    "roadmap": (
        {
            "visual_type": "roadmap",
            "title": "实施路线图",
            "phases": [
                {"code": "P1", "title": "阶段一"},
                {"code": "P2", "title": "阶段二"},
                {"code": "P3", "title": "阶段三"},
            ],
            "rows": [
                {
                    "label": "泳道甲",
                    "cells": [{"title": "格甲一"}, {"title": "格甲二"}, {"title": "格甲三"}],
                },
                {
                    "label": "泳道乙",
                    "cells": [
                        {"title": "格乙一", "emphasis": "bar"},
                        {"title": "格乙二"},
                        {"title": "格乙三"},
                    ],
                },
            ],
        },
        ["阶段一", "泳道甲", "格甲一", "泳道乙", "格乙一"],
    ),
    "pyramid": (
        {
            "visual_type": "pyramid",
            "title": "能力金字塔",
            "levels": [{"name": "顶层甲"}, {"name": "中层乙"}, {"name": "底层丙"}],
            "descriptions": [
                {"number": "1", "heading": "描述甲"},
                {"number": "2", "heading": "描述乙"},
                {"number": "3", "heading": "描述丙"},
            ],
        },
        ["顶层甲", "底层丙", "描述甲"],
    ),
    "heatmap-matrix": (
        {
            "visual_type": "heatmap-matrix",
            "title": "热力矩阵",
            "columns": ["列甲", "列乙", "列丙"],
            "rows": [
                {"label": "行甲", "scores": [3, 4, 5]},
                {"label": "行乙", "scores": [1, 2, 3]},
            ],
        },
        ["列甲", "列丙", "行甲", "行乙"],
    ),
    "thankyou": (
        {
            "visual_type": "thankyou",
            "eyebrow": "结语EB",
            "title": "致谢标题词",
            "emphasis": "强调结语词",
            "subtitle": "致谢副标题XY",
            "contacts": [{"label": "邮箱", "value": "联系值ZZ"}],
        },
        ["结语EB", "强调结语词", "致谢副标题XY", "联系值ZZ"],
    ),
    "cards-6": (
        {
            "visual_type": "cards-6",
            "title": "六大要点",
            "cards": [
                {"heading": "卡甲", "body": "内容甲"},
                {"heading": "卡乙", "body": "内容乙"},
                {"heading": "卡丙", "body": "内容丙"},
                {"heading": "卡丁", "body": "内容丁"},
                {"heading": "卡戊", "body": "内容戊"},
                {"heading": "卡己", "body": "内容己"},
            ],
        },
        ["卡甲", "卡乙", "卡丙", "卡丁", "卡戊", "卡己"],
    ),
    "rings": (
        {
            "visual_type": "rings",
            "title": "同心环",
            "rings": [{"label": "环甲"}, {"label": "环乙"}],
            "steps": [
                {"number": "01", "heading": "步骤标题甲", "body": "步骤体甲"},
                {"number": "02", "heading": "步骤标题乙"},
            ],
        },
        ["环甲", "步骤标题甲", "步骤标题乙"],
    ),
    "personas": (
        {
            "visual_type": "personas",
            "title": "角色画像",
            "personas": [
                {
                    "role": "角色甲",
                    "name": "姓名A",
                    "attrs": [{"label": "年龄", "value": "30"}],
                    "quote": "引言甲",
                },
                {"role": "角色乙", "name": "姓名B"},
            ],
        },
        # "年龄  30" = renderer 的 f"{label}  {value}" 正确形式; 若 attr(MetaItem)
        # 误走 str() 泄漏则是 "label='年龄' value='30'", 不含此子串 → 精确区分。
        ["角色甲", "姓名A", "年龄  30", "角色乙"],
    ),
    "risk-list": (
        {
            "visual_type": "risk-list",
            "title": "风险清单",
            "risks": [
                {"number": "R1", "title": "风险甲", "desc": "描述甲", "severity": "HIGH"},
                {"number": "R2", "title": "风险乙", "desc": "描述乙", "severity": "LOW"},
            ],
        },
        ["风险甲", "风险乙"],
    ),
    "governance": (
        {
            "visual_type": "governance",
            "title": "治理结构",
            "top": {"name": "决策委员会X"},
            "units": [
                {"code": "U1", "name": "工作组甲", "items": ["职责甲"]},
                {"code": "U2", "name": "工作组乙", "items": ["职责乙"]},
            ],
        },
        ["决策委员会X", "工作组甲", "工作组乙"],
    ),
}


def _build_plan(visual_type: str, variant: dict) -> dict:
    """构造一份仅含 1 张 slide 的 path X plan: data 全在 variant, content 只留必填 title。"""
    return {
        "meta": {"title": f"{visual_type} path-X 回归", "theme": "huawei"},
        "narrative": {"thesis": "path-X 回归", "arc": "smoke", "total_slides": 1},
        "slides": [
            {
                "id": 1,
                "role": "content",
                "visual_type": visual_type,
                # content 只放必填 title (故意与 variant 文本不同, 确保断言命中的是 variant 路径)
                "content": {"title": f"__CONTENT_TITLE_{visual_type}__"},
                "variant": variant,
            }
        ],
    }


def _gather_text(prs) -> str:
    """收集 PPTX 所有可见文本: 文本框 paragraphs + 表格单元格。"""
    chunks: list[str] = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    if p.text:
                        chunks.append(p.text)
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            chunks.append(cell.text)
    return "\n".join(chunks)


@pytest.fixture(scope="module")
def huawei_theme():
    return load_theme("huawei")


@pytest.mark.parametrize("visual_type", list(PATH_X_CASES.keys()), ids=list(PATH_X_CASES.keys()))
def test_path_x_variant_content_rendered(tmp_path, huawei_theme, visual_type):
    """path X variant 的专属内容必须全部出现在渲染输出里 (回归: 历史上多个版式静默丢内容)。"""
    from pptx import Presentation

    variant, expected_texts = PATH_X_CASES[visual_type]
    plan_data = _build_plan(visual_type, variant)
    plan = SlidePlan.model_validate(plan_data)

    # 确认确实走 path X: variant 字段解析成了 typed 子模型
    assert plan.slides[0].variant is not None, f"{visual_type} variant 未解析 (path X 前提失败)"

    out = tmp_path / f"{visual_type}-pathx.pptx"
    render_presentation(plan, huawei_theme, str(out))
    assert out.exists()

    prs = Presentation(str(out))
    all_text = _gather_text(prs)

    missing = [t for t in expected_texts if t not in all_text]
    assert not missing, (
        f"{visual_type} (path X) 渲染丢失 variant 内容: 缺 {missing}\n"
        f"期望全部出现: {expected_texts}\n"
        f"实际渲染文本:\n{all_text}"
    )

    # repr 泄漏防护: renderer 误对 pydantic 子模型做 str() 会渲染出 `label='x' field=None` 之类的
    # repr 文本 (rings 历史 bug)。这些 signature 不应出现在干净的演示文本里。
    repr_leaks = [
        sig for sig in ("=None", "visual_type=", "sublabel=", "style='outline'",
                        # MetaItem/Contact 等 pydantic 子模型被误 str() 时的 __str__ 片段
                        # (如 personas attrs: "label='年龄' value='30'"), 干净文本绝不含
                        "label='", "value='")
        if sig in all_text
    ]
    assert not repr_leaks, (
        f"{visual_type} (path X) 渲染泄漏 pydantic repr 片段 {repr_leaks} — "
        f"renderer 对 typed 子模型误用 str() 而非字段取值\n实际渲染文本:\n{all_text}"
    )


# ---------------------------------------------------------------------------
# 视觉特征回归断言 (S4 T12 从已删 test_variants_render_smoke 迁移, 改 path X 构造)
# path_x_render 主测文本内容; 这 3 条补几何/fill/字号区间的视觉特征覆盖,
# 防红条/深底/巨字回退 (各有历史 issue 来源, 见 docstring)。
# ---------------------------------------------------------------------------

def _render_path_x_first_slide(tmp_path, theme, visual_type: str):
    """用 PATH_X_CASES 的 path X 数据渲染单页, 返回首张 slide 对象。"""
    from pptx import Presentation

    variant = PATH_X_CASES[visual_type][0]
    plan = SlidePlan.model_validate(_build_plan(visual_type, variant))
    out = tmp_path / f"{visual_type}-visual.pptx"
    render_presentation(plan, theme, str(out))
    return Presentation(str(out)).slides[0]


def test_visual_cover_no_edge_red_bar(tmp_path, huawei_theme):
    """cover-left-bar 不应有贯穿边缘的细红条 (2026-05-31 全局规则: huawei 去红条, 用户选型 A)。"""
    from pptx.util import Emu
    from pptx.dml.color import RGBColor

    slide = _render_path_x_first_slide(tmp_path, huawei_theme, "cover-left-bar")
    red = RGBColor(0xC7, 0x00, 0x0B)
    for sh in slide.shapes:
        if sh.left is None or sh.width is None or sh.height is None:
            continue
        # 贯穿边缘红条 = 贴左边 (x≈0) + 窄 (<0.5") + 高 (≥5")
        edge_bar = (sh.left <= Emu(int(0.2 * 914400))
                    and sh.width < Emu(int(0.5 * 914400))
                    and sh.height >= Emu(int(5.0 * 914400)))
        if edge_bar:
            try:
                assert sh.fill.fore_color.rgb != red, "检测到贯穿边缘红条 — 违反 huawei 去红条规则"
            except (AttributeError, TypeError):
                pass


def test_visual_section_divider_light_bg(tmp_path, huawei_theme):
    """section-divider-dark 背景 = #FFFFFF (2026-05-31 浅底选型 B; 旧深底违反 P1/P3 用户评 0)。"""
    from pptx.dml.color import RGBColor

    slide = _render_path_x_first_slide(tmp_path, huawei_theme, "section-divider-dark")
    try:
        rgb = slide.background.fill.fore_color.rgb
    except Exception as e:
        pytest.skip(f"slide.background 不可读 (实现未设 solid fill): {e}")
    assert rgb == RGBColor(0xFF, 0xFF, 0xFF), f"背景期望浅底 #FFFFFF, 实际 #{rgb}"


def test_visual_thankyou_main_text_restrained(tmp_path, huawei_theme):
    """thankyou 标题字号收敛区间 [44,64]pt (2026-06-12 收敛 AP3 巨字, 旧 ≥72pt 占画面 1/3+)。"""
    slide = _render_path_x_first_slide(tmp_path, huawei_theme, "thankyou")
    sizes = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            if p.font.size is not None:
                sizes.append(p.font.size.pt)
            for r in p.runs:
                if r.font.size is not None:
                    sizes.append(r.font.size.pt)
    assert sizes, "thankyou slide 无任何 paragraph/run 带 font.size"
    assert 44 <= max(sizes) <= 64, (
        f"thankyou 标题最大字号 {max(sizes)} pt 不在克制区间 [44, 64]"
    )


def test_visual_arch_cell_highlight_primary_fill(tmp_path, huawei_theme):
    """architecture-layered 的 cell highlight=True 应渲染 primary 红底 (#C7000B)。

    回归 (2026-06-22 misc §5): ArchCell 历史无 highlight 字段, extra=forbid 下 path X
    cell 级高亮数据被拒, 高亮能力不可达 (非 never-fired, 是模型无法承载)。schema 加
    highlight 字段后, path X highlight cell 必须真出 primary 红底。
    """
    from pptx.dml.color import RGBColor

    slide = _render_path_x_first_slide(tmp_path, huawei_theme, "architecture-layered")
    primary = RGBColor(0xC7, 0x00, 0x0B)  # huawei tokens.colors.primary
    found = False
    for sh in slide.shapes:
        try:
            if sh.fill.fore_color.rgb == primary:
                found = True
                break
        except (AttributeError, TypeError):
            continue
    assert found, (
        "architecture-layered highlight cell 未渲染 primary 红底 — "
        "ArchCell.highlight 未被消费或 schema 未承载该字段"
    )
