# /// script
# requires-python = ">=3.11"
# dependencies = ["pytest>=7.0", "pyyaml>=6.0", "python-pptx>=1.0.0", "pydantic>=2.0", "lxml>=4.9", "Pillow>=10.0"]
# ///
"""S4 misc renderer 修复回归 (T14, fix-ppt-misc-medium-low §1).

覆盖 Dev-1 在 T14 交付的 4 个 renderer 行为点中 **T12-safe(走 path X typed variant)** 的 3 条:
- heatmap 0.5 档半块字符 (heatmap_matrix.py)
- swot 字母字号从 layouts.yaml 读 (swot.py + layouts.yaml)
- timeline desc/items 高度 safe_h guard (timeline_huawei.py + renderer_kit.safe_h)

第 4 条 (roadmap bars start/span guard) **未纳入**: 该 guard 只为 path Y 的 lanes/bars
legacy 结构兜底, typed variant RoadmapRow/RoadmapCell 无 start/span 字段。一旦 T12 把
SlideContent.extra='forbid', path Y lanes 被拒, 该 guard 无合法触发路径 (可能成死代码)。
处理方式待 team-lead + Dev-1 裁决 (见 QA 上报); 裁决后再决定是否补 test / 删 guard。

heatmap 半块字符: U+2588 (FULL BLOCK '█') / U+258C (LEFT HALF BLOCK '▌')。
为防源码全角/编码歧义, 字符用 \\uXXXX 转义常量。
"""
from __future__ import annotations

import sys
from pathlib import Path

PLUGIN_ROOT = Path(__file__).resolve().parent.parent
if str(PLUGIN_ROOT) not in sys.path:
    sys.path.insert(0, str(PLUGIN_ROOT))

FULL_BLOCK = "█"   # █
HALF_BLOCK = "▌"   # ▌

_META = {"title": "t", "preset": "research-report", "theme": "huawei",
         "generated_at": "2026-06-18T00:00:00"}
_NARR = {"thesis": "x", "arc": "a", "total_slides": 1}


def _render(plan_data: dict, out_path: Path):
    from engine.render import load_theme, render_presentation
    from schemas.slide_plan import SlidePlan
    from pptx import Presentation
    plan = SlidePlan.model_validate(plan_data)
    render_presentation(plan, load_theme("huawei"), str(out_path))
    return Presentation(str(out_path))


# ============================================================
# 1. heatmap 0.5 档半块字符 (path X: HeatmapRow.scores 含小数)
# ============================================================

def _heatmap_cell_text(prs) -> str:
    """heatmap 用 table 渲染, cell 文本在 table.rows[].cells[].text (不在 text_frame)。"""
    out = []
    for shape in prs.slides[0].shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    out.append(cell.text)
    return " ".join(out)


def test_heatmap_half_block_for_fractional_score(tmp_path):
    """score=3.5 → 3 full block + 1 half block; score=0.5 → 1 half block。"""
    plan = {"meta": _META, "narrative": _NARR, "slides": [{
        "id": 1, "role": "content", "visual_type": "heatmap-matrix",
        "content": {"title": "H"},
        "variant": {"visual_type": "heatmap-matrix", "title": "H",
                    "columns": ["A", "B", "C"],
                    "rows": [{"label": "r1", "scores": [3.5, 5, 0]},
                             {"label": "r2", "scores": [0.5, 2, 4]}]},
    }]}
    prs = _render(plan, tmp_path / "heatmap.pptx")
    blob = _heatmap_cell_text(prs)
    assert HALF_BLOCK in blob, f"0.5 档半块 U+258C 未出现; cells={blob!r}"
    assert (FULL_BLOCK * 3 + HALF_BLOCK) in blob, f"3.5 档应 3 full + 半块; cells={blob!r}"


# ============================================================
# 2. swot 字母字号读 layouts.yaml (path X: SwotQuadrants)
# ============================================================

def test_swot_letter_size_from_layouts(tmp_path):
    """象限字母 S/W/O/T 字号 = layouts.swot.letter_size_px(42) * _px_pt(0.6) = 25.2pt。"""
    plan = {"meta": _META, "narrative": _NARR, "slides": [{
        "id": 1, "role": "content", "visual_type": "swot",
        "content": {"title": "SWOT分析"},
        "variant": {"visual_type": "swot", "title": "SWOT分析", "quadrants": {
            "s": {"heading": "优势", "items": ["a"]},
            "w": {"heading": "劣势", "items": ["b"]},
            "o": {"heading": "机会", "items": ["c"]},
            "t": {"heading": "威胁", "items": ["d"]},
        }},
    }]}
    prs = _render(plan, tmp_path / "swot.pptx")
    sizes = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and shape.text_frame.text in ("S", "W", "O", "T"):
            sz = shape.text_frame.paragraphs[0].font.size
            if sz is not None:
                sizes.append(round(sz.pt, 2))
    assert len(sizes) == 4, f"应找到 4 个象限字母 shape, 实际 {len(sizes)}: {sizes}"
    assert all(abs(p - 25.2) < 0.01 for p in sizes), f"字母字号应 25.2pt, 实际 {sizes}"


# ============================================================
# 3. timeline desc/items 高度 safe_h guard (path X: TimelinePhase.items 超长)
# ============================================================

def test_timeline_long_items_no_negative_inches(tmp_path):
    """超长 items 压低剩余高度时 safe_h 兜底, render 不抛 Inches(负值) 异常。

    TimelinePhase typed 字段为 year/label/items(无 desc, extra=forbid), 故用 items 塞超长文本
    复现高度压缩场景 — 与 path Y 的 desc 超长等效触发 safe_h guard, 且 T12-safe。
    """
    long = "压力测试超长正文" * 30
    plan = {"meta": _META, "narrative": _NARR, "slides": [{
        "id": 1, "role": "content", "visual_type": "timeline-huawei",
        "content": {"title": "T"},
        "variant": {"visual_type": "timeline-huawei", "title": "T", "phases": [
            {"year": "2024", "label": "p1", "items": [long, long]},
            {"year": "2025", "label": "p2", "items": [long]},
            {"year": "2026", "label": "p3", "items": [long]},
        ]},
    }]}
    # 不抛异常即通过 (旧代码 Inches(负值) 会在 add shape 时报错)
    _render(plan, tmp_path / "timeline.pptx")
