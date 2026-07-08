# /// script
# requires-python = ">=3.11"
# dependencies = ["pytest>=7.0", "pyyaml>=6.0", "python-pptx>=1.0.0", "pydantic>=2.0"]
# ///
"""域 C — us-military-rotation 真实 E2E (S3-P0, test-specs TC-C1..C7).

PRD §11.5 核心载体: 第一个真实 E2E fixture 不再 skip 且 PASS。

走 `replay` 全链路 (subprocess, 与真实使用同路径), 断言:
- TC-C1: exit 0, PPTX 文件生成
- TC-C2: slide-plan 与 recorded 结构等价
- TC-C3: PPTX 页数 == expected/assertions.yaml
- TC-C4: visual_type 序列 == expected
- TC-C5: huawei 应用率 ≥ 基线阈值
- TC-C6: 无变量/模板泄露 (forbidden_substrings)
- TC-C7: E2E 内联敏感词零命中 (PPTX 文本 + fixture 三层)

断言基准全部读 expected/assertions.yaml (与 Dev 同真相源, 不硬编码序列)。

fixture/orchestrator 未就位时 skipif 守护, Dev 任务 3+9 交付后激活。
不再 skip 是硬目标 (退出准则 §5.2)。
"""
from __future__ import annotations

import subprocess
from pathlib import Path

import pytest

PLUGIN_ROOT = Path(__file__).resolve().parents[2]
ORCHESTRATOR = PLUGIN_ROOT / "engine" / "orchestrator.py"
FIXTURE_DIR = PLUGIN_ROOT / "tests" / "fixtures" / "us-military-rotation"
FIXTURE_INPUT = FIXTURE_DIR / "input"
FIXTURE_RECORDED = FIXTURE_DIR / "recorded"
FIXTURE_EXPECTED = FIXTURE_DIR / "expected"
ASSERTIONS_YAML = FIXTURE_EXPECTED / "assertions.yaml"
RECORDED_PLAN = FIXTURE_RECORDED / "plan-output.yaml"

# 脱敏扫描器路径 (TC-C7)
SCANNER_PATH = PLUGIN_ROOT / "tests" / "pixel" / "assert_no_sensitive.py"

ORCH_READY = ORCHESTRATOR.exists()
FIXTURE_READY = (
    FIXTURE_INPUT.exists()
    and RECORDED_PLAN.exists()
    and ASSERTIONS_YAML.exists()
)

pytestmark = [
    pytest.mark.e2e,
    pytest.mark.skipif(
        not ORCH_READY,
        reason=f"engine/orchestrator.py 尚未就绪 (Dev 任务 3 前): {ORCHESTRATOR}",
    ),
    pytest.mark.skipif(
        not FIXTURE_READY,
        reason=f"us-military fixture 尚未就绪 (Dev 任务 9 前): {FIXTURE_DIR}",
    ),
]


# ============================================================
# 公共辅助
# ============================================================

def _load_yaml(path: Path) -> dict:
    import yaml

    return yaml.safe_load(path.read_text(encoding="utf-8"))


def _slide_visual_types(plan_data: dict) -> list[str]:
    slides = plan_data.get("slides", []) if isinstance(plan_data, dict) else []
    return [s.get("visual_type") for s in slides]


def _slide_roles(plan_data: dict) -> list:
    slides = plan_data.get("slides", []) if isinstance(plan_data, dict) else []
    return [s.get("role") for s in slides]


def _extract_pptx_text(pptx_path: Path) -> str:
    """提取 PPTX 全部可见文本 (拼成一个大字符串, 供泄露/敏感词检查)。"""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    chunks.append(para.text)
    return "\n".join(chunks)


def _load_scanner():
    import importlib.util
    import sys

    if str(SCANNER_PATH.parent) not in sys.path:
        sys.path.insert(0, str(SCANNER_PATH.parent))
    spec = importlib.util.spec_from_file_location("assert_no_sensitive", SCANNER_PATH)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


@pytest.fixture(scope="module")
def replay_result(tmp_path_factory):
    """跑一次 replay 全链路 (module scope, 多个断言复用同一产物)。

    返回 (CompletedProcess, output_dir, pptx_path, run_dir)。
    """
    tmp = tmp_path_factory.mktemp("e2e_us_military")
    out = tmp / "deck.pptx"
    cmd = [
        "uv", "run", "--script", str(ORCHESTRATOR),
        "replay", str(FIXTURE_INPUT),
        "--fixture", str(FIXTURE_DIR),
        "--theme", "huawei",
        "--output", str(out),
    ]
    result = subprocess.run(
        cmd, cwd=str(PLUGIN_ROOT),
        capture_output=True, text=True, encoding="utf-8", errors="replace",
        timeout=300,
    )
    # 定位 run 目录
    runs = tmp / ".ppt-workdir" / "runs"
    run_dir = None
    if runs.is_dir():
        dirs = sorted(p for p in runs.iterdir() if p.is_dir())
        run_dir = dirs[-1] if dirs else None
    return result, tmp, out, run_dir


@pytest.fixture(scope="module")
def assertions():
    return _load_yaml(ASSERTIONS_YAML)


# ============================================================
# TC-C1: replay 全链路 exit 0 + PPTX 生成
# ============================================================

def test_replay_exits_0(replay_result):
    result, _, _, _ = replay_result
    assert result.returncode == 0, (
        f"replay 全链路应 exit 0, 实际 {result.returncode}\n"
        f"stdout={result.stdout}\nstderr={result.stderr}"
    )


def test_replay_produces_pptx(replay_result):
    _, _, pptx, _ = replay_result
    assert pptx.exists(), f"PPTX 未生成: {pptx}"
    assert pptx.stat().st_size > 10_000, (
        f"PPTX 体积 {pptx.stat().st_size} B 过小 (< 10 KB)"
    )


# ============================================================
# TC-C2: slide-plan 与 recorded 结构等价
# ============================================================

def test_slide_plan_matches_recorded(replay_result):
    """产出 03-plan/output.yaml 与 fixture recorded/plan-output.yaml 结构等价。"""
    _, _, _, run_dir = replay_result
    assert run_dir is not None, "未定位到 run 目录"
    produced = _load_yaml(run_dir / "03-plan" / "output.yaml")
    recorded = _load_yaml(RECORDED_PLAN)
    assert _slide_visual_types(produced) == _slide_visual_types(recorded), (
        "产出 slide-plan 的 visual_type 序列与 recorded 不一致"
    )
    assert _slide_roles(produced) == _slide_roles(recorded), (
        "产出 slide-plan 的 role 序列与 recorded 不一致"
    )


# ============================================================
# TC-C3: PPTX 页数 == expected
# ============================================================

def test_pptx_slide_count(replay_result, assertions):
    from pptx import Presentation

    _, _, pptx, _ = replay_result
    expected = assertions["expected_slide_count"]
    prs = Presentation(str(pptx))
    assert len(prs.slides) == expected, (
        f"PPTX 页数 {len(prs.slides)} != expected {expected}"
    )


# ============================================================
# TC-C4: visual_type 序列 == expected
# ============================================================

def test_visual_type_sequence(replay_result, assertions):
    _, _, _, run_dir = replay_result
    produced = _load_yaml(run_dir / "03-plan" / "output.yaml")
    expected_seq = assertions["expected_visual_type_sequence"]
    assert _slide_visual_types(produced) == expected_seq, (
        f"visual_type 序列与 expected 不一致\n"
        f"产出={_slide_visual_types(produced)}\nexpected={expected_seq}"
    )


# ============================================================
# TC-C5: huawei 应用率 ≥ 基线阈值
# ============================================================

def test_huawei_application_rate(replay_result, assertions):
    from schemas.variants import VARIANT_TYPES

    _, _, _, run_dir = replay_result
    produced = _load_yaml(run_dir / "03-plan" / "output.yaml")
    vts = _slide_visual_types(produced)
    total = len(vts)
    assert total > 0, "产出 plan 无 slides"
    huawei_count = sum(1 for vt in vts if vt in VARIANT_TYPES)
    ratio = huawei_count / total
    min_rate = assertions["min_application_rate"]
    assert ratio >= min_rate, (
        f"huawei 应用率 {ratio:.2%} < 基线 {min_rate:.2%} "
        f"(huawei={huawei_count}/{total})"
    )


# ============================================================
# TC-C6: 无变量/模板泄露
# ============================================================

def test_no_template_leakage(replay_result, assertions):
    """slide-plan.yaml + PPTX 文本检无 forbidden_substrings 泄露字样。"""
    _, _, pptx, run_dir = replay_result
    forbidden = assertions["forbidden_substrings"]
    plan_text = (run_dir / "03-plan" / "output.yaml").read_text(encoding="utf-8")
    pptx_text = _extract_pptx_text(pptx)
    combined = plan_text + "\n" + pptx_text
    leaked = [s for s in forbidden if s in combined]
    assert not leaked, f"检出模板泄露字样: {leaked}"


# ============================================================
# TC-C7: E2E 内联敏感词零命中 (PPTX 文本 + fixture 三层)
# ============================================================

def test_no_sensitive_in_pptx_and_fixture(replay_result):
    """PPTX 提取文本 + fixture 三层目录过 K01-K22 == 0。"""
    if not SCANNER_PATH.exists():
        pytest.skip(f"脱敏扫描器未就绪: {SCANNER_PATH}")
    _, out_dir, pptx, _ = replay_result
    scanner = _load_scanner()

    # fixture 三层目录扫描 (input/recorded/expected)
    fixture_hits = scanner.scan(FIXTURE_DIR)
    assert fixture_hits == 0, (
        f"fixture 目录含 {fixture_hits} 条敏感词命中 (零容忍)"
    )

    # PPTX 提取文本写临时 .txt 后扫 (PPTX 二进制内文本走 K01-K22)
    pptx_text = _extract_pptx_text(pptx)
    scratch = out_dir / "pptx-text-scan"
    scratch.mkdir(exist_ok=True)
    (scratch / "deck-text.txt").write_text(pptx_text, encoding="utf-8")
    pptx_hits = scanner.scan(scratch)
    assert pptx_hits == 0, (
        f"PPTX 提取文本含 {pptx_hits} 条敏感词命中 (零容忍)"
    )
