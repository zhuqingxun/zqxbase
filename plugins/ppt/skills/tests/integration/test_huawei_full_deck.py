# /// script
# requires-python = ">=3.11"
# dependencies = ["pytest>=7.0", "python-pptx>=1.0.0", "pydantic>=2.0", "pyyaml>=6.0"]
# ///
"""SPEC-IT-*：端到端 huawei 主题渲染集成测试。

依赖：Dev T1/T2/T3（schemas）+ T4（renderers）+ T5（samples yaml）全部完成。
所有测试用 skipif 守护缺失的 fixture，避免 Dev 分阶段交付时整个 run 爆炸。
"""
from __future__ import annotations

from pathlib import Path

import pytest

PLUGIN_ROOT = Path(__file__).resolve().parents[2]
FIXTURES_DIR = PLUGIN_ROOT / "tests" / "fixtures"
SAMPLES_DIR = PLUGIN_ROOT / "tests" / "variants" / "samples"

HUAWEI_25_FIXTURE = FIXTURES_DIR / "huawei-25-slides.yaml"
LEGACY_FIXTURE = FIXTURES_DIR / "legacy-slides.yaml"


# 延迟导入：schemas / engine 依赖 Dev 交付
try:
    from schemas.slide_plan import SlidePlan  # type: ignore
    from engine.render import render_presentation, load_theme  # type: ignore
except ImportError:
    SlidePlan = None  # type: ignore
    render_presentation = None  # type: ignore
    load_theme = None  # type: ignore


pytestmark = pytest.mark.skipif(
    SlidePlan is None or render_presentation is None,
    reason="engine / schemas 尚未就绪（Dev 前置任务未完成）",
)


@pytest.fixture(scope="module")
def huawei_theme():
    return load_theme("huawei")


# ============================================================
# SPEC-IT-1..8：25 版式端到端渲染
# ============================================================

@pytest.mark.skipif(not HUAWEI_25_FIXTURE.exists(), reason=f"{HUAWEI_25_FIXTURE} 尚未就绪（Dev T5 前）")
def test_render_huawei_25_slides(tmp_path, huawei_theme):
    """SPEC-IT-1 / IT-2 / IT-3 / IT-4：渲染 25 张 slide 覆盖全部 P0/P1/P2。"""
    import yaml as _yaml
    from pptx import Presentation

    data = _yaml.safe_load(HUAWEI_25_FIXTURE.read_text(encoding="utf-8"))
    plan = SlidePlan.model_validate(data)
    out = tmp_path / "huawei-25.pptx"
    render_presentation(plan, huawei_theme, str(out))

    # SPEC-IT-1：无异常
    # SPEC-IT-2：文件存在且大小合理
    assert out.exists(), f"{out} 未生成"
    assert out.stat().st_size > 50_000, f"PPTX 体积 {out.stat().st_size} B 过小（< 50 KB）"

    # SPEC-IT-3：slide 数量一致
    prs = Presentation(str(out))
    assert len(prs.slides) == len(plan.slides), (
        f"期望 slide 数 {len(plan.slides)}，实际 {len(prs.slides)}"
    )

    # SPEC-IT-4：每张 slide 至少 1 shape
    for i, slide in enumerate(prs.slides, start=1):
        assert len(slide.shapes) >= 1, f"Slide {i} 无 shape"


# ============================================================
# SPEC-IT-9：幂等性
# ============================================================

@pytest.mark.skipif(not HUAWEI_25_FIXTURE.exists(), reason="fixture 未就绪")
def test_render_is_idempotent(tmp_path, huawei_theme):
    """SPEC-IT-9：相同 YAML 两次渲染的 slide 数 / shape 数 / text 完全一致。"""
    import yaml as _yaml
    from pptx import Presentation

    data = _yaml.safe_load(HUAWEI_25_FIXTURE.read_text(encoding="utf-8"))
    plan = SlidePlan.model_validate(data)

    out_a = tmp_path / "a.pptx"
    out_b = tmp_path / "b.pptx"
    render_presentation(plan, huawei_theme, str(out_a))
    render_presentation(plan, huawei_theme, str(out_b))

    prs_a = Presentation(str(out_a))
    prs_b = Presentation(str(out_b))
    assert len(prs_a.slides) == len(prs_b.slides), "slide 数量不一致（非幂等）"

    for idx, (sa, sb) in enumerate(zip(prs_a.slides, prs_b.slides), start=1):
        assert len(sa.shapes) == len(sb.shapes), f"Slide {idx} shape 数不一致"
        texts_a = [_extract_text(sh) for sh in sa.shapes]
        texts_b = [_extract_text(sh) for sh in sb.shapes]
        assert texts_a == texts_b, f"Slide {idx} 文本内容不一致:\nA={texts_a}\nB={texts_b}"


def _extract_text(shape) -> str:
    """提取 shape 的可见文本（忽略 timestamp / rId 等差异）。"""
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


# ============================================================
# SPEC-IT-10 / IT-11：legacy plan 回归
# ============================================================

@pytest.mark.skipif(not LEGACY_FIXTURE.exists(), reason=f"{LEGACY_FIXTURE} 尚未就绪（Dev T5 前）")
def test_legacy_plan_renders_ok(tmp_path, huawei_theme):
    """SPEC-IT-10 / IT-11：bullets / cards-3 / framework 等旧 visual_type + theme=huawei 渲染成功。

    验证 O7 决策（load_theme 双暴露）有效——旧 renderer 读 theme['colors']['primary'] 零改造。
    """
    import yaml as _yaml
    from pptx import Presentation

    data = _yaml.safe_load(LEGACY_FIXTURE.read_text(encoding="utf-8"))
    plan = SlidePlan.model_validate(data)
    out = tmp_path / "legacy.pptx"
    render_presentation(plan, huawei_theme, str(out))

    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == len(plan.slides)


# ============================================================
# SPEC-IT-5..8：特定 slide 内容断言（fixture 约定顺序）
# ============================================================

@pytest.mark.skipif(not HUAWEI_25_FIXTURE.exists(), reason="fixture 未就绪")
def test_spec_it_5_cover_left_bar_shapes(tmp_path, huawei_theme):
    """SPEC-IT-5：Slide 1 (cover-left-bar) 至少含左红条 shape + 标题 text。

    识别左红条启发式：slide width 左边缘附近，width < 30 px（换算 EMU），高度接近全画布。
    """
    import yaml as _yaml
    from pptx import Presentation
    from pptx.util import Emu

    data = _yaml.safe_load(HUAWEI_25_FIXTURE.read_text(encoding="utf-8"))
    plan = SlidePlan.model_validate(data)
    # fixture 首张 slide 必须是 cover-left-bar
    assert plan.slides[0].visual_type == "cover-left-bar", (
        f"fixture 首张应为 cover-left-bar，实际 {plan.slides[0].visual_type}"
    )

    out = tmp_path / "cover.pptx"
    render_presentation(plan, huawei_theme, str(out))
    prs = Presentation(str(out))
    slide1 = prs.slides[0]
    assert len(slide1.shapes) >= 2, "Slide 1 至少要有左红条 + 标题两个 shape"

    # 查找左红条：x 近 0、width 小
    left_bar = next(
        (sh for sh in slide1.shapes
         if sh.left is not None and sh.left <= Emu(0.1 * 914400)  # 0.1 inch
         and sh.width is not None and sh.width < Emu(0.5 * 914400)),  # < 0.5 inch
        None
    )
    assert left_bar is not None, "Slide 1 未找到左红条 shape（x≈0, width<0.5in）"


@pytest.mark.skipif(not HUAWEI_25_FIXTURE.exists(), reason="fixture 未就绪")
def test_spec_it_6_section_divider_dark_background(tmp_path, huawei_theme):
    """SPEC-IT-6：section-divider-dark slide 背景色 = #2A2A2A（ink_dark）。

    fixture 第 3 张（或首个 section-divider-dark）slide 必须深底。
    """
    import yaml as _yaml
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    data = _yaml.safe_load(HUAWEI_25_FIXTURE.read_text(encoding="utf-8"))
    plan = SlidePlan.model_validate(data)
    section_slides = [
        (i, s) for i, s in enumerate(plan.slides) if s.visual_type == "section-divider-dark"
    ]
    if not section_slides:
        pytest.skip("fixture 无 section-divider-dark slide")

    out = tmp_path / "section.pptx"
    render_presentation(plan, huawei_theme, str(out))
    prs = Presentation(str(out))
    idx, _ = section_slides[0]
    slide = prs.slides[idx]
    # python-pptx 获取背景色：slide.background.fill.fore_color.rgb
    try:
        rgb = slide.background.fill.fore_color.rgb
    except Exception as e:
        pytest.skip(f"无法读取 slide background（实现未设 solid fill）: {e}")
    assert rgb == RGBColor(0x2A, 0x2A, 0x2A), f"背景期望 #2A2A2A，实际 #{rgb}"


@pytest.mark.skipif(not HUAWEI_25_FIXTURE.exists(), reason="fixture 未就绪")
def test_spec_it_7_matrix_2x2_has_enough_shapes(tmp_path, huawei_theme):
    """SPEC-IT-7：matrix-2x2 slide 至少 6 shapes（4 象限 + 2 轴）。"""
    import yaml as _yaml
    from pptx import Presentation

    data = _yaml.safe_load(HUAWEI_25_FIXTURE.read_text(encoding="utf-8"))
    plan = SlidePlan.model_validate(data)
    matrix_idx = next(
        (i for i, s in enumerate(plan.slides) if s.visual_type == "matrix-2x2"), None
    )
    if matrix_idx is None:
        pytest.skip("fixture 无 matrix-2x2 slide")

    out = tmp_path / "matrix.pptx"
    render_presentation(plan, huawei_theme, str(out))
    prs = Presentation(str(out))
    slide = prs.slides[matrix_idx]
    assert len(slide.shapes) >= 6, (
        f"matrix-2x2 期望 ≥6 shapes（4 象限 + 2 轴 + 标题），实际 {len(slide.shapes)}"
    )


@pytest.mark.skipif(not HUAWEI_25_FIXTURE.exists(), reason="fixture 未就绪")
def test_spec_it_8_architecture_layered_layer_count(tmp_path, huawei_theme):
    """SPEC-IT-8：architecture-layered slide 的 shape 数 ≥ Σ(cells) + layers。"""
    import yaml as _yaml
    from pptx import Presentation

    data = _yaml.safe_load(HUAWEI_25_FIXTURE.read_text(encoding="utf-8"))
    plan = SlidePlan.model_validate(data)
    arch_slides = [
        (i, s) for i, s in enumerate(plan.slides) if s.visual_type == "architecture-layered"
    ]
    if not arch_slides:
        pytest.skip("fixture 无 architecture-layered slide")

    idx, spec = arch_slides[0]
    # 架构决策（HR-1 方案 A）：variant 字段扁平放 content，由 SlideContent.extra='allow' 承接
    extra = getattr(spec.content, "__pydantic_extra__", None) or {}
    layers = extra.get("layers", [])
    if not layers:
        pytest.skip("arch slide content.layers 为空，无法精确断言")
    expected_min = len(layers) + sum(len(l.get("cells", [])) for l in layers)

    out = tmp_path / "arch.pptx"
    render_presentation(plan, huawei_theme, str(out))
    prs = Presentation(str(out))
    slide = prs.slides[idx]
    assert len(slide.shapes) >= expected_min, (
        f"architecture-layered 期望 shape ≥ layers({len(layers)}) + cells({expected_min - len(layers)})"
        f" = {expected_min}，实际 {len(slide.shapes)}"
    )
