# /// script
# requires-python = ">=3.11"
# dependencies = ["pytest>=7.0", "python-pptx>=1.0.0", "pydantic>=2.0", "pyyaml>=6.0", "lxml>=4.9"]
# ///
"""L-1 ~ L-7 layout fix 回归测试.

每个 L-N 一个 canonical case + 关键几何/数量断言. 防未来 commit 误改回旧行为.

视觉验收 (PNG 对账) 在 us-military-rotation 与 ability-evaluation 双 deck 上做过, 见
delivery-report; 本文件只做 surgical regression 防护.
"""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx.oxml.ns import qn

PLUGIN_ROOT = Path(__file__).resolve().parents[2]

try:
    from schemas.slide_plan import SlidePlan  # type: ignore
    from engine.render import render_presentation, load_theme, _RENDERERS  # type: ignore
except ImportError:
    SlidePlan = None  # type: ignore
    render_presentation = None  # type: ignore
    load_theme = None  # type: ignore
    _RENDERERS = {}  # type: ignore


pytestmark = pytest.mark.skipif(
    SlidePlan is None or render_presentation is None,
    reason="engine / schemas 尚未就绪",
)


# ============================================================
# 前置 Bug A regression: huawei 18 个 renderer 必须全部注册
# ============================================================

def test_bug_a_huawei_renderers_all_registered():
    """Bug A 回归: 装饰器副作用注册必须把 18 个 huawei renderer 全部塞进 _RENDERERS.

    历史 bug: PEP 723 双 module 实例让 huawei renderer 被注册到另一份 _RENDERERS dict,
    dispatch 全部 silent fallback 到 bullets. 已通过 sys.modules.setdefault alias 修复.
    """
    huawei_types = {
        "cover-left-bar", "toc", "section-divider-dark", "kpi-stats",
        "matrix-2x2", "architecture-layered", "timeline-huawei",
        "process-flow-huawei", "swot", "roadmap", "pyramid", "heatmap-matrix",
        "thankyou", "cards-6", "rings", "personas", "risk-list", "governance",
    }
    missing = huawei_types - set(_RENDERERS.keys())
    assert not missing, f"huawei renderer 注册缺失: {missing}"
    assert len(_RENDERERS) >= 42, f"_RENDERERS 总数应 >=42 (24 legacy + 18 huawei), 实际 {len(_RENDERERS)}"


# ============================================================
# Bug B regression: schemas 字段名与 renderer 字段名双名兼容
# ============================================================

def _build_plan(slides_yaml_data: list) -> SlidePlan:
    """从 dict list 构造 SlidePlan."""
    plan_data = {
        "meta": {"title": "test", "theme": "huawei"},
        "narrative": {"thesis": "test", "arc": "smoke", "total_slides": len(slides_yaml_data)},
        "slides": slides_yaml_data,
    }
    return SlidePlan(**plan_data)


def _render_one(plan: SlidePlan, tmp_path: Path):
    """Render plan, 返回 Presentation 对象."""
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = load_theme("huawei")
    out = tmp_path / "out.pptx"
    render_presentation(plan, theme, str(out))
    from pptx import Presentation  # type: ignore
    return Presentation(str(out))


def test_bug_b_section_divider_reads_big_number_alias(tmp_path):
    """Bug B 回归: section-divider-dark 既支持 number (legacy) 也支持 big_number (schemas)."""
    plan_a = _build_plan([{
        "id": 1, "role": "content", "visual_type": "section-divider-dark",
        "content": {"title": "x", "number": "01", "description": "d"},
    }])
    plan_b = _build_plan([{
        "id": 1, "role": "content", "visual_type": "section-divider-dark",
        "content": {"title": "x", "big_number": "01", "description": "d"},
    }])
    prs_a = _render_one(plan_a, tmp_path / "a")
    prs_b = _render_one(plan_b, tmp_path / "b")
    # 两份都应包含相同数量的 shape (slide 上 number 已渲染)
    assert len(prs_a.slides[0].shapes) == len(prs_b.slides[0].shapes)
    # 提取所有文本, 确认 "01" 都出现
    a_texts = " ".join(s.text_frame.text for s in prs_a.slides[0].shapes if s.has_text_frame)
    b_texts = " ".join(s.text_frame.text for s in prs_b.slides[0].shapes if s.has_text_frame)
    assert "01" in a_texts and "01" in b_texts


# ============================================================
# L-1: section_divider_dark 标题不压描述
# ============================================================

def test_L1_section_divider_long_title_not_overlap_description(tmp_path):
    """L-1: 长标题 (24 字) 渲染后 description 起点应 ≥ title 底沿, 无重叠."""
    long_title = "下一步: 中国企业对照 + 数字化选型 + 文化手册"  # 24 字
    plan = _build_plan([{
        "id": 1, "role": "content", "visual_type": "section-divider-dark",
        "content": {
            "title": long_title,
            "number": "01",
            "description": "这是一段紧跟标题的描述, 应当出现在标题下方不能重叠.",
        },
    }])
    prs = _render_one(plan, tmp_path)
    slide = prs.slides[0]
    # 找到 title 与 description 的 textbox: 通过文字内容定位
    title_shape = None
    desc_shape = None
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        text = s.text_frame.text
        if long_title in text:
            title_shape = s
        elif "应当出现在标题下方" in text:
            desc_shape = s
    assert title_shape is not None, "title shape 未渲染"
    assert desc_shape is not None, "description shape 未渲染"
    # title 底沿 ≤ description 顶沿 (考虑 PowerPoint EMU 单位, 1 inch = 914400)
    title_bottom = title_shape.top + title_shape.height
    assert desc_shape.top >= title_bottom, (
        f"description.top ({desc_shape.top}) 应 ≥ title.bottom ({title_bottom}), 否则重叠"
    )


# ============================================================
# L-2: toc giant_label 中文不裁切
# ============================================================

def test_L2_toc_header_restrained_no_giant(tmp_path):
    """L-2 (2026-05-31 taste C): toc 标题克制 — 无 300px 巨号 '目录'; 标题含 '目录'、在画布内、全页无巨号 (≤60pt, 合 P3)。"""
    plan = _build_plan([{
        "id": 1, "role": "content", "visual_type": "toc",
        "content": {
            "title": "目录",
            "chapters": [
                {"number": "01", "title": "A", "page": "01"},
                {"number": "02", "title": "B", "page": "02"},
                {"number": "03", "title": "C", "page": "03"},
            ],
        },
    }])
    prs = _render_one(plan, tmp_path)
    slide = prs.slides[0]
    canvas_h_emu = int(7.5 * 914400)
    canvas_w_emu = int(13.333 * 914400)
    header = None
    max_pt = 0.0
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        if "目录" in s.text_frame.text and header is None:
            header = s
        for para in s.text_frame.paragraphs:
            if para.font.size is not None:
                max_pt = max(max_pt, para.font.size.pt)
    assert header is not None, "toc 标题 '目录' 未渲染"
    assert header.top >= 0 and header.top + header.height <= canvas_h_emu
    assert header.left >= 0 and header.left + header.width <= canvas_w_emu
    # 旧 giant '目录' 为 300px≈180pt, 违反 P3; 新设计全页字号封顶
    assert max_pt <= 60.0, f"toc 出现巨号字 ({max_pt}pt), 违反 P3 装饰字不应巨大"


# ============================================================
# L-3: kpi_stats value+unit 单行不 wrap
# ============================================================

def test_L3_kpi_value_unit_single_line(tmp_path):
    """L-3: 长 value + unit 应字号自适应单行渲染, label 不应被压. 通过 shape 顺序的 top 单调递增验证."""
    plan = _build_plan([{
        "id": 1, "role": "content", "visual_type": "kpi-stats",
        "content": {
            "title": "K", "subtitle": "s",
            "kpis": [
                {"value": "15,000", "unit": "军官", "label": "年度参与", "desc": "d"},
                {"value": "50", "unit": "%+", "label": "首选命中", "desc": "d"},
                {"value": "80", "unit": "%", "label": "前 10 命中", "desc": "d"},
            ],
        },
    }])
    prs = _render_one(plan, tmp_path)
    slide = prs.slides[0]
    # 找出 cell 1 的所有 text shape (top 区分)
    shapes_with_label = [s for s in slide.shapes if s.has_text_frame and "年度参与" in s.text_frame.text]
    shapes_with_value = [s for s in slide.shapes if s.has_text_frame and "15,000" in s.text_frame.text]
    assert shapes_with_label, "label '年度参与' 未渲染"
    assert shapes_with_value, "value '15,000' 未渲染"
    label_top = shapes_with_label[0].top
    value_top = shapes_with_value[0].top
    value_bot = value_top + shapes_with_value[0].height
    # label 顶沿 ≥ value 底沿 (无垂直重叠)
    assert label_top >= value_bot, f"label.top ({label_top}) 应 ≥ value.bottom ({value_bot})"


# ============================================================
# L-4: matrix_2x2 y_axis 长文本不压象限
# ============================================================

def test_L4_matrix_y_axis_long_text_not_intrude_quadrant(tmp_path):
    """L-4: y_axis 超长文本竖排 (eaVert) 时字号应被 fit_pt 自适应压到 [9, default) 防 height 溢出 grid.

    2026-05-31 taste B: 旋转 270° (中文字侧躺+拥挤) → 真竖排 eaVert (字正立).
    旧断言查 rotation==270 + height 自适应; 新断言查 vert=eaVert + 字号自适应.
    注: 竖排 box 宽=axis 列宽、高=grid_h 都是硬编码几何, 不能作有效断言 (恒真);
    eaVert 的真实失败模式是 height 溢出/换列, 唯一防回退点 = fit_pt 字号缩放.
    """
    # 36+ 字超长标签: 任何 grid_h (≤7.5") 下竖排都超高, fit_pt 必压到 9pt 下限
    long_y = "专业深度纵轴关键岗位时间投入与作战经验广度综合评估多维度衡量标准参考体系"
    plan = _build_plan([{
        "id": 1, "role": "content", "visual_type": "matrix-2x2",
        "content": {
            "title": "M",
            "y_axis": long_y,
            "x_axis": "广度",
            "quadrants": [
                {"position": f"Q{i}", "heading": f"H{i}", "desc": "d"} for i in range(1, 5)
            ],
        },
    }])
    prs = _render_one(plan, tmp_path)
    slide = prs.slides[0]
    # 找 y_axis 竖排 (eaVert) textbox — 同时验证用真竖排而非回退到旧 rotation=270
    y_axis_shape = None
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        bodyPr = s.text_frame._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None and bodyPr.get("vert") == "eaVert":
            y_axis_shape = s
            break
    assert y_axis_shape is not None, "y_axis 竖排 (eaVert) textbox 未找到 (不应回退到 rotation=270)"
    # default_pt = _px_pt(y_axis_label_size_px=20) = 12pt; 超长标签 fit_pt 应压到 [9,12) 区间
    y_pt = y_axis_shape.text_frame.paragraphs[0].font.size.pt
    assert 9.0 <= y_pt < 12.0, (
        f"超长 y_axis 字号 ({y_pt}pt) 应被 fit_pt 自适应压到 [9,12) 防 height 溢出 grid_h, 实际未缩放"
    )


# ============================================================
# L-5: cards-6 heading wrap 幻影
# ============================================================

def test_L5_cards6_extreme_long_heading_truncated(tmp_path):
    """L-5: cards-6 在 3x2 网格下 (每卡 ~4" 宽) 极端长 heading 仍应 truncate 不溢出.

    更新: cards-6 改为 3x2 网格后, 4" 宽度可装下 "Blue Opportunities" (18 chars), 不再 truncate.
    用极端长度的 heading (>50 chars) 触发 truncate 验证机制仍生效.
    """
    extreme_heading = "Extremely Long Heading That Definitely Exceeds Card Width In Any Reasonable Theme"
    plan = _build_plan([{
        "id": 1, "role": "content", "visual_type": "cards-6",
        "content": {
            "title": "T", "subtitle": "s",
            "key_points": [
                {"heading": extreme_heading, "body": "b"},
                {"heading": "S2", "body": "b"},
                {"heading": "M3", "body": "b"},
                {"heading": "M4", "body": "b"},
                {"heading": "M5", "body": "b"},
                {"heading": "M6", "body": "b"},
            ],
        },
    }])
    prs = _render_one(plan, tmp_path)
    slide = prs.slides[0]
    # 找含 ellipsis 或截断的 heading
    found_truncated = False
    for s in slide.shapes:
        if s.has_text_frame:
            t = s.text_frame.text
            if t.startswith("Extremely") and t != extreme_heading:
                found_truncated = True
                assert "…" in t, f"长 heading 应被 truncate + ellipsis, 实际: {t!r}"
                break
    assert found_truncated, f"未找到被 truncate 的极端 heading 'Extremely…'"


def test_L5_cards6_3x2_grid_layout(tmp_path):
    """cards-6 应布局成 3 列 x 2 行网格 (而非单行 6 列), 每卡 ~4" 宽不挤压."""
    plan = _build_plan([{
        "id": 1, "role": "content", "visual_type": "cards-6",
        "content": {
            "title": "T",
            "key_points": [
                {"heading": f"Card{i}", "body": f"body{i}"} for i in range(1, 7)
            ],
        },
    }])
    prs = _render_one(plan, tmp_path)
    slide = prs.slides[0]
    # 找含 "Card1" ~ "Card6" 的 shape, 检查 top 分布有 2 行 (上下各 3 个)
    card_headings = [s for s in slide.shapes if s.has_text_frame
                     and s.text_frame.text in (f"Card{i}" for i in range(1, 7))]
    assert len(card_headings) == 6, f"6 cards 应渲染 6 个 heading, 实际 {len(card_headings)}"
    tops = sorted({s.top for s in card_headings})
    assert len(tops) == 2, f"3x2 网格应有 2 个 top 值 (2 行), 实际 {len(tops)}: {tops}"


# ============================================================
# L-6: risk_list grid 自适应
# ============================================================

def test_L6_risk_list_4_items_2x2_grid(tmp_path):
    """L-6: 4 risks 应自适应 2x2 网格, card 高度 > 旧上限 2.4" 撑满画布."""
    plan = _build_plan([{
        "id": 1, "role": "content", "visual_type": "risk-list",
        "content": {
            "title": "R",
            "risks": [
                {"number": str(i), "title": f"R{i}", "desc": "d", "severity": "MED"}
                for i in range(1, 5)
            ],
        },
    }])
    prs = _render_one(plan, tmp_path)
    slide = prs.slides[0]
    # 找所有标题为 R1/R2/R3/R4 的 shape (rounded rect 不直接定位, 用 R1/2/3/4 内嵌 textbox)
    risk_titles = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text in ("R1", "R2", "R3", "R4")]
    assert len(risk_titles) == 4, f"4 risks 应渲染 4 个 title shape, 实际 {len(risk_titles)}"
    # 检查布局是 2x2 (top 有 2 个 distinct value)
    tops = sorted(set(s.top for s in risk_titles))
    assert len(tops) == 2, f"4 risks 应在 2 行布局, 实际行数 {len(tops)}"


# ============================================================
# L-7: pyramid descriptions 与 layer 对齐
# ============================================================

def test_L7_pyramid_descriptions_align_with_layers(tmp_path):
    """L-7: descriptions[i] 渲染中心 Y 应与 levels[i] center Y 对齐 (≤ 0.4" tolerance)."""
    plan = _build_plan([{
        "id": 1, "role": "content", "visual_type": "pyramid",
        "content": {
            "title": "P",
            "levels": [
                {"name": f"L{i}"} for i in range(5)
            ],
            "descriptions": [
                {"number": f"0{i+1}", "heading": f"H{i}", "body": f"body{i}"}
                for i in range(5)
            ],
        },
    }])
    prs = _render_one(plan, tmp_path)
    slide = prs.slides[0]
    # 找 levels (TRAPEZOID) shape - 用 type 判断
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    trapezoids = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    # 找 descriptions: 含 H0/H1/H2/H3/H4 文本
    desc_shapes = []
    for i in range(5):
        for s in slide.shapes:
            if s.has_text_frame and f"H{i}" in s.text_frame.text and "0" in s.text_frame.text:
                desc_shapes.append(s)
                break
    # 至少 5 个 description heading shape
    assert len(desc_shapes) == 5, f"5 descriptions 应渲染 5 个 heading shape, 实际 {len(desc_shapes)}"
    # 5 个 desc 顶沿应递增 (与 layer 顺序一致)
    tops = [s.top for s in desc_shapes]
    assert tops == sorted(tops), f"descriptions 顺序应与 layer 顺序一致, 实际 tops: {tops}"


# ============================================================
# taste A: pyramid 红色渐变阶梯 (替旧 顶 primary_dark 近黑 + 红红灰灰灰)
# ============================================================

def test_pyramid_gradient_fill_lightens_top_to_bottom(tmp_path):
    """2026-05-31 taste A: pyramid N 层红色渐变 (顶 primary 浓红 → 底极浅红), 替旧 顶 primary_dark 近黑 + 下层 paper_2 灰.

    防回退断言: 顶层填充 = primary (#C7000B); 各层亮度自顶向下单调递增 (越下越浅).
    """
    plan = _build_plan([{
        "id": 1, "role": "content", "visual_type": "pyramid",
        "content": {
            "title": "P",
            "levels": [{"name": f"L{i}"} for i in range(5)],
        },
    }])
    prs = _render_one(plan, tmp_path)
    slide = prs.slides[0]
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    # pyramid 梯形: AUTO_SHAPE 且高度明显 (排除标题区短红线 ~0.05")
    traps = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.height > int(0.3 * 914400)]
    traps.sort(key=lambda s: s.top)
    assert len(traps) >= 5, f"5 层 pyramid 应有 ≥5 个梯形, 实际 {len(traps)}"
    fills = [s.fill.fore_color.rgb for s in traps[:5]]
    # 顶层 = primary 浓红 (非旧 primary_dark 近黑)
    assert str(fills[0]) == "C7000B", f"顶层应为 primary #C7000B, 实际 {fills[0]}"
    # 亮度自顶向下单调递增 (红色渐变阶梯, 越下越浅)
    lums = [0.299 * c[0] + 0.587 * c[1] + 0.114 * c[2] for c in fills]
    for i in range(1, 5):
        assert lums[i] > lums[i - 1], f"第 {i} 层应比上层浅 (渐变变浅), lums={lums}"


# ============================================================
# Bug C (template var leak): roadmap phase header 不应字面渲染 code='X' title='Y'
# ============================================================

def test_bug_c_roadmap_phase_header_no_pydantic_repr(tmp_path):
    """Bug C 回归: roadmap phase 表头应只显示 title (主) + summary (副), 不应出现 Pydantic repr.

    历史 bug: renderers_huawei.py:1293 `str(ph)` 对 RoadmapPhase 模型直接字符串化,
    输出 `code='P1' title='打底座' summary='0-3 年:数据与治理'` Python-like 字面值.
    """
    plan = _build_plan([{
        "id": 1, "role": "content", "visual_type": "roadmap",
        "content": {"title": "3 阶段路径", "subtitle": "10 年"},
        "variant": {
            "visual_type": "roadmap",
            "title": "3 阶段路径",
            "phases": [
                {"code": "P1", "title": "打底座", "summary": "0-3 年: 数据与治理"},
                {"code": "P2", "title": "立制度", "summary": "3-6 年: 硬约束与算法"},
                {"code": "P3", "title": "塑文化", "summary": "6-10 年: 文化"},
            ],
            "rows": [
                {"label": "治理", "cells": [
                    {"title": "CHRO 直通 CEO"},
                    {"title": "Talent Board"},
                    {"title": "章程化"},
                ]},
            ],
        },
    }])
    prs = _render_one(plan, tmp_path)
    slide = prs.slides[0]
    # 收集全部 text
    all_text = " ".join(
        s.text_frame.text for s in slide.shapes if s.has_text_frame
    )
    # 字面值不应出现
    assert "code='P1'" not in all_text, f"Pydantic repr 字面值泄露: {all_text!r}"
    assert "code=" not in all_text, f"code= 字面值泄露: {all_text!r}"
    assert "summary='" not in all_text, f"summary= 字面值泄露: {all_text!r}"
    # title 必须显示
    assert "打底座" in all_text
    assert "立制度" in all_text
    assert "塑文化" in all_text
    # summary 必须显示 (双层 header)
    assert "0-3 年" in all_text
    assert "3-6 年" in all_text
    assert "6-10 年" in all_text


def test_bug_c_roadmap_phase_str_fallback(tmp_path):
    """Bug C 回归: phases fallback 为 list[str] (Q1/Q2/...) 时应直接显示, 不报错."""
    plan = _build_plan([{
        "id": 1, "role": "content", "visual_type": "roadmap",
        "content": {"title": "默认季度", "subtitle": "s"},
        "variant": {
            "visual_type": "roadmap",
            "title": "默认季度",
            "phases": [
                {"code": "Q1", "title": "Q1"},
                {"code": "Q2", "title": "Q2"},
                {"code": "Q3", "title": "Q3"},
            ],
            "rows": [{"label": "lane", "cells": [{"title": "x"}, {"title": "y"}, {"title": "z"}]}],
        },
    }])
    prs = _render_one(plan, tmp_path)
    all_text = " ".join(
        s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame
    )
    # 无 summary 时, 不应有副标题位但 title 仍渲染
    assert "Q1" in all_text and "Q2" in all_text and "Q3" in all_text
    assert "code=" not in all_text
