"""renderer_kit.py: 共享 renderer helper 单一来源 (S4 god-module 拆分 Phase 1).

从 render.py (registry + theme helper + HELPER FUNCTIONS + COMMON COMPONENTS +
CARD HEADER + _render_cards) 与 renderers_huawei.py (单位/读值/布局工具) 抽取的纯
helper 集合, 加上 renderer 注册表 (register_renderer + _RENDERERS)。

设计目标: 消除 render.py <-> renderers_huawei.py 循环依赖。renderer_kit 是单向叶子:
不 import render.py, 不 import 任何 renderer 模块。render.py / renderers_huawei /
renderers/** 全部 `from engine.renderer_kit import ...`。

纯搬移, 函数签名与实现与原文件 100% 一致 (S4 Phase 1 零行为变化)。
"""
from __future__ import annotations

import logging
import sys
from functools import lru_cache
from pathlib import Path
from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# PEP 723 standalone 兼容: 确保 SKILLS_DIR 在 sys.path, 让 `from lib.* / from schemas.*`
# 在 render.py 以 __main__ 运行时也能解析 (renderer_kit 本身永远以 engine.renderer_kit 导入)。
_SKILLS_DIR = Path(__file__).resolve().parent.parent
if str(_SKILLS_DIR) not in sys.path:
    sys.path.insert(0, str(_SKILLS_DIR))

from lib.pptx_compat import qn, etree
from lib.font_fallback import resolve_font_for_pptx
from lib.margins import SafeArea
from schemas.slide_plan import SlideSpec, StructuredPoint

logger = logging.getLogger(__name__)

# 模块级 qn() 常量, 避免热路径 (add_rounded_rect 等) 每次重做 namespace lookup.
_QN_PRSTGEOM = qn("a:prstGeom")
_QN_AVLST = qn("a:avLst")
_QN_GD = qn("a:gd")


# ===========================================================================
# RENDERER REGISTRY (从 render.py 迁入, 单一来源)
# ===========================================================================
_RENDERERS: dict[str, callable] = {}

def register_renderer(visual_type: str):
    """Decorator to register a visual type renderer.

    重复注册时输出 INFO 而不静默覆盖, 避免 huawei renderer 与 legacy lambda 冲突静默合并
    (历史 cards-6 collision bug: range(2,7) lambda 与 huawei @register_renderer("cards-6")
    双注册, 净 dispatch 正确但完全靠 import 顺序, 任何顺序调整都让 cards-6 fallback 单行 6 列挤压).
    """
    def wrapper(fn):
        if visual_type in _RENDERERS:
            print(
                f"[render] renderer override: {visual_type} → {fn.__module__}.{fn.__name__}",
                file=sys.stderr,
            )
        _RENDERERS[visual_type] = fn
        return fn
    return wrapper


# ===========================================================================
# THEME HELPERS
# ===========================================================================

def _resolve_accent(spec: SlideSpec, theme: dict, idx: int = 0) -> str:
    """解析 accent 色: 优先 spec.design.accent_color, None 时用 theme.colors.primary,
    最终兜底 huawei 主红 (避免 theme 损坏时整体崩)。idx 给将来卡片轮转色预留。
    """
    if spec.design.accent_color:
        return spec.design.accent_color
    primary = theme.get("colors", {}).get("primary")
    if primary:
        return primary
    return "#C7000B"


def _resolve_title_color(spec: SlideSpec, theme: dict) -> str:
    """解析主标题色: 优先 spec.design.title_color, None 时用 theme.colors.ink, 兜底深灰。

    P3b/P5 收尾: 配套 P4 SlideDesign 默认色 None, 让 renderer 公共组件不传 None 给 hex_to_rgb。
    """
    if spec.design.title_color:
        return spec.design.title_color
    ink = theme.get("colors", {}).get("ink")
    if ink:
        return ink
    return "#1F1F1F"


def _resolve_body_color(spec: SlideSpec, theme: dict) -> str:
    """解析正文色: 优先 spec.design.body_color, None 时用 theme.colors.ink_soft, 兜底中灰。

    P3b/P5 收尾: 配套 P4 SlideDesign 默认色 None, 让 renderer 公共组件不传 None 给 hex_to_rgb。
    """
    if spec.design.body_color:
        return spec.design.body_color
    ink_soft = theme.get("colors", {}).get("ink_soft")
    if ink_soft:
        return ink_soft
    return "#4B4B4B"


def _ve(theme: dict) -> dict:
    """Get visual_elements from theme with defaults."""
    defaults = {
        "footer_show_page_number": True,
        "footer_show_chapter": True,
        "footer_height_inches": 0.35,
        "card_header_height_inches": 0.4,
        "card_header_style": "color_bar",
        "card_content_alignment": "top",
        "card_show_number": False,
        "card_number_style": "circle",
        "card_number_size_pt": 22,
        "metric_container_style": "rounded_rect",
        "metric_value_size_pt": 44,
        "metric_label_size_pt": 14,
        "description_size_pt": 14,
        "description_color": None,
        "divider_below_title": False,
        "divider_color": "#E2E8F0",
    }
    ve = theme.get("visual_elements", {})
    return {**defaults, **ve}

# ===========================================================================
# HELPER FUNCTIONS
# ===========================================================================

@lru_cache(maxsize=128)
def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#RRGGBB' to RGBColor.

    @lru_cache: 50 页 deck × 多个 shape 共用同一色值 (huawei 主题 token 颜色 < 20 unique),
    单次 cache miss 后剩余调用 O(1).
    """
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_textbox(slide, left, top, width, height, text, font_name, font_size_pt, font_color, alignment=PP_ALIGN.LEFT, bold=False):
    """Add a text box with standard formatting."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = resolve_font_for_pptx(font_name)
    p.font.size = Pt(font_size_pt)
    p.font.color.rgb = hex_to_rgb(font_color)
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_textbox_rich(slide, left, top, width, height, parts: list[dict], font_name: str, alignment=PP_ALIGN.LEFT):
    """Add a text box with mixed formatting (bold heading + regular body).

    parts: list of {"text": str, "size_pt": int, "color": str, "bold": bool}
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, part in enumerate(parts):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        run.text = part["text"]
        run.font.name = resolve_font_for_pptx(font_name)
        run.font.size = Pt(part["size_pt"])
        run.font.color.rgb = hex_to_rgb(part["color"])
        run.font.bold = part.get("bold", False)
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, corner_radius=0.1):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.fill.background()  # no border
    # Set corner radius via XML.
    # 2026-05-19 audit fix: MSO_SHAPE.ROUNDED_RECTANGLE 默认 avLst 已含 `<a:gd name="adj"/>`,
    # 直接 SubElement 又加一个会让 PowerPoint 取第一个 (默认 16.67% 圆角) → token 指定的 corner_radius 失效.
    # 必须先 remove 已有的 same-name gd 再 SubElement.
    sp = shape._element
    prstGeom = sp.find(_QN_PRSTGEOM)
    if prstGeom is not None:
        avLst = prstGeom.find(_QN_AVLST)
        if avLst is None:
            avLst = etree.SubElement(prstGeom, _QN_AVLST)
        # 清空已有的 gd 子元素 (避免重复 name="adj")
        for old_gd in list(avLst):
            if old_gd.tag == _QN_GD:
                avLst.remove(old_gd)
        gd = etree.SubElement(avLst, _QN_GD)
        radius_val = int(corner_radius / min(width, height) * 50000)
        gd.set("name", "adj")
        gd.set("fmla", f"val {radius_val}")
    return shape

def estimate_content_height(text: str, width_inches: float, font_size_pt: int,
                            line_spacing: float = 1.15, padding: float = 0.4) -> float:
    """Estimate the height needed to display text in a given width.

    2026-05-19 audit fix (ARCH-007): 委托给 lib/text_metrics.py 单一来源 (中英文加权).
    padding 在此函数语义是上下总 padding (不像 text_metrics padding_inches 是单边), 折半传入.
    """
    from lib.text_metrics import estimate_text_height_inches
    return estimate_text_height_inches(
        text=text,
        width_inches=width_inches,
        font_size_pt=font_size_pt,
        line_spacing=line_spacing,
        padding_inches=padding / 2,
    )


def compute_card_height(points: list[str], card_width: float, font_size_pt: int,
                        max_height: float, min_height: float = 1.5) -> float:
    """Compute card height that fits the tallest point, clamped to bounds."""
    usable_w = card_width - 0.3  # internal padding
    tallest = max(
        (estimate_content_height(p, usable_w, font_size_pt) for p in points),
        default=min_height,
    )
    return max(min_height, min(tallest, max_height))


def set_slide_background(slide, color_hex: str):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def normalize_point(p) -> StructuredPoint:
    """Normalize a key_point (str or dict or StructuredPoint) into StructuredPoint."""
    if isinstance(p, StructuredPoint):
        return p
    if isinstance(p, dict):
        return StructuredPoint(**p)
    if isinstance(p, str):
        if ": " in p:
            heading, body = p.split(": ", 1)
            return StructuredPoint(heading=heading.strip(), body=body.strip())
        return StructuredPoint(body=p)
    return StructuredPoint(body=str(p))


def _adaptive_point_to_structured(p) -> StructuredPoint:
    """AdaptivePoint (cards/process/comparison 自适应母版 path X 条目) → StructuredPoint。

    AdaptivePoint 字段 heading/body/icon; StructuredPoint 无 icon (cards/process/comparison
    renderer 当前不渲染 icon, 忽略)。元素可能是 AdaptivePoint 对象或 dict (防御性)。
    """
    if isinstance(p, StructuredPoint):
        return p
    if isinstance(p, dict):
        return StructuredPoint(heading=p.get("heading"), body=p.get("body") or "")
    # AdaptivePoint / 其它带属性对象
    return StructuredPoint(
        heading=getattr(p, "heading", None),
        body=getattr(p, "body", None) or "",
    )


def get_points(spec: SlideSpec) -> list[StructuredPoint]:
    """Extract and normalize points from spec, path Y (content.key_points) 优先, path X (variant.points) fallback。

    收敛后无后缀母版 (cards/process/comparison) 的 typed variant 把条目放在 `variant.points`
    (list[AdaptivePoint]); 而 path Y 仍用 `content.key_points`。本函数统一两路: key_points 非空
    时走原路径 (零行为变化); 为空时 fallback 读 variant.points (AdaptivePoint.heading/body)。
    这是 T8 新建的 path X→points 桥, AdaptivePoint docstring 所述"双路兼容"由此落地。
    """
    if spec.content.key_points:
        return [normalize_point(p) for p in spec.content.key_points]
    # path X fallback: typed CardsContent/ProcessContent/ComparisonContent.points
    variant = getattr(spec, "variant", None)
    variant_points = getattr(variant, "points", None) if variant is not None else None
    if variant_points:
        return [_adaptive_point_to_structured(p) for p in variant_points]
    return []


def get_point_bodies(points: list[StructuredPoint]) -> list[str]:
    """Extract body text from normalized points (for height calculation)."""
    return [p.body for p in points]

# ===========================================================================
# COMMON COMPONENTS
# ===========================================================================

def render_title_zone(slide, spec: SlideSpec, theme: dict, safe: SafeArea) -> float:
    """Render title + description. Returns total height used."""
    ve = _ve(theme)
    y = safe.top
    # Title
    add_textbox(
        slide, safe.left, y, safe.width, 0.6,
        spec.content.title, spec.design.font_family,
        spec.design.title_size_pt, _resolve_title_color(spec, theme),
        bold=True,
    )
    y += 0.6

    # Description (if present)
    desc = spec.content.description
    if desc:
        desc_color = ve["description_color"] or theme.get("colors", {}).get("text_secondary", _resolve_body_color(spec, theme))
        desc_size = ve["description_size_pt"]
        desc_h = estimate_content_height(desc, safe.width, desc_size, padding=0.15)
        desc_h = min(desc_h, 0.8)  # cap description height
        add_textbox(
            slide, safe.left, y + 0.05, safe.width, desc_h,
            desc, spec.design.font_family,
            desc_size, desc_color,
        )
        y += desc_h + 0.1

    # Optional divider line below title zone
    if ve.get("divider_below_title"):
        divider_color = ve.get("divider_color", "#E2E8F0")
        line = slide.shapes.add_connector(
            1,  # MSO_CONNECTOR.STRAIGHT
            Inches(safe.left), Inches(y + 0.02),
            Inches(safe.left + safe.width), Inches(y + 0.02),
        )
        line.line.color.rgb = hex_to_rgb(divider_color)
        line.line.width = Pt(0.75)
        y += 0.08

    return y - safe.top


def render_footer(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> float:
    """Render footnote + page number footer. Returns footer height."""
    ve = _ve(theme)
    if not ve["footer_show_page_number"] and not spec.content.footnote:
        return 0.0

    footer_h = ve["footer_height_inches"]
    footer_top = safe.bottom - footer_h
    colors = theme.get("colors", {})
    footer_color = colors.get("text_secondary", _resolve_body_color(spec, theme))
    caption_size = theme.get("typography", {}).get("caption_size_pt", 10)

    # Footnote (left)
    if spec.content.footnote:
        add_textbox(
            slide, safe.left, footer_top, safe.width * 0.7, footer_h,
            spec.content.footnote, spec.design.font_family,
            caption_size, footer_color,
        )

    # Page number (right)
    if ve["footer_show_page_number"]:
        page_text = f"{spec.id} / {total_slides}"
        add_textbox(
            slide, safe.left + safe.width * 0.7, footer_top,
            safe.width * 0.3, footer_h,
            page_text, spec.design.font_family,
            caption_size, footer_color,
            alignment=PP_ALIGN.RIGHT,
        )

    return footer_h


def get_content_zone(safe: SafeArea, title_h: float, footer_h: float) -> tuple[float, float]:
    """Return (content_top, content_height) after title and footer."""
    content_top = safe.top + title_h + 0.15
    content_h = safe.height - title_h - footer_h - 0.3
    return content_top, max(content_h, 1.0)


# ===========================================================================
# CARD HEADER RENDERERS
# ===========================================================================

def _get_card_accent(theme: dict, spec: SlideSpec, index: int | None) -> str:
    """Pick accent color for a card: theme.colors.card_left_bar_colors[i] if available,
    else spec.design.accent_color, else theme.colors.primary (P5: 走 token).
    """
    if index is None:
        return _resolve_accent(spec, theme)
    bars = theme.get("colors", {}).get("card_left_bar_colors")
    if bars:
        return bars[index % len(bars)]
    return _resolve_accent(spec, theme)


def _truncate_for_single_line(text: str, max_width_in: float, font_size_pt: float) -> str:
    """L-5 fix: 强制单行, 超长 truncate + ellipsis. 区分中英文等价字宽.

    避免 add_textbox 的 word_wrap=True 让长 heading 在 header 区外溢出造成"幻影"白色文字.
    使用保守 buffer: 实际可用宽度 -= 0.15", 防止 PowerPoint 渲染时仍判定超宽自动扩 textbox.
    """
    if not text:
        return ""
    text_str = str(text)
    safe_width_in = max(0.3, max_width_in - 0.15)  # 留 0.15" buffer 防 PowerPoint auto-grow
    chars_per_in = 72.0 / max(font_size_pt, 1.0)
    max_chars = max(2, int(safe_width_in * chars_per_in))

    def weighted(s: str) -> float:
        # 英文字符按 0.6 (略保守, 加宽以防 LibreOffice/PowerPoint 字号渲染稍宽)
        return sum(1.0 if "一" <= ch <= "鿿" else 0.6 for ch in s)

    if weighted(text_str) <= max_chars:
        return text_str
    # 二分 truncate (text_str[:mid] + "…")
    lo, hi = 1, len(text_str)
    while lo < hi:
        mid = (lo + hi + 1) // 2
        if weighted(text_str[:mid] + "…") <= max_chars:
            lo = mid
        else:
            hi = mid - 1
    return text_str[:lo].rstrip() + "…"


def _add_textbox_no_wrap(slide, left, top, width, height, text, font_name, font_size_pt, font_color, bold=False, alignment=None):
    """L-5 helper: 强制 word_wrap=False 的 textbox, 用于 card header 单行场景."""
    if alignment is None:
        alignment = PP_ALIGN.LEFT
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = resolve_font_for_pptx(font_name)
    p.font.size = Pt(font_size_pt)
    p.font.color.rgb = hex_to_rgb(font_color)
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def render_card_header(slide, x, y, width, heading: str, theme: dict, spec: SlideSpec, index: int | None = None) -> float:
    """Render card header based on theme style. Returns header height (0 if no heading).

    Index enables per-card accent color from theme.colors.card_left_bar_colors.

    L-5 fix: heading 强制单行 + 超长 truncate, 防止 wrap 后白色文字溢出 header 在浅色 body 上形成幻影.
    """
    if not heading:
        return 0.0
    ve = _ve(theme)
    style = ve["card_header_style"]
    header_h = ve["card_header_height_inches"]
    accent = _get_card_accent(theme, spec, index)

    if style == "color_bar":
        # Full-width color bar with white text
        add_rounded_rect(slide, x, y, width, header_h, accent,
                         corner_radius=theme.get("visual_preferences", {}).get("corner_radius_inches", 0.1))
        text_w = width - 0.2
        truncated = _truncate_for_single_line(heading, text_w, spec.design.body_size_pt)
        _add_textbox_no_wrap(
            slide, x + 0.1, y + 0.05, text_w, header_h - 0.1,
            truncated, spec.design.font_family,
            spec.design.body_size_pt, "#FFFFFF",
            bold=True,
        )
        return header_h

    elif style == "left_bar":
        # Left accent bar + heading text on card background
        bar_w = 0.06
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(bar_w), Inches(header_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(accent)
        shape.line.fill.background()
        text_w = width - bar_w - 0.2
        font_pt = spec.design.body_size_pt + 1
        truncated = _truncate_for_single_line(heading, text_w, font_pt)
        _add_textbox_no_wrap(
            slide, x + bar_w + 0.1, y + 0.05, text_w, header_h - 0.1,
            truncated, spec.design.font_family,
            font_pt, _resolve_title_color(spec, theme),
            bold=True,
        )
        return header_h

    else:  # "none" — bold text only
        text_w = width - 0.3
        truncated = _truncate_for_single_line(heading, text_w, spec.design.body_size_pt)
        _add_textbox_no_wrap(
            slide, x + 0.15, y + 0.05, text_w, header_h - 0.1,
            truncated, spec.design.font_family,
            spec.design.body_size_pt, _resolve_title_color(spec, theme),
            bold=True,
        )
        return header_h


def render_card_number_badge(slide, x, y, theme: dict, spec: SlideSpec, index: int) -> None:
    """Render a number badge (01/02/03...) at the top-right corner of a card."""
    ve = _ve(theme)
    if not ve["card_show_number"]:
        return
    accent = _get_card_accent(theme, spec, index)
    size_pt = ve["card_number_size_pt"]
    style = ve["card_number_style"]
    num_text = f"{index + 1:02d}"

    if style == "circle":
        diameter = 0.5
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y), Inches(diameter), Inches(diameter),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(accent)
        shape.line.fill.background()
        add_textbox(
            slide, x, y + 0.03, diameter, diameter - 0.06,
            num_text, spec.design.font_family,
            14, "#FFFFFF",
            alignment=PP_ALIGN.CENTER, bold=True,
        )
    else:  # plain
        add_textbox(
            slide, x, y, 0.7, 0.5,
            num_text, spec.design.font_family,
            size_pt, accent,
            bold=True,
        )


# ===========================================================================
# CARDS RENDERER (从 render.py 迁入; legacy cards-2..5 + huawei cards-6 共享底座)
# ===========================================================================

def _render_cards(slide, spec: SlideSpec, theme: dict, safe: SafeArea, n_cards: int, total_slides: int):
    """N-card layout with optional card headers.

    Layout 选择:
    - n_cards in {2, 3, 4, 5}: 单行 N 列 (传统横排)
    - n_cards == 6: **3 列 x 2 行** 网格 (符合 Cards6Content 文档"六宫格 (3 列 x 2 行)" 规范, 不再单行 6 列挤压)
    """
    set_slide_background(slide, spec.design.background)
    title_h = render_title_zone(slide, spec, theme, safe)
    footer_h = render_footer(slide, spec, theme, safe, total_slides)
    content_top, content_h = get_content_zone(safe, title_h, footer_h)

    ve = _ve(theme)
    colors = theme.get("colors", {}).get("card_fills", ["#F4F4F4"] * 5)
    gap = theme.get("spacing", {}).get("element_gap_inches", 0.25)
    corner_r = theme.get("visual_preferences", {}).get("corner_radius_inches", 0.1)

    # 自适应布局派生 (T8 收敛: legacy cards-2..5 单行 + huawei cards-6 网格统一为按 n 派生):
    #   n in {2, 3}: 单行 N 列 (横排)
    #   n == 4:      双行 2x2
    #   n in {5, 6}: 3 列 x 2 行网格 (n=5 时末格留空, 复用 cards-6 网格底座)
    # 视觉统一会改变旧 cards-4 (4 列单行→2x2) / cards-5 (5 列单行→3x2) 外观 = 蓝图决策2 拍板接受。
    if n_cards <= 3:
        n_cols, n_rows = n_cards, 1
    elif n_cards == 4:
        n_cols, n_rows = 2, 2
    else:  # 5, 6
        n_cols, n_rows = 3, 2
    card_width = (safe.width - gap * (n_cols - 1)) / n_cols

    points = get_points(spec)
    # path X: content.key_points 为空时, 从 variant.cards (Cards6Content.cards: heading/body)
    # 或 content.cards (path Y 扁平) 构造 points。修复 cards-6 在 path X 下只渲染空卡框的 bug。
    if not points:
        cards = None
        variant = getattr(spec, "variant", None)
        if variant is not None:
            cards = getattr(variant, "cards", None)
        if not cards:
            cards = (getattr(spec.content, "__pydantic_extra__", None) or {}).get("cards")
        if cards:
            def _card_field(c, key):
                if isinstance(c, dict):
                    return c.get(key)
                return getattr(c, key, None)
            points = [
                StructuredPoint(
                    heading=(_card_field(c, "heading") or None),
                    body=(_card_field(c, "body") or ""),
                )
                for c in cards
            ]
    bodies = get_point_bodies(points)
    has_headers = any(p.heading for p in points)
    header_h = ve["card_header_height_inches"] if has_headers else 0.0

    # 垂直铺满 content 区 (T9 视觉迭代修头重脚轻):
    # 旧逻辑 card_height=min(内容高度, 网格高度) 在文案短时取矮卡 + 顶部对齐 + 固定 gap 行距,
    # 导致多行网格挤 content 上半、下半大量空白。现改为按行均分 content 区垂直铺满。
    # row_pitch = 每行占用的垂直步距 (含行间距); card_height = row_pitch - 行间距, 卡片更高更舒展。
    # 卡高直接取网格均分高度 (不再 min 到内容高度), 内容靠卡内 body textbox 自适应渲染。
    row_gap = gap if n_rows > 1 else 0.0
    row_pitch = content_h / n_rows
    card_height = max(row_pitch - row_gap, 0.5)

    # 网格整体顶端对齐 content_top; 卡高已按行均分铺满, 自然填充到接近 footer。
    grid_top = content_top

    # 单行场景 (cards-2/3) 卡高上限 (T9 v4): 用户审美判据"内外留白平衡" — 卡片矮下来,
    # 卡外留白 (卡到 content 区上下边) 与卡内文字留白视觉相当。单行卡封顶到 content_h 的 50%
    # (~2.67"), 封顶后整个单行卡带在 content 区垂直居中, 卡外上下留白各约 1.34"。
    # body 文本框 MIDDLE 锚点让文字在矮卡内垂直居中, 内外留白趋于平衡。多行不受影响。
    if n_rows == 1:
        single_row_max = content_h * 0.5
        if card_height > single_row_max:
            card_height = single_row_max
            grid_top = content_top + (content_h - card_height) / 2

    show_number = ve["card_show_number"]
    for i in range(n_cards):
        r = i // n_cols
        c = i % n_cols
        x = safe.left + c * (card_width + gap)
        card_top = grid_top + r * row_pitch
        fill = colors[i % len(colors)]
        use_rounded = theme.get("visual_preferences", {}).get("rounded_corners", True)

        # Card background
        if use_rounded:
            add_rounded_rect(slide, x, card_top, card_width, card_height, fill, corner_r)
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(card_top),
                Inches(card_width), Inches(card_height),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill)
            shape.line.fill.background()

        # Card header (per-card accent from theme.colors.card_left_bar_colors)
        pt = points[i] if i < len(points) else StructuredPoint(body="")
        h_used = render_card_header(slide, x, card_top, card_width, pt.heading, theme, spec, index=i)

        # Number badge (top-right corner, overlays on card)
        if show_number:
            badge_size = 0.5
            render_card_number_badge(
                slide, x + card_width - badge_size - 0.12, card_top - 0.12,
                theme, spec, index=i,
            )

        # Card body text — 垂直居中 (T9 v3, 用户诉求"文字整体在中间, 不顶在上面更舒服")。
        # header 仍在卡顶, header 下方 body 区里文本框垂直锚点 MIDDLE, 让文字在可用空间上下留白均衡。
        text = pt.body
        body_top = card_top + h_used + 0.12
        body_avail = card_height - h_used - 0.22
        body_box = add_textbox(
            slide, x + 0.2, body_top,
            card_width - 0.4, body_avail,
            text, spec.design.font_family,
            spec.design.body_size_pt, _resolve_body_color(spec, theme),
        )
        body_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


# ===========================================================================
# HUAWEI 单位/读值/布局工具 (从 renderers_huawei.py 迁入)
# ===========================================================================
# px → pt 换算（研究报告 §问题 3 验证公式）
# 画布参考 1920×1080 对 slide 13.333"×7.5"，pt ≈ px × 0.6
def _px_pt(px: int | float) -> float:
    return float(px) * 0.6


# px → inch 换算（画布 1920 px → slide 13.333 inch → 1 inch ≈ 144 px）
# 为兼容研究报告 "/96" 的 Web DPI 表述，此处保留 96 DPI 换算便于对照。
def _px_in(px: int | float) -> float:
    return float(px) / 96.0


def _tokens(theme: dict) -> dict:
    return theme.get("tokens", theme)


def _layouts(theme: dict) -> dict:
    return theme.get("_layouts", theme.get("layouts", {}))


def _color(theme: dict, token_name: str, fallback: str = "#000000") -> str:
    return _tokens(theme).get("colors", {}).get(token_name, fallback)


def _lerp_hex(c1: str, c2: str, t: float) -> str:
    """两个 hex 色按 t∈[0,1] 线性插值, 返回大写 #RRGGBB. 用于 pyramid 红色渐变阶梯 (taste A)."""
    a, b = c1.lstrip("#"), c2.lstrip("#")
    r1, g1, b1 = int(a[0:2], 16), int(a[2:4], 16), int(a[4:6], 16)
    r2, g2, b2 = int(b[0:2], 16), int(b[2:4], 16), int(b[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _is_dark(hex_color: str, threshold: float = 140.0) -> bool:
    """按感知亮度 (sRGB 加权) 判断底色是否偏暗, 用于自动选白/黑文字色."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return (0.299 * r + 0.587 * g + 0.114 * b) < threshold


def _type_px(theme: dict, name: str, fallback_px: int = 24) -> int:
    return _tokens(theme).get("type_scale_px", {}).get(name, fallback_px)


def _layout(theme: dict, visual_type_key: str) -> dict:
    return _layouts(theme).get(visual_type_key, {})


def _extra(source, *fields: str, default=None):
    """读 extra 字段 (path X 双查兼容 + 多 alias 字段名).

    优先级 (按 fields 顺序逐一尝试, 第一个非空就返回):
    1. 若 source 有 .variant 属性 (即 SlideSpec, path X 推荐): 先查 variant 直属字段, 再查 variant.__pydantic_extra__
    2. 回退到 content (SlideSpec.content 或 source 自身, path Y 兜底): 查直属, 再查 __pydantic_extra__

    支持单字段: `_extra(spec, "title")`, 也支持多 alias: `_extra(spec, "number", "big_number")`
    (按顺序找第一个非空, schemas 字段名与 renderer 历史读的字段名共存时使用)。

    Args:
        source: SlideSpec 实例 (推荐) 或 SlideContent 实例 (兼容)
        *fields: 字段名, 支持多个 alias 按顺序查找
        default: 兜底值

    Returns:
        字段值, 全部缺失时返回 default
    """
    variant = getattr(source, "variant", None)
    content = getattr(source, "content", source)
    for field in fields:
        if variant is not None:
            v = getattr(variant, field, None)
            if v is not None:
                return v
            vextras = getattr(variant, "__pydantic_extra__", None) or {}
            if field in vextras and vextras[field] is not None:
                return vextras[field]
        # content 直属声明字段 (title/subtitle/description/key_points 等) — 保留。
        val = getattr(content, field, None)
        if val is not None:
            return val
        # T13 (ARCH-002 阶段3): content.__pydantic_extra__ path Y 分支已删 —— SlideContent
        # extra='forbid' 后 __pydantic_extra__ 恒为 None/空, 此分支永不命中 (死代码)。
        # 华为版式专属字段统一走 variant (path X) 上方分支读取。
    return default


def _extra_list_merge(source, field: str) -> list:
    """对 list 类字段, 同时从 path X (variant) 和 path Y (content) 各读一份, 按 index 字段 union 合并.

    Why: LLM 生成 plan 时常同时写 spec.content.levels 和 spec.variant.levels,
    但因 schemas 字段名与 renderer 历史读字段名不同, 两侧数据可能字段互补 (variant 给 name, content 给 title+desc).
    单纯 _extra 优先 path X 会丢 path Y 的字段. 此 helper 按 index 把两侧字段 union, 避免数据丢失.

    Args:
        source: SlideSpec 实例
        field: list 字段名, 如 "levels" / "phases" / "steps" / "layers" / "units" / "descriptions"

    Returns:
        list of dict, 每个 dict 包含两侧 union 后的字段 (path X 优先, 但缺失字段从 path Y 补)
    """
    variant = getattr(source, "variant", None)
    content = getattr(source, "content", source)

    def _extract_list(obj, fname):
        if obj is None:
            return []
        v = getattr(obj, fname, None)
        if isinstance(v, list):
            return v
        extras = getattr(obj, "__pydantic_extra__", None) or {}
        v2 = extras.get(fname)
        return v2 if isinstance(v2, list) else []

    list_v = _extract_list(variant, field)
    # T13 (ARCH-002 阶段3) DEFER 标注: SlideContent extra='forbid' 后, levels/descriptions 等
    # 华为专属 list 字段既非 content 声明字段、又被 forbid 拒, 故 list_c 恒空 —— 下方 merge
    # 实际已退化为只读 variant (path X)。content 侧为死代码, 但简化 merge 循环收益低 (仅
    # pyramid 调用) 且重构有丢字段风险, 按 team-lead 指示 defer 不强求清理, 保留 + 标注。
    list_c = _extract_list(content, field)

    def _to_dict(obj) -> dict:
        if obj is None:
            return {}
        if isinstance(obj, dict):
            return {k: v for k, v in obj.items() if v not in (None, "")}
        if isinstance(obj, str):
            return {"_value": obj}
        # Pydantic BaseModel
        d = {}
        # 已声明字段
        for k in obj.__class__.model_fields if hasattr(obj.__class__, "model_fields") else ():
            v = getattr(obj, k, None)
            if v not in (None, "", []):
                d[k] = v
        # extra 字段
        extras = getattr(obj, "__pydantic_extra__", None) or {}
        for k, v in extras.items():
            if v not in (None, "") and k not in d:
                d[k] = v
        return d

    n = max(len(list_v), len(list_c))
    merged = []
    for i in range(n):
        item_c = _to_dict(list_c[i]) if i < len(list_c) else {}
        item_v = _to_dict(list_v[i]) if i < len(list_v) else {}
        # path Y 在底, path X 覆盖 (相同 key)
        m = dict(item_c)
        for k, v in item_v.items():
            if v not in (None, "") or k not in m:
                m[k] = v
        merged.append(m)
    return merged


def _estimate_text_height_in(text: str, width_in: float, font_size_pt: float,
                             line_spacing: float = 1.25, padding_in: float = 0.15) -> float:
    """估算 wrapped text 在 textbox 中实际占用高度 (inches). 委托 lib/text_metrics 单一来源.

    2026-05-19 audit fix (ARCH-007): 公式 (cn 1.0 / en 0.55 加权) 已沉淀到 lib/text_metrics.py,
    这里仅做 width 兜底 (>=0.5) 后委托, 保留旧 API.
    """
    if not text:
        return 0.0
    from lib.text_metrics import estimate_text_height_inches
    return estimate_text_height_inches(
        text=str(text),
        width_inches=max(0.5, float(width_in)),
        font_size_pt=font_size_pt,
        line_spacing=line_spacing,
        padding_inches=padding_in / 2,  # text_metrics padding 是单边, 这里语义是单边总
    )


def _dget(d, *keys, default=None):
    """对 dict 或 Pydantic BaseModel 实例做多 key alias fallback.

    第一个非空就返回, 都缺时返 default.
    支持 path Y (yaml 顶层 content extra='allow', 解析为 list of dict)
    和 path X (yaml 用 variant, 解析为 list of Pydantic 子模型) 双兼容.
    """
    if d is None:
        return default
    is_dict = isinstance(d, dict)
    for k in keys:
        if is_dict:
            v = d.get(k)
        else:
            v = getattr(d, k, None)
            # Pydantic extra fields 在 __pydantic_extra__ 里
            if v is None:
                extras = getattr(d, "__pydantic_extra__", None) or {}
                v = extras.get(k)
        if v is not None and v != "":
            return v
    return default


def _add_rect(slide, left_in: float, top_in: float, width_in: float, height_in: float,
              fill_hex: str | None, line_hex: str | None = None, line_width_pt: float = 0.0):
    """添加纯色矩形（无圆角）。fill_hex=None 表示无填充；line_hex=None 表示无描边。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    if fill_hex is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(line_hex)
        shape.line.width = Pt(line_width_pt)
    return shape


def clamp(value: int, lo: int, hi: int) -> int:
    """整数钳制到 [lo, hi]。自适应母版按 n=clamp(len(points), lo, hi) 派生布局时复用。"""
    return max(lo, min(hi, value))


def safe_h(height_in: float, minimum: float = 0.1) -> float:
    """高度下限钳制: 防 textbox/shape 计算出负高或近零高传给 Inches() 渲染异常。

    多个 huawei renderer 用 `area_bottom - cursor - padding` 推算剩余高度, 当内容区被上方
    元素挤占时该差值可能为负 (timeline desc / roadmap lane 等)。统一走本 helper 兜底,
    保证 >= minimum (默认 0.1")，避免 Inches(负值) 产出异常 shape。
    """
    return max(minimum, height_in)


def _canvas_dims_in(theme: dict) -> tuple[float, float]:
    """从 layouts.canvas 读画布 px 尺寸并换算 inch。fallback 到 16:9 标准。"""
    canvas = _layouts(theme).get("canvas", {})
    w_px = canvas.get("width_px", 1920)
    h_px = canvas.get("height_px", 1080)
    scale = 13.333 / max(w_px, 1)
    return w_px * scale, h_px * scale


def _canvas_padding_in(theme: dict) -> tuple[float, float, float, float]:
    """返回 (top, right, bottom, left) inches。"""
    canvas = _layouts(theme).get("canvas", {})
    w_px = canvas.get("width_px", 1920)
    scale = 13.333 / max(w_px, 1)
    top = canvas.get("padding_top_px", 88) * scale
    right = canvas.get("padding_right_px", 96) * scale
    bottom = canvas.get("padding_bottom_px", 76) * scale
    left = canvas.get("padding_left_px", 96) * scale
    return top, right, bottom, left


def _slide_dims_in() -> tuple[float, float]:
    """标准 16:9 slide dims."""
    return 13.333, 7.5


# 趋势色 token 决策（Plan §O5）
TREND_COLOR_TOKEN: dict[str, str] = {"up": "ok", "down": "warn", "flat": "ink_soft"}


def _font_family(spec: SlideSpec, theme: dict) -> str:
    """选择字体并接入 resolve_font_for_pptx, 让 Mac/Linux 上有 CJK fallback.

    历史教训 (2026-05-19 audit, INT-002): 此函数原直接返回 tokens.fonts.sans[0] = "Microsoft YaHei",
    huawei renderer 内 7 处 `xx.font.name = font` raw 写入 PPTX. Mac 上 PowerPoint 找不到 Microsoft YaHei
    走自身字体 fallback (Hiragino/PingFang), 与设计意图严重偏离, 严重时中文豆腐方块.
    现在统一接入 resolve_font_for_pptx, FONT_FALLBACK_CHAIN + CJK_FALLBACK_BY_OS 保证每平台都有 CJK 字体.
    """
    sans = _tokens(theme).get("fonts", {}).get("sans", [])
    preferred = sans[0] if sans else (spec.design.font_family or "Microsoft YaHei")
    return resolve_font_for_pptx(preferred)
