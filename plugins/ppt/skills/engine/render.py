# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=1.0.0", "pydantic>=2.0", "pyyaml>=6.0", "lxml>=4.9"]
# ///
"""render.py: Stage 4 - Deterministic PPTX rendering.

Consumes slide-plan.yaml, produces .pptx. NO AI involvement.
Same YAML input always produces same PPTX output.

Uses registry pattern: each visual type has a dedicated renderer function.

Usage:
    uv run --script engine/render.py <slide-plan.yaml> --theme <theme> --output output.pptx
    uv run --script engine/render.py <slide-plan.yaml> --theme <theme> --output output.pptx \
        --base-pptx existing.pptx --only-slides 3,5,7
"""

import argparse
import logging
import sys
from functools import lru_cache
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

SKILLS_DIR = Path(__file__).resolve().parent.parent
# PEP 723 standalone script: add parent dir so `from lib.*` and `from schemas.*`
# resolve when executed via `uv run --script`. Not needed if imported as a library.
sys.path.insert(0, str(SKILLS_DIR))

from lib.pptx_compat import delete_slide
from lib.font_fallback import resolve_font_for_pptx
from lib.margins import get_safe_area, enforce_margins, SafeArea
from lib.content_fitter import suggest_font_size, estimate_text_overflow
from schemas.slide_plan import SlidePlan, SlideSpec, SlideRole, StructuredPoint

# ===========================================================================
# THEME LOADING (delegated to engine.theme_loader, single source of truth)
# ===========================================================================

from engine.theme_loader import (
    load_theme as _tl_load_theme,
    _REMOVED_THEMES as _TL_REMOVED,
    _DEFAULT_THEME as _TL_DEFAULT,
)

# 保留模块级常量供外部 test 直接 import (tests/variants/test_theme_loader.py:19 已用)
_REMOVED_THEMES: set[str] = _TL_REMOVED
_DEFAULT_THEME: str = _TL_DEFAULT


def load_theme(name: str | None) -> dict:
    """已迁移到 engine.theme_loader.load_theme。本函数仅为向后兼容保留。"""
    return _tl_load_theme(name)


# ===========================================================================
# 共享 helper + renderer 注册表 — 迁至 engine.renderer_kit (单一来源, 解循环依赖)
# ===========================================================================
# render.py <-> renderers_huawei.py 旧循环依赖已消除: 二者都 import renderer_kit (单向)。
# 旧 sys.modules.setdefault("engine.render", ...) PEP723 双 module hack 不再需要 —
# huawei 模块改 import renderer_kit 后, render.py 以 __main__ 运行时不再触发 engine.render 二次加载。
from engine.renderer_kit import (  # noqa: E402
    register_renderer, _RENDERERS,
    _resolve_accent, _resolve_title_color, _resolve_body_color, _ve,
    hex_to_rgb, add_textbox, add_textbox_rich, add_rounded_rect,
    estimate_content_height, compute_card_height, set_slide_background,
    normalize_point, get_points, get_point_bodies,
    render_title_zone, render_footer, get_content_zone,
    _get_card_accent, _truncate_for_single_line, _add_textbox_no_wrap,
    render_card_header, render_card_number_badge,
    _render_cards,
)

# ===========================================================================
# VISUAL TYPE RENDERERS
# ===========================================================================

@register_renderer("hero-statement")
def render_hero_statement(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int):
    """Single punchy statement, centered. No footer for title-type slides."""
    set_slide_background(slide, spec.design.background)
    add_textbox(
        slide, safe.left, safe.top + safe.height * 0.3,
        safe.width, safe.height * 0.3,
        spec.content.title, spec.design.font_family,
        spec.design.title_size_pt + 4, _resolve_title_color(spec, theme),
        alignment=PP_ALIGN.CENTER, bold=True,
    )
    if spec.content.subtitle:
        add_textbox(
            slide, safe.left, safe.top + safe.height * 0.6,
            safe.width, 0.5,
            spec.content.subtitle, spec.design.font_family,
            spec.design.body_size_pt, _resolve_body_color(spec, theme),
            alignment=PP_ALIGN.CENTER,
        )

@register_renderer("quote-hero")
def render_quote_hero(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int):
    """Powerful quote with attribution."""
    set_slide_background(slide, spec.design.background)
    add_textbox(
        slide, safe.left + 1, safe.top + safe.height * 0.25,
        safe.width - 2, safe.height * 0.35,
        f'\u201c{spec.content.title}\u201d', spec.design.font_family,
        spec.design.title_size_pt, _resolve_accent(spec, theme),
        alignment=PP_ALIGN.CENTER, bold=False,
    )
    if spec.content.subtitle:
        add_textbox(
            slide, safe.left + 1, safe.top + safe.height * 0.65,
            safe.width - 2, 0.5,
            f"-- {spec.content.subtitle}", spec.design.font_family,
            spec.design.body_size_pt, _resolve_body_color(spec, theme),
            alignment=PP_ALIGN.CENTER,
        )

@register_renderer("bullets")
def render_bullets(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int):
    """Bulleted content with title zone, bullet markers, and footer."""
    set_slide_background(slide, spec.design.background)
    title_h = render_title_zone(slide, spec, theme, safe)
    footer_h = render_footer(slide, spec, theme, safe, total_slides)
    content_top, content_h = get_content_zone(safe, title_h, footer_h)

    points = get_points(spec)
    if not points:
        return

    bullet_w = safe.width - 0.5
    gap = 0.12
    y = content_top

    for pt in points:
        # Bullet marker (small filled circle)
        marker_size = 0.08
        marker_y = y + spec.design.body_size_pt / 72 * 0.4
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(safe.left + 0.15), Inches(marker_y),
            Inches(marker_size), Inches(marker_size),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(_get_card_accent(theme, spec, 0))
        shape.line.fill.background()

        # Text
        if pt.heading:
            text = f"{pt.heading}\n{pt.body}"
            parts = [
                {"text": pt.heading, "size_pt": spec.design.body_size_pt, "color": _resolve_title_color(spec, theme), "bold": True},
                {"text": f"\n{pt.body}", "size_pt": spec.design.body_size_pt, "color": _resolve_body_color(spec, theme), "bold": False},
            ]
            h = estimate_content_height(text, bullet_w, spec.design.body_size_pt)
            add_textbox_rich(slide, safe.left + 0.5, y, bullet_w, h, parts, spec.design.font_family)
        else:
            h = estimate_content_height(pt.body, bullet_w, spec.design.body_size_pt)
            add_textbox(
                slide, safe.left + 0.5, y, bullet_w, h,
                pt.body, spec.design.font_family,
                spec.design.body_size_pt, _resolve_body_color(spec, theme),
            )
        y += h + gap
        if y > content_top + content_h:
            break

# 注册 cards-2 到 cards-5 (单行布局); cards-6 由 renderers_huawei.@register_renderer("cards-6") 提供 3x2 huawei 实现
for n in range(2, 6):
    _RENDERERS[f"cards-{n}"] = lambda s, sp, t, sa, ts, _n=n: _render_cards(s, sp, t, sa, _n, ts)

# Comparison renderer (comparison-2 through comparison-5)
def _render_comparison(slide, spec: SlideSpec, theme: dict, safe: SafeArea, n_cols: int, total_slides: int):
    """N-column comparison layout with header row."""
    set_slide_background(slide, spec.design.background)
    title_h = render_title_zone(slide, spec, theme, safe)
    footer_h = render_footer(slide, spec, theme, safe, total_slides)
    content_top, content_h = get_content_zone(safe, title_h, footer_h)

    colors = theme.get("colors", {}).get("card_fills", ["#F4F4F4"] * 5)
    gap = theme.get("spacing", {}).get("element_gap_inches", 0.25)
    col_width = (safe.width - gap * (n_cols - 1)) / n_cols
    corner_r = theme.get("visual_preferences", {}).get("corner_radius_inches", 0.1)

    points = get_points(spec)
    bodies = get_point_bodies(points)

    header_height = 0.5
    max_body_height = content_h - header_height - 0.15
    body_height = compute_card_height(
        bodies, col_width, spec.design.body_size_pt, max_body_height,
    ) if bodies else max_body_height

    header_top = content_top
    body_top = header_top + header_height + 0.1

    for i in range(n_cols):
        x = safe.left + i * (col_width + gap)
        fill = colors[i % len(colors)]
        pt = points[i] if i < len(points) else StructuredPoint(body="")

        # Header bar
        header_text = pt.heading or f"Option {i+1}"
        add_rounded_rect(slide, x, header_top, col_width, header_height, _get_card_accent(theme, spec, i), corner_r)
        add_textbox(
            slide, x + 0.1, header_top + 0.05,
            col_width - 0.2, header_height - 0.1,
            header_text, spec.design.font_family,
            spec.design.body_size_pt, "#FFFFFF",
            alignment=PP_ALIGN.CENTER, bold=True,
        )
        # Body card
        add_rounded_rect(slide, x, body_top, col_width, body_height, fill, corner_r)
        body_text = pt.body
        add_textbox(
            slide, x + 0.15, body_top + 0.15,
            col_width - 0.3, body_height - 0.3,
            body_text, spec.design.font_family,
            spec.design.body_size_pt, _resolve_body_color(spec, theme),
        )

for n in range(2, 6):
    _RENDERERS[f"comparison-{n}"] = lambda s, sp, t, sa, ts, _n=n: _render_comparison(s, sp, t, sa, _n, ts)

# Process renderer (process-2-phase through process-5-phase)
def _render_process(slide, spec: SlideSpec, theme: dict, safe: SafeArea, n_phases: int, total_slides: int):
    """N-phase process flow with arrow connectors."""
    set_slide_background(slide, spec.design.background)
    title_h = render_title_zone(slide, spec, theme, safe)
    footer_h = render_footer(slide, spec, theme, safe, total_slides)
    content_top, content_h = get_content_zone(safe, title_h, footer_h)

    colors = theme.get("colors", {}).get("card_fills", ["#F4F4F4"] * 5)
    gap = theme.get("spacing", {}).get("element_gap_inches", 0.25)
    arrow_width = 0.3
    total_arrow_space = arrow_width * (n_phases - 1)
    phase_width = (safe.width - gap * (n_phases - 1) - total_arrow_space) / n_phases
    corner_r = theme.get("visual_preferences", {}).get("corner_radius_inches", 0.1)

    points = get_points(spec)
    bodies = get_point_bodies(points)

    # Content-aware phase height
    max_phase_height = content_h
    phase_height = compute_card_height(
        bodies, phase_width, spec.design.body_size_pt, max_phase_height, min_height=2.0,
    ) + 0.6 if bodies else max_phase_height
    phase_height = min(phase_height, max_phase_height)
    phase_top = content_top

    for i in range(n_phases):
        x = safe.left + i * (phase_width + gap + arrow_width)
        fill = colors[i % len(colors)]
        pt = points[i] if i < len(points) else StructuredPoint(body="")

        # Phase box
        add_rounded_rect(slide, x, phase_top, phase_width, phase_height, fill, corner_r)

        # Phase heading (from StructuredPoint.heading or fallback to "Phase N")
        heading = pt.heading or f"Phase {i+1}"
        add_textbox(
            slide, x + 0.15, phase_top + 0.1,
            phase_width - 0.3, 0.4,
            heading, spec.design.font_family,
            spec.design.body_size_pt + 2, _resolve_accent(spec, theme),
            bold=True,
        )
        # Phase body
        text = pt.body
        body_start = phase_top + 0.6
        body_avail = phase_height - 0.7
        add_textbox(
            slide, x + 0.15, body_start,
            phase_width - 0.3, body_avail,
            text, spec.design.font_family,
            spec.design.body_size_pt, _resolve_body_color(spec, theme),
        )
        # Arrow connector (except after last phase)
        if i < n_phases - 1:
            arrow_x = x + phase_width + gap * 0.3
            arrow_y = phase_top + phase_height / 2 - 0.2
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(arrow_y),
                Inches(arrow_width), Inches(0.4),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(_resolve_accent(spec, theme))
            shape.line.fill.background()

for n in range(2, 6):
    _RENDERERS[f"process-{n}-phase"] = lambda s, sp, t, sa, ts, _n=n: _render_process(s, sp, t, sa, _n, ts)

@register_renderer("data-contrast")
def render_data_contrast(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int):
    """Two large metrics side by side with containers and body text."""
    set_slide_background(slide, spec.design.background)
    title_h = render_title_zone(slide, spec, theme, safe)
    footer_h = render_footer(slide, spec, theme, safe, total_slides)
    content_top, content_h = get_content_zone(safe, title_h, footer_h)

    ve = _ve(theme)
    colors = theme.get("colors", {}).get("card_fills", ["#F4F4F4"] * 5)
    gap = theme.get("spacing", {}).get("element_gap_inches", 0.25)
    corner_r = theme.get("visual_preferences", {}).get("corner_radius_inches", 0.1)

    points = get_points(spec)
    half_w = (safe.width - gap) / 2
    metric_size = ve["metric_value_size_pt"]
    label_size = ve["metric_label_size_pt"]
    container_style = ve["metric_container_style"]

    for i in range(min(2, len(points) if points else 0)):
        pt = points[i]
        x = safe.left + i * (half_w + gap)

        if pt.metric_value and container_style != "none":
            # Metric container
            container_h = 2.2
            fill = colors[i % len(colors)]
            if container_style == "circle":
                # Circle container
                circle_size = min(half_w * 0.6, container_h)
                cx = x + (half_w - circle_size) / 2
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(cx), Inches(content_top), Inches(circle_size), Inches(circle_size),
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(fill)
                shape.line.fill.background()
            else:
                # Rounded rect container
                add_rounded_rect(slide, x, content_top, half_w, container_h, fill, corner_r)

            # Metric value (large, centered)
            add_textbox(
                slide, x, content_top + 0.3, half_w, 1.0,
                pt.metric_value, spec.design.font_family,
                metric_size, _resolve_accent(spec, theme),
                alignment=PP_ALIGN.CENTER, bold=True,
            )
            # Metric label (small, centered, below value)
            label = pt.metric_label or pt.heading or ""
            if label:
                add_textbox(
                    slide, x + 0.1, content_top + 1.3, half_w - 0.2, 0.6,
                    label, spec.design.font_family,
                    label_size, _resolve_body_color(spec, theme),
                    alignment=PP_ALIGN.CENTER,
                )
            # Body text below container
            if pt.body:
                body_top = content_top + container_h + 0.2
                body_h = content_h - container_h - 0.3
                if body_h > 0.3:
                    add_textbox(
                        slide, x + 0.1, body_top, half_w - 0.2, body_h,
                        pt.body, spec.design.font_family,
                        spec.design.body_size_pt, _resolve_body_color(spec, theme),
                    )
        else:
            # Fallback: large text (backwards compat with plain strings)
            add_textbox(
                slide, x, content_top + 0.5, half_w, 2.0,
                pt.metric_value or pt.body, spec.design.font_family,
                metric_size, _resolve_accent(spec, theme),
                alignment=PP_ALIGN.CENTER, bold=True,
            )
            if pt.body and pt.metric_value:
                add_textbox(
                    slide, x + 0.1, content_top + 2.8, half_w - 0.2, content_h - 3.0,
                    pt.body, spec.design.font_family,
                    spec.design.body_size_pt, _resolve_body_color(spec, theme),
                )

@register_renderer("table")
def render_table(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int):
    """Table from table_data."""
    set_slide_background(slide, spec.design.background)
    title_h = render_title_zone(slide, spec, theme, safe)
    footer_h = render_footer(slide, spec, theme, safe, total_slides)
    content_top, content_h = get_content_zone(safe, title_h, footer_h)

    td = spec.content.table_data
    if not td:
        return
    rows = len(td.rows) + 1  # +1 for header
    cols = len(td.headers) if td.headers else (len(td.rows[0]) if td.rows else 1)
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(safe.left), Inches(content_top),
        Inches(safe.width), Inches(min(content_h, rows * 0.5)),
    )
    tbl = table_shape.table
    for j, h in enumerate(td.headers):
        cell = tbl.cell(0, j)
        cell.text = h
    for i, row in enumerate(td.rows):
        for j, val in enumerate(row):
            if j < cols:
                tbl.cell(i + 1, j).text = str(val)

@register_renderer("comparison-tables")
def render_comparison_tables(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int):
    """Side-by-side comparison tables."""
    set_slide_background(slide, spec.design.background)
    title_h = render_title_zone(slide, spec, theme, safe)
    footer_h = render_footer(slide, spec, theme, safe, total_slides)
    content_top, content_h = get_content_zone(safe, title_h, footer_h)

    if spec.content.key_points:
        points = get_points(spec)
        y = content_top
        for pt in points:
            text = f"{pt.heading}: {pt.body}" if pt.heading else pt.body
            add_textbox(
                slide, safe.left + 0.3, y,
                safe.width - 0.3, 0.4,
                f"  {text}", spec.design.font_family,
                spec.design.body_size_pt, _resolve_body_color(spec, theme),
            )
            y += 0.5

@register_renderer("timeline-horizontal")
def render_timeline(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int):
    """Horizontal timeline with nodes."""
    set_slide_background(slide, spec.design.background)
    title_h = render_title_zone(slide, spec, theme, safe)
    footer_h = render_footer(slide, spec, theme, safe, total_slides)
    content_top, content_h = get_content_zone(safe, title_h, footer_h)

    points = get_points(spec)
    n = len(points) or 1
    node_y = content_top + content_h * 0.3
    for i, pt in enumerate(points):
        x = safe.left + i * safe.width / n + safe.width / n / 2 - 0.3
        # Node circle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(node_y), Inches(0.6), Inches(0.6),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(_resolve_accent(spec, theme))
        shape.line.fill.background()
        # Label
        label = pt.heading or pt.body
        add_textbox(
            slide, x - 0.5, node_y + 0.8, 1.6, 0.8,
            label, spec.design.font_family,
            spec.design.body_size_pt - 2, _resolve_body_color(spec, theme),
            alignment=PP_ALIGN.CENTER,
        )

@register_renderer("story-card")
def render_story_card(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int):
    """Full-bleed background with text overlay."""
    set_slide_background(slide, spec.design.background)
    overlay_height = safe.height * 0.4
    overlay_top = safe.top + safe.height - overlay_height
    add_rounded_rect(
        slide, safe.left, overlay_top, safe.width, overlay_height,
        "#000000", corner_radius=0.0,
    )
    add_textbox(
        slide, safe.left + 0.5, overlay_top + 0.3,
        safe.width - 1, 0.6,
        spec.content.title, spec.design.font_family,
        spec.design.title_size_pt, "#FFFFFF",
        bold=True,
    )
    if spec.content.body_text or spec.content.key_points:
        points = get_points(spec)
        text = spec.content.body_text or "\n".join(p.body for p in points)
        add_textbox(
            slide, safe.left + 0.5, overlay_top + 1.0,
            safe.width - 1, overlay_height - 1.3,
            text, spec.design.font_family,
            spec.design.body_size_pt, "#CCCCCC",
        )

# Framework renderers
def _render_framework(slide, spec: SlideSpec, theme: dict, safe: SafeArea, n_cols: int, total_slides: int):
    """Framework layout with optional columns and card headers."""
    set_slide_background(slide, spec.design.background)
    title_h = render_title_zone(slide, spec, theme, safe)
    footer_h = render_footer(slide, spec, theme, safe, total_slides)
    content_top, content_h = get_content_zone(safe, title_h, footer_h)

    ve = _ve(theme)
    colors = theme.get("colors", {}).get("card_fills", ["#F4F4F4"] * 5)
    gap = theme.get("spacing", {}).get("element_gap_inches", 0.25)
    col_width = (safe.width - gap * (n_cols - 1)) / n_cols
    corner_r = theme.get("visual_preferences", {}).get("corner_radius_inches", 0.1)

    points = get_points(spec)
    bodies = get_point_bodies(points)
    has_headers = any(p.heading for p in points)
    header_h = ve["card_header_height_inches"] if has_headers else 0.0

    max_col_height = content_h
    body_height = compute_card_height(
        bodies, col_width, spec.design.body_size_pt, max_col_height - header_h,
    ) if bodies else 2.0
    col_height = min(body_height + header_h, max_col_height)
    col_top = content_top

    for i in range(n_cols):
        x = safe.left + i * (col_width + gap)
        fill = colors[i % len(colors)]
        add_rounded_rect(slide, x, col_top, col_width, col_height, fill, corner_r)

        pt = points[i] if i < len(points) else StructuredPoint(body="")
        h_used = render_card_header(slide, x, col_top, col_width, pt.heading, theme, spec)

        add_textbox(
            slide, x + 0.2, col_top + h_used + 0.1,
            col_width - 0.4, col_height - h_used - 0.2,
            pt.body, spec.design.font_family,
            spec.design.body_size_pt, _resolve_body_color(spec, theme),
        )

# bare "framework" 列数自适应 point 数 (1-4): 原固定 n_cols=1 时 `for i in range(n_cols)`
# 只渲染 points[0], 多 point 静默丢内容. clamp 到 4 与 framework-4col 上限一致.
_RENDERERS["framework"] = lambda s, sp, t, sa, ts: _render_framework(
    s, sp, t, sa, max(1, min(len(get_points(sp)), 4)), ts
)
_RENDERERS["framework-2col"] = lambda s, sp, t, sa, ts: _render_framework(s, sp, t, sa, 2, ts)
_RENDERERS["framework-3col"] = lambda s, sp, t, sa, ts: _render_framework(s, sp, t, sa, 3, ts)
_RENDERERS["framework-4col"] = lambda s, sp, t, sa, ts: _render_framework(s, sp, t, sa, 4, ts)

# 华为主题 renderer 注册：必须放在 _RENDERERS 所有旧条目填充完之后，避免循环 import
# （renderers_huawei 需要 register_renderer + helpers，render.py 再反 import 它）
from engine.renderers import huawei  # noqa: F401, E402  触发 18 huawei renderer 注册

# ===========================================================================
# MAIN RENDERING PIPELINE
# ===========================================================================

def _get_blank_layout(prs: Presentation):
    """Find a blank slide layout, tolerant of custom templates."""
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[len(prs.slide_layouts) - 1]


def _compute_safe_area(theme: dict) -> SafeArea:
    """计算 theme 对应的 SafeArea. 单独抽出来让 render_presentation 入口算一次后透传."""
    margin = theme.get("spacing", {}).get("slide_margin_inches", 0.5)
    return get_safe_area({"left": margin, "top": margin, "right": margin, "bottom": margin})


def render_slide(prs: Presentation, spec: SlideSpec, theme: dict, total_slides: int,
                 slide_layout=None, safe: SafeArea | None = None):
    """Render a single slide based on its visual_type.

    2026-05-19 audit fix (PERF): slide_layout + safe 改为可选参数, 由 render_presentation 入口
    一次算好透传. 老 caller 不传仍能 fallback 自己算 (向后兼容). 50 页 deck 上省 200+ 次 dict.get
    链 + N 次 layout 名 lowercase 扫描.
    """
    if slide_layout is None:
        slide_layout = _get_blank_layout(prs)
    if safe is None:
        safe = _compute_safe_area(theme)
    slide = prs.slides.add_slide(slide_layout)
    renderer = _RENDERERS.get(spec.visual_type)
    if renderer is None:
        raise ValueError(
            f"no renderer for visual_type={spec.visual_type!r} (slide {spec.id}). "
            f"Registered: {sorted(_RENDERERS.keys())}. "
            f"schemas/slide_plan.py VALID_VISUAL_TYPES 与 _RENDERERS 不同步."
        )
    renderer(slide, spec, theme, safe, total_slides)
    return slide


def render_presentation(plan: SlidePlan, theme: dict, output_path: str,
                        base_pptx: str = None, only_slides: list[int] = None):
    """Render full presentation from slide plan."""
    if base_pptx:
        prs = Presentation(base_pptx)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    if base_pptx and only_slides:
        indices_to_delete = sorted(
            [sid - 1 for sid in only_slides if 0 < sid <= len(prs.slides)],
            reverse=True,
        )
        for idx in indices_to_delete:
            delete_slide(prs, idx)

    # 一次性算 layout + safe area, 传给每个 render_slide 调用 — 避免 N 次重复 dict.get / layout 扫描
    blank_layout = _get_blank_layout(prs)
    safe_area = _compute_safe_area(theme)

    total_slides = plan.narrative.total_slides
    for spec in plan.slides:
        if only_slides and spec.id not in only_slides:
            continue
        render_slide(prs, spec, theme, total_slides, slide_layout=blank_layout, safe=safe_area)

    # Windows file lock protection — 尝试 9 个备用名 (output_v2.pptx ... output_v9.pptx)
    try:
        prs.save(output_path)
    except PermissionError as e0:
        import os
        base, ext = os.path.splitext(output_path)
        tried: list[str] = []
        for v in range(2, 10):
            alt = f"{base}_v{v}{ext}"
            tried.append(alt)
            try:
                prs.save(alt)
                print(f"Saved to {alt} (original path locked)")
                return
            except PermissionError:
                continue
        raise PermissionError(
            f"All output paths locked: {output_path} + tried {tried}. "
            f"Close PowerPoint or unlock these files and retry."
        ) from e0


def _assert_renderer_registry():
    """启动时验证 renderer 注册完整, silent fallback 改 fail-fast.

    比对集合而非总数 — 任何 visual_type ∈ VALID_VISUAL_TYPES 但缺 renderer 立即报错, 含具体名字.
    旧版本 PEP 723 双 module 实例 bug 让 huawei 18 个版式全部 silent 走 bullets fallback,
    用户拿到"标题+页码"的空白 deck 而无任何 WARN. 这个 assert 是最后一道防线.
    """
    from schemas.slide_plan import VALID_VISUAL_TYPES
    registered = set(_RENDERERS.keys())
    missing = VALID_VISUAL_TYPES - registered
    extra = registered - VALID_VISUAL_TYPES
    if missing or extra:
        raise RuntimeError(
            f"renderer registry 与 VALID_VISUAL_TYPES 不一致: "
            f"missing={sorted(missing)}, extra={sorted(extra)}. "
            f"可能原因: huawei renderer 模块未加载 / 双 module 实例 / schemas 加了 visual_type 但忘写 renderer."
        )
    print(f"[render] loaded {len(registered)} renderers", file=sys.stderr)


def _resolve_image_refs(plan: SlidePlan, base_dir: Path) -> None:
    """相对 image_ref 以 slide-plan 文件所在目录为基准解析为绝对路径.

    renderer 函数拿不到 plan 文件路径, 解析必须在这里 (唯一知道路径的层) 做完.
    绝对路径原样保留; 存在性检查留给 renderer (缺失降级无图, 不阻塞).
    """
    for spec in plan.slides:
        ref = getattr(spec.content, "image_ref", None)
        if ref:
            p = Path(ref)
            if not p.is_absolute():
                spec.content.image_ref = str((base_dir / p).resolve())


def main():
    _assert_renderer_registry()
    parser = argparse.ArgumentParser(description="Render PPTX from slide-plan.yaml")
    parser.add_argument("slide_plan", help="Path to slide-plan.yaml")
    parser.add_argument("--theme", default=None,
                        help="theme name (default: huawei if omitted)")
    parser.add_argument("--output", default="output.pptx")
    parser.add_argument("--base-pptx", default=None)
    parser.add_argument("--only-slides", default=None,
                        help="Comma-separated slide IDs to render")
    args = parser.parse_args()

    plan = SlidePlan.from_yaml(args.slide_plan)
    _resolve_image_refs(plan, Path(args.slide_plan).resolve().parent)
    theme = load_theme(args.theme)
    only = [int(x) for x in args.only_slides.split(",")] if args.only_slides else None
    render_presentation(plan, theme, args.output, args.base_pptx, only)
    print(f"Rendered {len(plan.slides)} slides to {args.output}")


if __name__ == "__main__":
    main()
