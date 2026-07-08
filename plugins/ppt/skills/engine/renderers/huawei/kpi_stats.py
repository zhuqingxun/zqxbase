"""华为主题 renderer: kpi-stats (S4 Phase 1 从 renderers_huawei.py 拆出)。

helper 单一来源在 engine.renderer_kit; 本模块仅含单个 @register_renderer。
"""
from __future__ import annotations

from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from lib.font_fallback import resolve_font_for_pptx
from lib.margins import SafeArea
from schemas.slide_plan import SlideSpec, StructuredPoint

from engine.renderer_kit import (
    register_renderer, _RENDERERS, RendererContext,
    hex_to_rgb, add_textbox, add_textbox_rich, add_rounded_rect, set_slide_background,
    render_title_zone, render_footer, get_content_zone, get_points, get_point_bodies,
    _render_cards, _resolve_accent, _resolve_title_color, _resolve_body_color, _ve,
    compute_card_height, estimate_content_height, normalize_point,
    render_card_header, render_card_number_badge, _get_card_accent,
    _truncate_for_single_line, _add_textbox_no_wrap,
    _px_pt, _px_in, _tokens, _layouts, _color, _lerp_hex, _is_dark, _type_px, _layout,
    _extra, _extra_list_merge, _estimate_text_height_in, _dget, _add_rect,
    _canvas_dims_in, _canvas_padding_in, _slide_dims_in, _font_family, TREND_COLOR_TOKEN,
)


@register_renderer("kpi-stats")
def render_kpi_stats(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """KPI 统计：N 列（默认 4）横排，每列 value(108px) + label(18px) + desc + trend(14px)。"""
    ctx = RendererContext.build(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "kpi_stats")
    columns_cfg = cfg.get("columns", 4)
    value_px = cfg.get("value_size_px", 108)  # Plan §O6 修正
    unit_px = cfg.get("unit_size_px", 32)
    label_px = cfg.get("label_size_px", 18)
    desc_px = cfg.get("desc_size_px", 16)
    trend_px = cfg.get("trend_size_px", 14)
    sep_px = cfg.get("column_separator_width_px", 2)

    sw, sh = ctx.sw, ctx.sh
    right_pad, bottom_pad = ctx.right_pad, ctx.bottom_pad
    font = ctx.font

    kpis = _extra(spec, "kpis", default=None) or []
    if not kpis:
        # 退化：从 key_points 构造
        for p in get_points(spec):
            kpis.append({
                "value": p.metric_value or "",
                "label": p.metric_label or p.heading or "",
                "desc": p.body or "",
            })
    n = min(len(kpis), columns_cfg) if kpis else columns_cfg

    # taste B (2026-05-31): value 为定性词 (不含数字) 时改"竖排三层"版式, 避免 kpi 大数字版式
    # 硬塞定性标签 → 巨红词 (AP3) + 稀疏 (AP4)。真数字 KPI (15,000/50%/21) 仍走下方横排。
    _vals = [str(_dget(k, "value", default="")) for k in kpis] if kpis else []
    _qualitative = bool(_vals) and sum(1 for v in _vals if any(c.isdigit() for c in v)) * 2 < len(_vals)
    if _qualitative:
        red = _color(theme, "primary", "#C7000B")
        ink = _color(theme, "ink", "#1F1F1F")
        ink_soft = _color(theme, "ink_soft", "#4B4B4B")
        rule = _color(theme, "rule_soft", "#ECECEC")
        x0 = safe.left
        idx_pt = _px_pt(cfg.get("tier_index_size_px", 50))   # ~30pt 序号
        term_pt = _px_pt(cfg.get("tier_value_size_px", 50))  # ~30pt 定性词 (中等, 非巨大)
        name_pt = _px_pt(cfg.get("tier_label_size_px", 33))  # ~20pt 层名
        # 标题下短红线 + 副标题 (从 variant 取, 历史只读 content.subtitle, path X 下丢)
        _add_rect(slide, x0 + 0.03, 1.42, 1.4, 0.05, red)
        _subtitle = _extra(spec, "subtitle", default=None)
        if _subtitle:
            add_textbox(slide, x0, 1.6, sw - x0 - right_pad, 0.4,
                        str(_subtitle), font, _px_pt(label_px), ink_soft)
        m = len(kpis)
        list_top = 2.25
        list_bottom = sh - bottom_pad - 0.2
        row_h = min(1.15, (list_bottom - list_top) / max(m, 1))
        for i, item in enumerate(kpis):
            if isinstance(item, str):
                item = {"value": item}
            ry = list_top + i * row_h
            cells = [
                (x0, 1.0, f"{i + 1:02d}", idx_pt, red, True),
                (x0 + 1.1, 1.7, str(_dget(item, "label", default="")), name_pt, ink, True),
                (x0 + 2.95, 2.1, str(_dget(item, "value", default="")), term_pt, red, True),
                (x0 + 5.15, sw - (x0 + 5.15) - right_pad,
                 str(_dget(item, "desc", "description", default="")), _px_pt(desc_px), ink_soft, False),
            ]
            for cx, cw, ctext, cpt, ccolor, cbold in cells:
                tb = add_textbox(slide, cx, ry, cw, row_h, ctext, font, cpt, ccolor, bold=cbold)
                tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i < m - 1:
                _add_rect(slide, x0, ry + row_h - 0.04, sw - x0 - right_pad, 0.012, rule)
        return

    area_left, area_top, area_w, area_h = ctx.area_left, ctx.area_top, ctx.area_w, ctx.area_h
    col_w = area_w / max(n, 1)

    # kpi 内容只占上半 → 整块垂直居中 + 分隔线缩到内容高, 消除下半大片留空与悬空竖线
    _has_desc = any((not isinstance(k, str)) and _dget(k, "desc", "description", default="") for k in kpis)
    kpi_content_h = min(area_h, 3.0 if _has_desc else 2.0)
    block_top = area_top + max(0.0, (area_h - kpi_content_h) / 2)

    value_pt = _px_pt(value_px)
    unit_pt = _px_pt(unit_px)
    label_pt = _px_pt(label_px)
    desc_pt = _px_pt(desc_px)
    trend_pt = _px_pt(trend_px)
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    primary = _color(theme, "primary", "#C7000B")

    for i in range(n):
        item: Any = kpis[i] if i < len(kpis) else {}
        if isinstance(item, str):
            item = {"value": item, "label": ""}
        x = area_left + i * col_w
        cy = block_top

        value_text = str(_dget(item, "value", default=""))
        unit_text = str(_dget(item, "unit", default=""))
        # L-3 fix: value 强制单行 (KPI 数字本质单行展示), 字号超 col_w 时自适应缩放,
        # 避免 word_wrap=True 让长 value (如 "15,000" + unit) 折行后压住下方 label.
        # 算 value+unit 总宽度等价中文字数, 超 col_w 时按比例缩字号
        cn_v = sum(1 for ch in value_text if "一" <= ch <= "鿿")
        wt_v = cn_v + (len(value_text) - cn_v) * 0.55
        cn_u = sum(1 for ch in unit_text if "一" <= ch <= "鿿")
        wt_u = cn_u + (len(unit_text) - cn_u) * 0.55
        # 留 25% buffer 防 LibreOffice/PowerPoint 实际渲染字宽比理论大 (中文等宽字体在 PowerPoint 实测约 1.1-1.2 em)
        value_max_w = (col_w - 0.2) * 0.85
        # value 主体宽 = wt_v * value_pt/72; unit 宽 = wt_u * unit_pt/72
        # 不超 value_max_w 时维持 default; 超时按比例缩 value_pt (留 buffer)
        full_w = wt_v * value_pt / 72 + (wt_u + 1) * unit_pt / 72  # +1 char 间距
        if full_w > value_max_w and full_w > 0:
            scale = value_max_w / full_w
            value_pt_eff = max(20.0, value_pt * scale)
            unit_pt_eff = max(12.0, unit_pt * scale)
        else:
            value_pt_eff = value_pt
            unit_pt_eff = unit_pt
        value_actual_h = value_pt_eff / 72 * 1.25 + 0.15
        if unit_text:
            parts = [
                {"text": value_text, "size_pt": value_pt_eff, "color": primary, "bold": True},
                {"text": f" {unit_text}", "size_pt": unit_pt_eff, "color": primary, "bold": False},
            ]
            txbox = add_textbox_rich(slide, x, cy, col_w, value_actual_h, parts, font)
            txbox.text_frame.word_wrap = False  # 关键: 强制单行
        else:
            txbox = add_textbox(
                slide, x, cy, col_w, value_actual_h,
                value_text, font, value_pt_eff,
                primary, bold=True,
            )
            txbox.text_frame.word_wrap = False
        cy += value_actual_h + max(0.08, value_pt_eff / 72 * 0.12)  # gutter 按字号缩放

        label_text = str(_dget(item, "label", default=""))
        if label_text:
            label_actual_h = max(
                label_pt / 72 * 1.4 + 0.1,
                _estimate_text_height_in(label_text, col_w, label_pt, line_spacing=1.4, padding_in=0.08),
            )
            add_textbox(
                slide, x, cy, col_w, label_actual_h,
                label_text, font, label_pt,
                ink, bold=True,
            )
            cy += label_actual_h + 0.05

        desc_text = str(_dget(item, "desc", "description", default=""))
        if desc_text:
            desc_actual_h = max(
                desc_pt / 72 * 1.5 + 0.3,
                _estimate_text_height_in(desc_text, col_w, desc_pt, line_spacing=1.4, padding_in=0.15),
            )
            add_textbox(
                slide, x, cy, col_w, desc_actual_h,
                desc_text, font, desc_pt,
                ink_soft,
            )
            cy += desc_actual_h + 0.1

        trend_text = str(_dget(item, "trend_text", default=""))
        trend_dir = _dget(item, "trend", "trend_direction", default="")
        # flat 不画箭头 (避免"→"指向空白的悬空箭头, 2026-05-31 taste 诊断 AP4 项)
        arrow = {"up": "↑", "down": "↓"}.get(str(trend_dir).lower(), "")
        if trend_text or arrow:
            token = TREND_COLOR_TOKEN.get(str(trend_dir).lower(), "ink_soft")
            tcolor = _color(theme, token, ink_soft)
            add_textbox(
                slide, x, cy, col_w, trend_pt / 72 * 1.5 + 0.15,
                f"{arrow} {trend_text}".strip(), font, trend_pt,
                tcolor,
            )

        # 列分隔线（非末列）
        if i < n - 1:
            sep_x = x + col_w - _px_in(sep_px) / 2
            _add_rect(
                slide, sep_x, block_top, max(_px_in(sep_px), 0.01), kpi_content_h,
                _color(theme, "rule", "#D9D9D9"),
            )


