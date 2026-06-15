"""华为主题 renderer: swot (S4 Phase 1 从 renderers_huawei.py 拆出)。

helper 单一来源在 engine.renderer_kit; 本模块仅含单个 @register_renderer。
"""
from __future__ import annotations

from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from lib.font_fallback import resolve_font_for_pptx
from lib.margins import SafeArea
from schemas.slide_plan import SlideSpec, StructuredPoint

from engine.renderer_kit import (
    register_renderer, _RENDERERS,
    hex_to_rgb, add_textbox, add_textbox_rich, add_rounded_rect, set_slide_background,
    render_title_zone, render_footer, get_content_zone, get_points, get_point_bodies,
    _render_cards, _resolve_accent, _resolve_title_color, _resolve_body_color, _ve,
    compute_card_height, estimate_content_height, normalize_point,
    render_card_header, render_card_number_badge, _get_card_accent,
    _truncate_for_single_line, _add_textbox_no_wrap,
    _px_pt, _px_in, _tokens, _layouts, _color, _lerp_hex, _is_dark, _type_px, _layout,
    _extra, _extra_list_merge, _estimate_text_height_in, _dget, _add_rect,
    _canvas_dims_in, _canvas_padding_in, _slide_dims_in, _font_family, TREND_COLOR_TOKEN,
)


@register_renderer("swot")
def render_swot(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """SWOT 2x2：四象限 (S/W/O/T) + 中心红十字分隔线。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "swot")
    cross_w_px = cfg.get("cross_line_width_px", 2)
    cross_token = cfg.get("cross_line_token", "primary")

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5
    cell_w = area_w / 2
    cell_h = area_h / 2

    primary = _color(theme, "primary", "#C7000B")
    cross_color = _color(theme, cross_token, primary)
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    paper_2 = _color(theme, "paper_2", "#F4F4F4")

    # (legacy_key 顶层 list, short_key=schema quadrants 子键, letter 角标, default_label 默认标题)
    quadrant_map = [
        ("strengths", "s", "S", "优势"),
        ("weaknesses", "w", "W", "劣势"),
        ("opportunities", "o", "O", "机会"),
        ("threats", "t", "T", "威胁"),
    ]
    body_pt = _px_pt(_tokens(theme).get("type_scale_px", {}).get("body", 24))
    heading_pt = _px_pt(_tokens(theme).get("type_scale_px", {}).get("h4", 28))

    # path X: variant.quadrants.{s,w,o,t} (SwotQuadrant: heading + items)。
    # 历史 bug: renderer 只读顶层 strengths/weaknesses/...(path Y), schema 的嵌套 quadrants 静默丢失,
    # path X 下四象限正文全空。现优先消费 quadrants 嵌套, 缺失时退回顶层 list。
    quadrants_obj = _extra(spec, "quadrants", default=None)

    for i, (legacy_key, short_key, letter, default_label) in enumerate(quadrant_map):
        row, col = divmod(i, 2)
        x = area_left + col * cell_w
        y = area_top + row * cell_h
        _add_rect(slide, x, y, cell_w, cell_h, paper_2)

        # 取该象限的 items + 可选 heading 覆盖
        items = None
        heading_override = None
        if quadrants_obj is not None:
            q = _dget(quadrants_obj, short_key, default=None)
            if q is not None:
                items = _dget(q, "items", default=None)
                heading_override = _dget(q, "heading", default=None)
        if items is None:
            items = _extra(spec, legacy_key, default=None)
        items = items or []
        if isinstance(items, str):
            items = [items]
        label = str(heading_override) if heading_override else default_label

        add_textbox(
            slide, x + 0.2, y + 0.15, 1.0, heading_pt / 72 * 1.3 + 0.1,
            letter, font, heading_pt * 1.5,
            primary, bold=True,
        )
        add_textbox(
            slide, x + 1.2, y + 0.3, cell_w - 1.4, heading_pt / 72 * 1.3 + 0.1,
            label, font, heading_pt,
            ink, bold=True,
        )
        text = "\n".join(f"• {str(it)}" for it in items[:5]) if items else ""
        if text:
            add_textbox(
                slide, x + 0.3, y + 1.1, cell_w - 0.5, cell_h - 1.3,
                text, font, body_pt,
                ink_soft,
            )

    # 中心红十字线
    line_pt = max(cross_w_px * 0.5, 1.0)
    # 横线
    hl = slide.shapes.add_connector(
        1,
        Inches(area_left), Inches(area_top + cell_h),
        Inches(area_left + area_w), Inches(area_top + cell_h),
    )
    hl.line.color.rgb = hex_to_rgb(cross_color)
    hl.line.width = Pt(line_pt)
    # 竖线
    vl = slide.shapes.add_connector(
        1,
        Inches(area_left + cell_w), Inches(area_top),
        Inches(area_left + cell_w), Inches(area_top + area_h),
    )
    vl.line.color.rgb = hex_to_rgb(cross_color)
    vl.line.width = Pt(line_pt)


