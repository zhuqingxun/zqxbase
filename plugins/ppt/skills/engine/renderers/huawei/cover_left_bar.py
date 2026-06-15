"""华为主题 renderer: cover-left-bar (S4 Phase 1 从 renderers_huawei.py 拆出)。

helper 单一来源在 engine.renderer_kit; 本模块仅含单个 @register_renderer。
"""
from __future__ import annotations

from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from lib.font_fallback import resolve_font_for_pptx
from lib.margins import SafeArea
from schemas.slide_plan import SlideSpec, StructuredPoint

from engine.renderer_kit import (
    register_renderer, _RENDERERS,
    hex_to_rgb, add_textbox, add_textbox_rich, add_rounded_rect, set_slide_background,
    render_title_zone, render_footer, get_content_zone, get_points, get_point_bodies,
    _render_cards, _resolve_accent, _resolve_title_color, _resolve_body_color, _ve,
    compute_card_height, estimate_content_height, normalize_point,
    render_card_header, render_card_number_badge, _get_card_accent,
    _truncate_for_single_line, _add_textbox_no_wrap,
    _px_pt, _px_in, _tokens, _layouts, _color, _lerp_hex, _is_dark, _type_px, _layout,
    _extra, _extra_list_merge, _estimate_text_height_in, _dget, _add_rect,
    _canvas_dims_in, _canvas_padding_in, _slide_dims_in, _font_family, TREND_COLOR_TOKEN,
)


def _add_cover_image(slide, spec: SlideSpec) -> None:
    """cover 消费 content.image_ref: 右侧区插入图片 (codex-fusion F3, 封面摄影版).

    - 必须先于文字框调用 (python-pptx 无 z-order API, 先加的 shape 在底层, 文字保持在上)
    - 位置取 design.image_position, 未提供时用默认右侧区 (避让左栏文字)
    - 等比缩放装进指定框 (用 add_picture 后的 native EMU 尺寸算 scale, 不依赖 PIL)
    - 文件缺失 → stderr 警告 + 渲染无图 cover (失败不阻塞, PRD F3 约束)
    """
    import sys as _sys
    from pathlib import Path as _Path

    path = _Path(spec.content.image_ref)
    if not path.is_file():
        print(f"[render] WARNING: cover image_ref 不存在, 渲染无图封面: {path}",
              file=_sys.stderr)
        return
    pos = spec.design.image_position
    if pos is not None:
        x, y, w, h = pos.x, pos.y, pos.width, pos.height
    else:
        x, y, w, h = 6.2, 1.2, 6.3, 5.1  # 默认右侧区 (slide 13.333"x7.5")
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y))
    scale = min(Inches(w) / pic.width, Inches(h) / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    # 框内居中
    pic.left = Inches(x) + int((Inches(w) - pic.width) / 2)
    pic.top = Inches(y) + int((Inches(h) - pic.height) / 2)


@register_renderer("cover-left-bar")
def render_cover_left_bar(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """封面（浅底无边条版, 2026-05-31 taste 选型 A 去红条）: eyebrow + 克制标题 + 副标题 + 底部 meta 栏。

    历史: 旧版有 full-height 左红条 (bar_width_px) + 108px 巨标题。用户全局规则: huawei 风格
    不要贯穿边缘的细红条; 标题降到 ~52pt; 红色仅留在 eyebrow 文字。构图垂直居中偏上, 底部 meta。
    """
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))

    # codex-fusion F3: image_ref 非空才进图片分支 (无 image_ref 时零行为变化)
    if spec.content.image_ref:
        _add_cover_image(slide, spec)

    cfg = _layout(theme, "cover_left_bar")
    title_px = cfg.get("title_size_px", 88)        # ~52pt (旧 108px)
    subtitle_px = cfg.get("subtitle_size_px", 36)
    eyebrow_px = cfg.get("eyebrow_size_px", 26)
    meta_label_px = cfg.get("meta_label_size_px", 20)
    meta_value_px = cfg.get("meta_size_px", 26)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)
    red = _color(theme, "primary", "#C7000B")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    ink_mute = _color(theme, "ink_mute", "#9A9A9A")
    rule = _color(theme, "rule", "#D9D9D9")

    x = left_pad
    body_w = sw - x - right_pad

    # eyebrow + 标题 + 副标题 (垂直居中偏上)
    eyebrow = _extra(spec, "eyebrow", default=spec.chapter or "")
    y = 2.0
    if eyebrow:
        eb_pt = _px_pt(eyebrow_px)
        eb_h = eb_pt / 72 * 1.4 + 0.1
        add_textbox(slide, x, y, body_w, eb_h, str(eyebrow), font, eb_pt, red, bold=True)
        y += eb_h + 0.15

    title_pt = _px_pt(title_px)
    title_h = max(
        title_pt / 72 * 1.3 + 0.2,
        _estimate_text_height_in(spec.content.title or "", body_w, title_pt, line_spacing=1.1, padding_in=0.15),
    )
    add_textbox(slide, x, y, body_w, title_h, spec.content.title or "", font, title_pt, ink, bold=True)
    y += title_h + 0.25

    # emphasis: 红色强调短语 (variant CoverLeftBarContent.emphasis), 历史未渲染; 置于标题与副标题之间
    emphasis_text = _extra(spec, "emphasis", default=None)
    if emphasis_text:
        emph_pt = _px_pt(cfg.get("emphasis_size_px", subtitle_px))
        emph_h = emph_pt / 72 * 1.5 + 0.15
        add_textbox(slide, x, y, body_w, emph_h, str(emphasis_text), font, emph_pt, red, bold=True)
        y += emph_h + 0.1

    # subtitle 从 variant 取 (历史只读 content.subtitle, path X 下 subtitle 在 variant 里被丢)
    subtitle_text = _extra(spec, "subtitle", default=None)
    if subtitle_text:
        sub_pt = _px_pt(subtitle_px)
        add_textbox(slide, x, y, body_w, sub_pt / 72 * 1.6 + 0.2,
                    str(subtitle_text), font, sub_pt, ink_soft)

    # 底部 meta 栏 (label/value 竖排, 上方细线)
    meta = (_extra(spec, "meta", default=None) or _extra(spec, "meta_items", default=None) or [])[:4]
    if meta:
        rule_y = sh - bottom_pad - 1.3
        _add_rect(slide, x, rule_y, body_w, 0.012, rule)
        meta_y = rule_y + 0.25
        col_w = body_w / max(len(meta), 1)
        lbl_pt = _px_pt(meta_label_px)
        val_pt = _px_pt(meta_value_px)
        for i, item in enumerate(meta):
            label_v = _dget(item, "label", default="")
            value_v = _dget(item, "value", default="")
            if not (label_v or value_v) and isinstance(item, str):
                value_v = item
            cx = x + i * col_w
            if label_v:
                add_textbox(slide, cx, meta_y, col_w - 0.3, 0.3, str(label_v), font, lbl_pt, ink_mute, bold=True)
            add_textbox(slide, cx, meta_y + 0.3, col_w - 0.3, 0.4, str(value_v), font, val_pt, ink_soft)


