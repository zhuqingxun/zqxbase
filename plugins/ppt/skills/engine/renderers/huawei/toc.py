"""华为主题 renderer: toc (S4 Phase 1 从 renderers_huawei.py 拆出)。

helper 单一来源在 engine.renderer_kit; 本模块仅含单个 @register_renderer。
"""
from __future__ import annotations

from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from lib.font_fallback import resolve_font_for_pptx
from lib.margins import SafeArea
from schemas.slide_plan import SlideSpec, StructuredPoint

from engine.renderer_kit import (
    register_renderer, _RENDERERS,
    hex_to_rgb, add_textbox, add_textbox_rich, add_rounded_rect, set_slide_background,
    render_title_zone, render_footer, get_content_zone, get_points, get_point_bodies,
    _render_cards, _resolve_accent, _resolve_title_color, _resolve_body_color, _ve,
    compute_card_height, estimate_content_height, normalize_point,
    render_card_header, render_card_number_badge, _get_card_accent,
    _truncate_for_single_line, _add_textbox_no_wrap,
    _px_pt, _px_in, _tokens, _layouts, _color, _lerp_hex, _is_dark, _type_px, _layout,
    _extra, _extra_list_merge, _estimate_text_height_in, _dget, _add_rect,
    _canvas_dims_in, _canvas_padding_in, _slide_dims_in, _font_family, TREND_COLOR_TOKEN,
)


@register_renderer("toc")
def render_toc(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """目录（极简版, 2026-05-31 taste 选型 C 去页码）: 克制标题 + 短红线 + 红编号列表 (无巨号/无页码/无分隔线)。

    历史: 旧版左侧 300px 巨号"目录" (AP3 巨型装饰字)。现重排为顶部克制标题 ("目录 · CONTENTS")
    + 红色编号 (强调但不巨大) + 章节标题列表; 按用户选型去掉页码列与行分隔线。
    """
    sw, sh = _slide_dims_in()
    cfg = _layout(theme, "toc")
    set_slide_background(slide, _color(theme, cfg.get("background_token", "paper"), "#FFFFFF"))

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)
    red = _color(theme, "primary", "#C7000B")
    ink = _color(theme, "ink", "#1F1F1F")

    header_pt = _px_pt(cfg.get("header_size_px", 66))   # ~40pt (旧 giant 300px 违反 P3)
    number_pt = _px_pt(cfg.get("number_size_px", 57))   # ~34pt 红编号 (强调但不巨大)
    title_pt = _px_pt(cfg.get("item_size_px", 33))      # ~20pt 章节标题

    x = left_pad + 0.3
    # 标题 + 可选英文后缀
    base_title = _extra(spec, "giant_label", default=None) or spec.content.title or "目录"
    suffix = cfg.get("header_suffix", "· CONTENTS")
    header_text = f"{base_title} {suffix}".strip() if suffix else base_title
    add_textbox(
        slide, x, 0.62, sw - x - right_pad, header_pt / 72 * 1.4 + 0.1,
        header_text, font, header_pt, ink, bold=True,
    )
    _add_rect(slide, x + 0.03, 1.5, 1.4, 0.05, red)

    # 章节列表 (编号 + 标题, 无页码无分隔线)
    items = _extra(spec, "items", default=None) or _extra(spec, "chapters", default=None) or []
    if not items:
        for p in get_points(spec):
            items.append({"title": p.heading or p.body})

    n = max(len(items), 1)
    list_top = 2.0
    list_bottom = sh - bottom_pad - 0.1
    row_h = min(0.82, (list_bottom - list_top) / n)
    num_col = 1.5
    title_w = sw - (x + num_col) - right_pad
    for i, item in enumerate(items):
        cy = list_top + i * row_h
        title_text = item if isinstance(item, str) else str(_dget(item, "title", default=""))
        number_text = f"{i + 1:02d}"
        # 编号 (红, 垂直居中)
        nbx = slide.shapes.add_textbox(Inches(x), Inches(cy), Inches(num_col), Inches(row_h))
        ntf = nbx.text_frame
        ntf.word_wrap = False
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ntf.margin_top = 0
        ntf.margin_bottom = 0
        npa = ntf.paragraphs[0]
        npa.text = number_text
        npa.font.name = font
        npa.font.size = Pt(number_pt)
        npa.font.bold = True
        npa.font.color.rgb = hex_to_rgb(red)
        # 标题 (ink, 垂直居中对齐编号)
        tbx = slide.shapes.add_textbox(Inches(x + num_col), Inches(cy), Inches(title_w), Inches(row_h))
        ttf = tbx.text_frame
        ttf.word_wrap = True
        ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ttf.margin_top = 0
        ttf.margin_bottom = 0
        tpa = ttf.paragraphs[0]
        tpa.text = title_text
        tpa.font.name = font
        tpa.font.size = Pt(title_pt)
        tpa.font.color.rgb = hex_to_rgb(ink)


