"""华为主题 renderer: heatmap-matrix (S4 Phase 1 从 renderers_huawei.py 拆出)。

helper 单一来源在 engine.renderer_kit; 本模块仅含单个 @register_renderer。
"""
from __future__ import annotations

from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from lib.font_fallback import resolve_font_for_pptx
from lib.margins import SafeArea
from schemas.slide_plan import SlideSpec, StructuredPoint

from engine.renderer_kit import (
    register_renderer, _RENDERERS, RendererContext,
    hex_to_rgb, add_textbox, add_textbox_rich, add_rounded_rect, set_slide_background,
    render_title_zone, render_footer, get_content_zone, get_points, get_point_bodies,
    _render_cards, _resolve_accent, _resolve_title_color, _resolve_body_color, _ve,
    compute_card_height, estimate_content_height, normalize_point,
    render_card_header, render_card_number_badge, _get_card_accent,
    _truncate_for_single_line, _add_textbox_no_wrap,
    _px_pt, _px_in, _tokens, _layouts, _color, _lerp_hex, _is_dark, _type_px, _layout,
    _extra, _extra_list_merge, _estimate_text_height_in, _dget, _add_rect,
    _canvas_dims_in, _canvas_padding_in, _slide_dims_in, _font_family, TREND_COLOR_TOKEN,
)


@register_renderer("heatmap-matrix")
def render_heatmap_matrix(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """热力图矩阵：N+1 行 × (cols+2) 列表格。

    Unicode 方案 A：score_chars = 5 个 U+2588 Full Block，按 score 截取前 N 个显示为红色。
    末列为 total 合计。
    """
    ctx = RendererContext.build(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "heatmap_matrix")
    row_label_col_px = cfg.get("row_label_column_px", 220)
    total_col_px = cfg.get("total_column_px", 220)
    # Unicode 方案 A：5 个 U+2588 Full Block。fallback 兜底防 layouts 缺字段或误填退格符。
    score_chars = cfg.get("score_chars") or "█████"
    # 0.5 档精度: 小数部分 >= 0.5 时用 U+258C Left Half Block 占一格 (如 3.5 → ███▌░)。
    half_char = cfg.get("half_block_char") or "▌"

    font = ctx.font
    area_left, area_top, area_w, area_h = ctx.area_left, ctx.area_top, ctx.area_w, ctx.area_h

    rows_data = _extra(spec, "rows", default=None) or []
    columns = _extra(spec, "columns", default=None) or []
    if not rows_data:
        for p in get_points(spec):
            rows_data.append({
                "label": p.heading or "",
                "scores": [0, 0, 0, 0],
                "total": p.body or "",
            })
    if not columns:
        columns = [f"C{i + 1}" for i in range(len(rows_data[0].get("scores", [])) if rows_data else 4)]

    n_rows = len(rows_data) + 1  # +1 header
    n_cols = len(columns) + 2    # label + N + total

    # 列宽
    row_label_w = _px_in(row_label_col_px)
    total_w = _px_in(total_col_px)
    middle_w = max(area_w - row_label_w - total_w, 1.0)
    mid_col_w = middle_w / max(len(columns), 1)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(area_left), Inches(area_top),
        Inches(area_w), Inches(min(area_h, n_rows * 0.55)),
    )
    tbl = table_shape.table
    # 设列宽
    tbl.columns[0].width = Emu(int(row_label_w * 914400))
    for ci in range(len(columns)):
        tbl.columns[ci + 1].width = Emu(int(mid_col_w * 914400))
    tbl.columns[n_cols - 1].width = Emu(int(total_w * 914400))

    primary = _color(theme, "primary", "#C7000B")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_mute = _color(theme, "ink_mute", "#9A9A9A")
    paper_3 = _color(theme, "paper_3", "#EAEAEA")
    body_pt = _px_pt(_type_px(theme, "body", 24))
    small_pt = _px_pt(_type_px(theme, "small", 20))

    def _set_cell(cell, text: str, color_hex: str, bold: bool = False, fill_hex: str | None = None,
                  size_pt: float = body_pt, align=PP_ALIGN.CENTER) -> None:
        if fill_hex:
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(fill_hex)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        p = tf.paragraphs[0]
        p.alignment = align
        if p.runs:
            run = p.runs[0]
            run.font.name = font
            run.font.size = Pt(size_pt)
            run.font.color.rgb = hex_to_rgb(color_hex)
            run.font.bold = bold

    # Header 行
    _set_cell(tbl.cell(0, 0), "", ink, fill_hex=paper_3, bold=True, size_pt=small_pt)
    for ci, col_name in enumerate(columns):
        _set_cell(tbl.cell(0, ci + 1), str(col_name), ink, bold=True,
                  fill_hex=paper_3, size_pt=small_pt)
    _set_cell(tbl.cell(0, n_cols - 1), "合计", ink, bold=True,
              fill_hex=paper_3, size_pt=small_pt)

    # 数据行
    for ri, row in enumerate(rows_data):
        if isinstance(row, str):
            row = {"label": row, "scores": [0] * len(columns), "total": ""}
        _set_cell(tbl.cell(ri + 1, 0), str(_dget(row, "label", default="")), ink,
                  bold=True, align=PP_ALIGN.LEFT, size_pt=small_pt)
        scores = _dget(row, "scores", default=[]) or []
        for ci in range(len(columns)):
            # 支持 0.5 档: score 可为 float (3.5). 非数值/缺失回退 0, clamp 到 [0, 5].
            try:
                sc = float(scores[ci]) if ci < len(scores) else 0.0
            except (ValueError, TypeError):
                sc = 0.0
            sc = max(0.0, min(sc, 5.0))
            full = int(sc)                      # 整数档满块数
            has_half = (sc - full) >= 0.5       # 小数 >=0.5 追加一个半块
            filled = score_chars[:full] + (half_char if has_half else "")
            empty = "░" * (5 - full - (1 if has_half else 0))  # U+2591 Light Shade 补满 5 格
            cell = tbl.cell(ri + 1, ci + 1)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            if filled:
                r1 = p.runs[0] if p.runs else p.add_run()
                r1.text = filled
                r1.font.name = font
                r1.font.size = Pt(body_pt)
                r1.font.color.rgb = hex_to_rgb(primary)
                r1.font.bold = True
            if empty:
                r2 = p.add_run()
                r2.text = empty
                r2.font.name = font
                r2.font.size = Pt(body_pt)
                r2.font.color.rgb = hex_to_rgb(ink_mute)
        total_text = str(_dget(row, "total", default=""))
        _set_cell(tbl.cell(ri + 1, n_cols - 1), total_text, ink_mute,
                  align=PP_ALIGN.LEFT, size_pt=small_pt)


