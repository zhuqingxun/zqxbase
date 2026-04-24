# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=1.0.0", "pydantic>=2.0", "pyyaml>=6.0", "lxml>=4.9"]
# ///
"""test_render.py: 回归测试 — 为每种 visual type 生成最小 slide 并渲染。

验证:
1. 所有 visual type 均无崩溃
2. 输出 slide 数与 plan 一致
3. 每页至少有 1 个 shape

Usage:
    uv run --script engine/test_render.py
    uv run --script engine/test_render.py --keep   # 保留生成的 pptx 供检查
"""

from __future__ import annotations

import argparse
import sys
import tempfile
from datetime import datetime
from pathlib import Path

import yaml
from pptx import Presentation

# PEP 723: resolve sibling packages
SKILLS_DIR = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(SKILLS_DIR))

from engine.render import render_presentation, load_theme, _RENDERERS
from schemas.slide_plan import SlidePlan, SlideSpec, SlideRole, SlideContent, SlideDesign, PlanMeta, Narrative

# ---------------------------------------------------------------------------
# 测试夹具：为每种 visual type 生成最小 SlideSpec
# ---------------------------------------------------------------------------

_SAMPLE_POINTS_3 = [
    "第一个要点: 这是用于回归测试的示例内容，包含足够的中文字符来验证渲染引擎的文本处理能力和容器高度计算。",
    "第二个要点: 回归测试需要覆盖所有 visual type，确保每种渲染器都能正常工作，不会因为边界条件而崩溃。",
    "第三个要点: 内容量达到 80 字以上的 key_point 才能通过 validate_plan 的内容量门禁检查，因此测试数据也需满足此要求。",
]

_SAMPLE_POINTS_5 = _SAMPLE_POINTS_3 + [
    "第四个要点: 额外的测试内容用于验证 4 列和 5 列布局的渲染器能正确处理不同数量的卡片和容器高度自适应。",
    "第五个要点: 最后一个要点确保 cards-5、comparison-5、process-5-phase 等宽布局的渲染器在极端情况下不会溢出。",
]

_COMPARISON_POINTS_2 = [
    "方案A — 优势: 这是用于对比型 visual type 的测试内容，冒号分隔标题和正文，验证 header/body 拆分逻辑。正文部分需要足够长以测试换行。",
    "方案B — 劣势: 另一个对比维度的测试内容，同样使用冒号分隔格式，确保两列渲染器能正确处理不同长度的文本内容并垂直居中。",
]


def _make_spec(slide_id: int, visual_type: str) -> SlideSpec:
    """为指定 visual type 生成最小 SlideSpec。"""
    content = SlideContent(title=f"测试: {visual_type}")

    # 根据 visual type 填充合适的 content 字段
    if visual_type in ("hero-statement", "story-card"):
        content.subtitle = "副标题测试文本 — 回归测试自动生成"
    elif visual_type == "quote-hero":
        content.title = "这是一段用于测试引言渲染器的引言文本，需要足够长以验证文字居中和换行效果"
        content.subtitle = "测试引言来源"
    elif visual_type == "table":
        from schemas.slide_plan import TableData
        content.table_data = TableData(
            headers=["列A", "列B", "列C"],
            rows=[["数据1", "数据2", "数据3"], ["数据4", "数据5", "数据6"]],
        )
    elif visual_type.startswith("comparison"):
        n = int(visual_type.split("-")[1]) if "-" in visual_type and visual_type.split("-")[1].isdigit() else 2
        content.key_points = (_COMPARISON_POINTS_2 * 3)[:n]
    elif visual_type == "data-contrast":
        content.key_points = ["2024年: 增长率 +15.3%", "2023年: 增长率 +8.7%"]
    elif any(visual_type.startswith(p) for p in ("cards-", "process-", "framework")):
        # 从 visual type 名推断列数
        parts = visual_type.replace("-phase", "").split("-")
        try:
            n = int(parts[-1].replace("col", ""))
        except ValueError:
            n = 3
        content.key_points = (_SAMPLE_POINTS_5)[:n]
    else:
        # bullets, comparison-tables, timeline-horizontal, 兜底
        content.key_points = _SAMPLE_POINTS_3

    role = SlideRole.TITLE if visual_type == "hero-statement" else SlideRole.CONTENT
    return SlideSpec(
        id=slide_id,
        role=role,
        chapter="回归测试",
        visual_type=visual_type,
        content=content,
        design=SlideDesign(),
    )


# ---------------------------------------------------------------------------
# 测试执行
# ---------------------------------------------------------------------------

def run_tests(keep: bool = False) -> bool:
    """运行全部 visual type 回归测试。返回 True 表示全部通过。"""
    visual_types = sorted(_RENDERERS.keys())
    print(f"回归测试: {len(visual_types)} 种 visual type\n")

    slides = [_make_spec(i + 1, vt) for i, vt in enumerate(visual_types)]
    plan = SlidePlan(
        meta=PlanMeta(
            title="回归测试",
            preset="research-report",
            theme="clean-light",
            generated_at=datetime.now(),
        ),
        narrative=Narrative(
            thesis="回归测试",
            arc="test",
            total_slides=len(slides),
        ),
        slides=slides,
    )

    theme = load_theme("clean-light")
    output_path = str(SKILLS_DIR / "engine" / "_test_regression.pptx") if keep else tempfile.mktemp(suffix=".pptx")

    # 渲染
    try:
        render_presentation(plan, theme, output_path)
    except Exception as e:
        print(f"FAIL  渲染崩溃: {e}")
        return False

    # 验证
    prs = Presentation(output_path)
    actual_slides = len(prs.slides)
    expected_slides = len(visual_types)
    all_pass = True

    if actual_slides != expected_slides:
        print(f"FAIL  Slide 数量: 期望 {expected_slides}, 实际 {actual_slides}")
        all_pass = False

    for i, (slide, vt) in enumerate(zip(prs.slides, visual_types)):
        shape_count = len(slide.shapes)
        if shape_count == 0:
            print(f"FAIL  Slide {i+1} ({vt}): 0 shapes")
            all_pass = False
        else:
            print(f"PASS  Slide {i+1:2d} ({vt:25s}): {shape_count} shapes")

    print(f"\n{'='*50}")
    if all_pass:
        print(f"ALL PASS — {actual_slides} slides, {len(visual_types)} visual types")
    else:
        print("SOME TESTS FAILED")

    if keep:
        print(f"输出保留在: {output_path}")
    else:
        Path(output_path).unlink(missing_ok=True)

    return all_pass


def main():
    parser = argparse.ArgumentParser(description="render.py 回归测试")
    parser.add_argument("--keep", action="store_true", help="保留生成的 pptx")
    args = parser.parse_args()
    ok = run_tests(keep=args.keep)
    sys.exit(0 if ok else 1)


if __name__ == "__main__":
    main()
