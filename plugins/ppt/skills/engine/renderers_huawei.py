"""17 个华为主题 renderer + cards-6 复用。

模式：@register_renderer 装饰器 + 从 theme['tokens'] 读色值 / theme['_layouts'] 读布局。
禁止硬编码色值/字号（必须从 tokens/layouts 读；layouts 缺失字段时用 Plan §O6 表的决议值作为 fallback）。

对应 PRD §7.1 V01-V17 共 17 个新 visual_type（加上 V17 的变体汇总）。
模板参考：themes/huawei/reference/templates/<NN>-*.html
"""
from __future__ import annotations

from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from engine.render import (
    register_renderer, hex_to_rgb, add_textbox, add_textbox_rich,
    add_rounded_rect, set_slide_background,
    render_title_zone, render_footer, get_content_zone, get_points,
    _render_cards,
)
from lib.font_fallback import resolve_font_for_pptx
from lib.margins import SafeArea
from schemas.slide_plan import SlideSpec, StructuredPoint

# ===========================================================================
# 单位/读值工具
# ===========================================================================

# px → pt 换算（研究报告 §问题 3 验证公式）
# 画布参考 1920×1080 对 slide 13.333"×7.5"，pt ≈ px × 0.6
def _px_pt(px: int | float) -> float:
    return float(px) * 0.6


# px → inch 换算（画布 1920 px → slide 13.333 inch → 1 inch ≈ 144 px）
# 为兼容研究报告 "/96" 的 Web DPI 表述，此处保留 96 DPI 换算便于对照。
def _px_in(px: int | float) -> float:
    return float(px) / 96.0


def _tokens(theme: dict) -> dict:
    return theme.get("tokens", theme)


def _layouts(theme: dict) -> dict:
    return theme.get("_layouts", theme.get("layouts", {}))


def _color(theme: dict, token_name: str, fallback: str = "#000000") -> str:
    return _tokens(theme).get("colors", {}).get(token_name, fallback)


def _type_px(theme: dict, name: str, fallback_px: int = 24) -> int:
    return _tokens(theme).get("type_scale_px", {}).get(name, fallback_px)


def _layout(theme: dict, visual_type_key: str) -> dict:
    return _layouts(theme).get(visual_type_key, {})


def _extra(source, *fields: str, default=None):
    """读 extra 字段 (path X 双查兼容 + 多 alias 字段名).

    优先级 (按 fields 顺序逐一尝试, 第一个非空就返回):
    1. 若 source 有 .variant 属性 (即 SlideSpec, path X 推荐): 先查 variant 直属字段, 再查 variant.__pydantic_extra__
    2. 回退到 content (SlideSpec.content 或 source 自身, path Y 兜底): 查直属, 再查 __pydantic_extra__

    支持单字段: `_extra(spec, "title")`, 也支持多 alias: `_extra(spec, "number", "big_number")`
    (按顺序找第一个非空, schemas 字段名与 renderer 历史读的字段名共存时使用)。

    Args:
        source: SlideSpec 实例 (推荐) 或 SlideContent 实例 (兼容)
        *fields: 字段名, 支持多个 alias 按顺序查找
        default: 兜底值

    Returns:
        字段值, 全部缺失时返回 default
    """
    variant = getattr(source, "variant", None)
    content = getattr(source, "content", source)
    for field in fields:
        if variant is not None:
            v = getattr(variant, field, None)
            if v is not None:
                return v
            vextras = getattr(variant, "__pydantic_extra__", None) or {}
            if field in vextras and vextras[field] is not None:
                return vextras[field]
        val = getattr(content, field, None)
        if val is not None:
            return val
        extras = getattr(content, "__pydantic_extra__", None) or {}
        if field in extras and extras[field] is not None:
            return extras[field]
    return default


def _extra_list_merge(source, field: str) -> list:
    """对 list 类字段, 同时从 path X (variant) 和 path Y (content) 各读一份, 按 index 字段 union 合并.

    Why: LLM 生成 plan 时常同时写 spec.content.levels 和 spec.variant.levels,
    但因 schemas 字段名与 renderer 历史读字段名不同, 两侧数据可能字段互补 (variant 给 name, content 给 title+desc).
    单纯 _extra 优先 path X 会丢 path Y 的字段. 此 helper 按 index 把两侧字段 union, 避免数据丢失.

    Args:
        source: SlideSpec 实例
        field: list 字段名, 如 "levels" / "phases" / "steps" / "layers" / "units" / "descriptions"

    Returns:
        list of dict, 每个 dict 包含两侧 union 后的字段 (path X 优先, 但缺失字段从 path Y 补)
    """
    variant = getattr(source, "variant", None)
    content = getattr(source, "content", source)

    def _extract_list(obj, fname):
        if obj is None:
            return []
        v = getattr(obj, fname, None)
        if isinstance(v, list):
            return v
        extras = getattr(obj, "__pydantic_extra__", None) or {}
        v2 = extras.get(fname)
        return v2 if isinstance(v2, list) else []

    list_v = _extract_list(variant, field)
    list_c = _extract_list(content, field)

    def _to_dict(obj) -> dict:
        if obj is None:
            return {}
        if isinstance(obj, dict):
            return {k: v for k, v in obj.items() if v not in (None, "")}
        if isinstance(obj, str):
            return {"_value": obj}
        # Pydantic BaseModel
        d = {}
        # 已声明字段
        for k in obj.__class__.model_fields if hasattr(obj.__class__, "model_fields") else ():
            v = getattr(obj, k, None)
            if v not in (None, "", []):
                d[k] = v
        # extra 字段
        extras = getattr(obj, "__pydantic_extra__", None) or {}
        for k, v in extras.items():
            if v not in (None, "") and k not in d:
                d[k] = v
        return d

    n = max(len(list_v), len(list_c))
    merged = []
    for i in range(n):
        item_c = _to_dict(list_c[i]) if i < len(list_c) else {}
        item_v = _to_dict(list_v[i]) if i < len(list_v) else {}
        # path Y 在底, path X 覆盖 (相同 key)
        m = dict(item_c)
        for k, v in item_v.items():
            if v not in (None, "") or k not in m:
                m[k] = v
        merged.append(m)
    return merged


def _estimate_text_height_in(text: str, width_in: float, font_size_pt: float,
                             line_spacing: float = 1.25, padding_in: float = 0.15) -> float:
    """估算 wrapped text 在 textbox 中实际占用高度 (inches). 委托 lib/text_metrics 单一来源.

    2026-05-19 audit fix (ARCH-007): 公式 (cn 1.0 / en 0.55 加权) 已沉淀到 lib/text_metrics.py,
    这里仅做 width 兜底 (>=0.5) 后委托, 保留旧 API.
    """
    if not text:
        return 0.0
    from lib.text_metrics import estimate_text_height_inches
    return estimate_text_height_inches(
        text=str(text),
        width_inches=max(0.5, float(width_in)),
        font_size_pt=font_size_pt,
        line_spacing=line_spacing,
        padding_inches=padding_in / 2,  # text_metrics padding 是单边, 这里语义是单边总
    )


def _dget(d, *keys, default=None):
    """对 dict 或 Pydantic BaseModel 实例做多 key alias fallback.

    第一个非空就返回, 都缺时返 default.
    支持 path Y (yaml 顶层 content extra='allow', 解析为 list of dict)
    和 path X (yaml 用 variant, 解析为 list of Pydantic 子模型) 双兼容.
    """
    if d is None:
        return default
    is_dict = isinstance(d, dict)
    for k in keys:
        if is_dict:
            v = d.get(k)
        else:
            v = getattr(d, k, None)
            # Pydantic extra fields 在 __pydantic_extra__ 里
            if v is None:
                extras = getattr(d, "__pydantic_extra__", None) or {}
                v = extras.get(k)
        if v is not None and v != "":
            return v
    return default


def _add_rect(slide, left_in: float, top_in: float, width_in: float, height_in: float,
              fill_hex: str | None, line_hex: str | None = None, line_width_pt: float = 0.0):
    """添加纯色矩形（无圆角）。fill_hex=None 表示无填充；line_hex=None 表示无描边。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    if fill_hex is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(line_hex)
        shape.line.width = Pt(line_width_pt)
    return shape


def _canvas_dims_in(theme: dict) -> tuple[float, float]:
    """从 layouts.canvas 读画布 px 尺寸并换算 inch。fallback 到 16:9 标准。"""
    canvas = _layouts(theme).get("canvas", {})
    w_px = canvas.get("width_px", 1920)
    h_px = canvas.get("height_px", 1080)
    scale = 13.333 / max(w_px, 1)
    return w_px * scale, h_px * scale


def _canvas_padding_in(theme: dict) -> tuple[float, float, float, float]:
    """返回 (top, right, bottom, left) inches。"""
    canvas = _layouts(theme).get("canvas", {})
    w_px = canvas.get("width_px", 1920)
    scale = 13.333 / max(w_px, 1)
    top = canvas.get("padding_top_px", 88) * scale
    right = canvas.get("padding_right_px", 96) * scale
    bottom = canvas.get("padding_bottom_px", 76) * scale
    left = canvas.get("padding_left_px", 96) * scale
    return top, right, bottom, left


def _slide_dims_in() -> tuple[float, float]:
    """标准 16:9 slide dims."""
    return 13.333, 7.5


# 趋势色 token 决策（Plan §O5）
TREND_COLOR_TOKEN: dict[str, str] = {"up": "ok", "down": "warn", "flat": "ink_soft"}


def _font_family(spec: SlideSpec, theme: dict) -> str:
    """选择字体并接入 resolve_font_for_pptx, 让 Mac/Linux 上有 CJK fallback.

    历史教训 (2026-05-19 audit, INT-002): 此函数原直接返回 tokens.fonts.sans[0] = "Microsoft YaHei",
    huawei renderer 内 7 处 `xx.font.name = font` raw 写入 PPTX. Mac 上 PowerPoint 找不到 Microsoft YaHei
    走自身字体 fallback (Hiragino/PingFang), 与设计意图严重偏离, 严重时中文豆腐方块.
    现在统一接入 resolve_font_for_pptx, FONT_FALLBACK_CHAIN + CJK_FALLBACK_BY_OS 保证每平台都有 CJK 字体.
    """
    sans = _tokens(theme).get("fonts", {}).get("sans", [])
    preferred = sans[0] if sans else (spec.design.font_family or "Microsoft YaHei")
    return resolve_font_for_pptx(preferred)


# ===========================================================================
# P0 八版式
# ===========================================================================

@register_renderer("cover-left-bar")
def render_cover_left_bar(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """封面：左红条 16px + 大标题 108px + 副标题 34px + eyebrow 18px + 底部 meta 栏。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))

    cfg = _layout(theme, "cover_left_bar")
    bar_px = cfg.get("bar_width_px", 16)
    title_px = cfg.get("title_size_px", 108)
    subtitle_px = cfg.get("subtitle_size_px", 34)
    eyebrow_px = cfg.get("eyebrow_size_px", 18)
    meta_px = cfg.get("meta_size_px", 16)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)

    # 左红条
    bar_w_in = _px_in(bar_px)
    _add_rect(slide, 0, 0, bar_w_in, sh, _color(theme, "primary", "#C7000B"))

    font = _font_family(spec, theme)
    text_left = left_pad

    # eyebrow（小标签，可选）
    eyebrow = _extra(spec, "eyebrow", default=spec.chapter or "")
    cur_y = top_pad
    if eyebrow:
        eb_h = _px_pt(eyebrow_px) / 72 + 0.1
        add_textbox(
            slide, text_left, cur_y, sw - text_left - right_pad, eb_h,
            str(eyebrow), font, _px_pt(eyebrow_px),
            _color(theme, "primary", "#C7000B"), bold=True,
        )
        cur_y += eb_h + 0.15

    # 主标题
    title_pt = _px_pt(title_px)
    title_h = title_pt / 72 * 1.3 + 0.1
    add_textbox(
        slide, text_left, cur_y, sw - text_left - right_pad, title_h,
        spec.content.title, font, title_pt,
        _color(theme, "ink", "#1F1F1F"), bold=True,
    )
    cur_y += title_h + 0.1

    # 副标题
    if spec.content.subtitle:
        sub_pt = _px_pt(subtitle_px)
        sub_h = sub_pt / 72 * 1.4 + 0.2
        add_textbox(
            slide, text_left, cur_y, sw - text_left - right_pad, sub_h,
            spec.content.subtitle, font, sub_pt,
            _color(theme, "ink_soft", "#4B4B4B"),
        )

    # 底部 meta 栏（4 列横排）
    meta = _extra(spec, "meta", default=None) or _extra(spec, "meta_items", default=None) or []
    if meta:
        meta_y = sh - bottom_pad - _px_pt(meta_px) / 72 * 1.3 - 0.1
        col_w = (sw - text_left - right_pad) / max(len(meta), 1)
        meta_pt = _px_pt(meta_px)
        meta_h = meta_pt / 72 * 1.3 + 0.1
        for i, item in enumerate(meta[:4]):
            # 用 _dget 双兼容: dict (path Y) 或 Pydantic MetaItem (path X)
            label_v = _dget(item, "label", default="")
            value_v = _dget(item, "value", default="")
            if label_v or value_v:
                text = f"{label_v} {value_v}".strip()
            elif isinstance(item, str):
                text = item
            else:
                text = ""
            add_textbox(
                slide, text_left + i * col_w, meta_y, col_w, meta_h,
                text, font, meta_pt,
                _color(theme, "ink_mute", "#9A9A9A"),
            )


@register_renderer("toc")
def render_toc(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """目录：左侧大字 300px + 右侧章节列表（编号 × 标题 × 页码）。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))

    cfg = _layout(theme, "toc")
    giant_px = cfg.get("giant_title_size_px", 300)
    chapter_px = cfg.get("chapter_size_px", 22)
    page_px = cfg.get("page_size_px", 20)
    number_col_px = cfg.get("number_column_px", 160)
    item_gap_px = cfg.get("item_gap_px", 24)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    # 左侧大字（两栏布局中的左栏约 40% 宽）
    left_col_w = sw * 0.40 - left_pad
    giant_pt_default = _px_pt(giant_px)
    giant_text = _extra(spec, "giant_label", default=None) or spec.content.title or "目录"

    # L-2 fix: giant_label 字号需同时满足"宽度"和"高度"约束, 避免中文短文本被强制 wrap 后超出画布:
    # - 宽度: 中文字宽 ≈ font_pt/72 inch; 英文 ≈ 0.55 em. 总等价宽度不能超 left_col_w
    # - 高度: textbox 高度 ≈ font_pt/72 * 1.3 inches, 不能超画布可用区
    canvas_h = sh - top_pad - bottom_pad
    cn_chars = sum(1 for ch in giant_text if "一" <= ch <= "鿿")
    other_chars = len(giant_text) - cn_chars
    weighted_chars = max(1, cn_chars + other_chars * 0.55)
    max_pt_by_width = (left_col_w - 0.3) * 72 / weighted_chars
    max_pt_by_height = canvas_h * 72 / 1.4  # 单行高度 ≈ font_pt/72 * 1.3, 加 0.1 安全 margin
    giant_pt = min(giant_pt_default, max_pt_by_width, max_pt_by_height)

    # 字号确定后, textbox 高度按字号算
    box_h = giant_pt / 72 * 1.3 + 0.3
    gtb = slide.shapes.add_textbox(
        Inches(left_pad), Inches(top_pad),
        Inches(left_col_w), Inches(box_h),
    )
    gtf = gtb.text_frame
    gtf.word_wrap = False
    gp = gtf.paragraphs[0]
    gp.text = giant_text
    gp.font.name = font
    gp.font.size = Pt(giant_pt)
    gp.font.bold = True
    gp.font.color.rgb = hex_to_rgb(_color(theme, "primary", "#C7000B"))

    # 右侧章节列表
    items = _extra(spec, "items", default=None) or _extra(spec, "chapters", default=None) or []
    if not items:
        # 退化：用 key_points
        for p in get_points(spec):
            items.append({"title": p.heading or p.body, "page": ""})

    right_x = sw * 0.40
    right_w = sw - right_x - right_pad
    cur_y = top_pad
    row_h = _px_pt(chapter_px) / 72 * 1.8 + _px_in(item_gap_px)
    number_col_in = _px_in(number_col_px)

    chap_pt = _px_pt(chapter_px)
    page_pt = _px_pt(page_px)
    for i, item in enumerate(items):
        if cur_y + row_h > sh - bottom_pad:
            break
        if isinstance(item, str):
            title_text = item
            page_text = ""
        else:
            title_text = str(_dget(item, "title", default=""))
            page_text = str(_dget(item, "page", default=""))

        number_text = f"{i + 1:02d}"
        # 编号（红色大字）
        add_textbox(
            slide, right_x, cur_y, number_col_in, chap_pt / 72 * 1.4 + 0.1,
            number_text, font, chap_pt,
            _color(theme, "primary", "#C7000B"), bold=True,
        )
        # 标题
        add_textbox(
            slide, right_x + number_col_in, cur_y,
            right_w - number_col_in - 0.8, chap_pt / 72 * 1.4 + 0.1,
            title_text, font, chap_pt,
            _color(theme, "ink", "#1F1F1F"),
        )
        # 页码右对齐
        if page_text:
            add_textbox(
                slide, right_x + right_w - 0.8, cur_y + 0.02,
                0.8, page_pt / 72 * 1.4 + 0.1,
                page_text, font, page_pt,
                _color(theme, "ink_mute", "#9A9A9A"),
                alignment=PP_ALIGN.RIGHT,
            )
        # 分割线
        line_y = cur_y + row_h - 0.1
        _add_rect(
            slide, right_x, line_y, right_w, 0.01,
            _color(theme, "rule_soft", "#ECECEC"),
        )
        cur_y += row_h


@register_renderer("section-divider-dark")
def render_section_divider_dark(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """章节分隔页（深底）：左侧超大数字 360px + 右侧章节标题 88px + 描述。"""
    sw, sh = _slide_dims_in()
    bg_token = _layout(theme, "section_divider_dark").get("background_token", "ink_dark")
    set_slide_background(slide, _color(theme, bg_token, "#2A2A2A"))

    cfg = _layout(theme, "section_divider_dark")
    big_px = cfg.get("big_number_size_px", 360)
    title_px = cfg.get("title_size_px", 88)
    desc_px = cfg.get("description_size_px", 30)
    eyebrow_px = cfg.get("eyebrow_size_px", 22)
    left_col_px = cfg.get("grid_left_column_px", 380)
    gap_px = cfg.get("grid_gap_px", 80)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)
    paper = _color(theme, "paper", "#FFFFFF")
    primary = _color(theme, "primary", "#C7000B")
    ink_mute = _color(theme, "ink_mute", "#9A9A9A")

    # 左侧大数字
    big_pt = _px_pt(big_px)
    left_col_in = _px_in(left_col_px)
    number_text = _extra(spec, "number", "big_number") or str(spec.id).zfill(2)
    nbx = slide.shapes.add_textbox(
        Inches(left_pad), Inches(top_pad),
        Inches(max(left_col_in, 3.4)),
        Inches(big_pt / 72 * 1.0 + 0.3),
    )
    ntf = nbx.text_frame
    ntf.word_wrap = False
    np_ = ntf.paragraphs[0]
    np_.text = str(number_text)
    np_.font.name = font
    np_.font.size = Pt(big_pt)
    np_.font.bold = False
    np_.font.color.rgb = hex_to_rgb(primary)

    # 右侧（eyebrow / 标题 / 描述）
    right_x = left_pad + max(left_col_in, 3.4) + _px_in(gap_px)
    right_w = sw - right_x - right_pad
    cur_y = top_pad + 1.2  # 视觉上与数字中部对齐

    eyebrow = _extra(spec, "eyebrow", default=None) or spec.chapter
    if eyebrow:
        eb_pt = _px_pt(eyebrow_px)
        eb_h = eb_pt / 72 * 1.4 + 0.1
        add_textbox(
            slide, right_x, cur_y, right_w, eb_h,
            str(eyebrow), font, eb_pt,
            primary, bold=True,
        )
        cur_y += eb_h + 0.1

    title_pt = _px_pt(title_px)
    title_text_str = spec.content.title or ""
    # L-1 fix: 用估算的 wrapped 高度作 textbox height, 让下游 description 能链式定位
    title_h = max(
        title_pt / 72 * 1.3 + 0.2,
        _estimate_text_height_in(title_text_str, right_w, title_pt, line_spacing=1.18, padding_in=0.15),
    )
    add_textbox(
        slide, right_x, cur_y, right_w, title_h,
        title_text_str, font, title_pt,
        paper, bold=True,
    )
    # gutter: title 字号 88pt 时约 0.3", 字号小时按比例
    gutter_in = max(0.2, title_pt / 72 * 0.25)
    cur_y += title_h + gutter_in

    desc = spec.content.description or spec.content.subtitle
    if desc:
        desc_pt = _px_pt(desc_px)
        desc_h = max(
            desc_pt / 72 * 1.5 + 0.3,
            _estimate_text_height_in(desc, right_w, desc_pt, line_spacing=1.4, padding_in=0.2),
        )
        # 防止溢出底边: 下边界限制
        max_desc_h = sh - bottom_pad - cur_y - 0.2
        desc_h = min(desc_h, max_desc_h)
        add_textbox(
            slide, right_x, cur_y, right_w, desc_h,
            desc, font, desc_pt,
            ink_mute,
        )


@register_renderer("kpi-stats")
def render_kpi_stats(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """KPI 统计：N 列（默认 4）横排，每列 value(108px) + label(18px) + desc + trend(14px)。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "kpi_stats")
    columns_cfg = cfg.get("columns", 4)
    value_px = cfg.get("value_size_px", 108)  # Plan §O6 修正
    unit_px = cfg.get("unit_size_px", 32)
    label_px = cfg.get("label_size_px", 18)
    desc_px = cfg.get("desc_size_px", 16)
    trend_px = cfg.get("trend_size_px", 14)
    sep_px = cfg.get("column_separator_width_px", 2)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    kpis = _extra(spec, "kpis", default=None) or []
    if not kpis:
        # 退化：从 key_points 构造
        for p in get_points(spec):
            kpis.append({
                "value": p.metric_value or "",
                "label": p.metric_label or p.heading or "",
                "desc": p.body or "",
            })
    n = min(len(kpis), columns_cfg) if kpis else columns_cfg

    area_left = left_pad
    area_top = top_pad + 1.4  # title zone 占位
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5
    col_w = area_w / max(n, 1)

    value_pt = _px_pt(value_px)
    unit_pt = _px_pt(unit_px)
    label_pt = _px_pt(label_px)
    desc_pt = _px_pt(desc_px)
    trend_pt = _px_pt(trend_px)
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    primary = _color(theme, "primary", "#C7000B")

    for i in range(n):
        item: Any = kpis[i] if i < len(kpis) else {}
        if isinstance(item, str):
            item = {"value": item, "label": ""}
        x = area_left + i * col_w
        cy = area_top

        value_text = str(_dget(item, "value", default=""))
        unit_text = str(_dget(item, "unit", default=""))
        # L-3 fix: value 强制单行 (KPI 数字本质单行展示), 字号超 col_w 时自适应缩放,
        # 避免 word_wrap=True 让长 value (如 "15,000" + unit) 折行后压住下方 label.
        # 算 value+unit 总宽度等价中文字数, 超 col_w 时按比例缩字号
        cn_v = sum(1 for ch in value_text if "一" <= ch <= "鿿")
        wt_v = cn_v + (len(value_text) - cn_v) * 0.55
        cn_u = sum(1 for ch in unit_text if "一" <= ch <= "鿿")
        wt_u = cn_u + (len(unit_text) - cn_u) * 0.55
        # 留 25% buffer 防 LibreOffice/PowerPoint 实际渲染字宽比理论大 (中文等宽字体在 PowerPoint 实测约 1.1-1.2 em)
        value_max_w = (col_w - 0.2) * 0.85
        # value 主体宽 = wt_v * value_pt/72; unit 宽 = wt_u * unit_pt/72
        # 不超 value_max_w 时维持 default; 超时按比例缩 value_pt (留 buffer)
        full_w = wt_v * value_pt / 72 + (wt_u + 1) * unit_pt / 72  # +1 char 间距
        if full_w > value_max_w and full_w > 0:
            scale = value_max_w / full_w
            value_pt_eff = max(20.0, value_pt * scale)
            unit_pt_eff = max(12.0, unit_pt * scale)
        else:
            value_pt_eff = value_pt
            unit_pt_eff = unit_pt
        value_actual_h = value_pt_eff / 72 * 1.25 + 0.15
        if unit_text:
            parts = [
                {"text": value_text, "size_pt": value_pt_eff, "color": primary, "bold": True},
                {"text": f" {unit_text}", "size_pt": unit_pt_eff, "color": primary, "bold": False},
            ]
            txbox = add_textbox_rich(slide, x, cy, col_w, value_actual_h, parts, font)
            txbox.text_frame.word_wrap = False  # 关键: 强制单行
        else:
            txbox = add_textbox(
                slide, x, cy, col_w, value_actual_h,
                value_text, font, value_pt_eff,
                primary, bold=True,
            )
            txbox.text_frame.word_wrap = False
        cy += value_actual_h + max(0.08, value_pt_eff / 72 * 0.12)  # gutter 按字号缩放

        label_text = str(_dget(item, "label", default=""))
        if label_text:
            label_actual_h = max(
                label_pt / 72 * 1.4 + 0.1,
                _estimate_text_height_in(label_text, col_w, label_pt, line_spacing=1.4, padding_in=0.08),
            )
            add_textbox(
                slide, x, cy, col_w, label_actual_h,
                label_text, font, label_pt,
                ink, bold=True,
            )
            cy += label_actual_h + 0.05

        desc_text = str(_dget(item, "desc", "description", default=""))
        if desc_text:
            desc_actual_h = max(
                desc_pt / 72 * 1.5 + 0.3,
                _estimate_text_height_in(desc_text, col_w, desc_pt, line_spacing=1.4, padding_in=0.15),
            )
            add_textbox(
                slide, x, cy, col_w, desc_actual_h,
                desc_text, font, desc_pt,
                ink_soft,
            )
            cy += desc_actual_h + 0.1

        trend_text = str(_dget(item, "trend_text", default=""))
        trend_dir = _dget(item, "trend", "trend_direction", default="")
        if trend_text or trend_dir:
            token = TREND_COLOR_TOKEN.get(str(trend_dir).lower(), "ink_soft")
            tcolor = _color(theme, token, ink_soft)
            arrow = {"up": "↑", "down": "↓", "flat": "→"}.get(str(trend_dir).lower(), "")
            add_textbox(
                slide, x, cy, col_w, trend_pt / 72 * 1.5 + 0.15,
                f"{arrow} {trend_text}".strip(), font, trend_pt,
                tcolor,
            )

        # 列分隔线（非末列）
        if i < n - 1:
            sep_x = x + col_w - _px_in(sep_px) / 2
            _add_rect(
                slide, sep_x, area_top, max(_px_in(sep_px), 0.01), area_h,
                _color(theme, "rule", "#D9D9D9"),
            )


@register_renderer("matrix-2x2")
def render_matrix_2x2(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """2x2 矩阵：左 Y 轴 + 底 X 轴 + 4 象限卡片，可按 highlight 标记高亮象限。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "matrix_2x2")
    y_col_px = cfg.get("axis_y_column_px", 140)
    x_row_px = cfg.get("axis_x_row_px", 140)
    gap_px = cfg.get("grid_gap_px", 20)
    border_px = cfg.get("border_width_px", 2)
    border_token = cfg.get("border_token", "ink")
    heading_px = cfg.get("quadrant_heading_size_px", 22)
    y_label_px = cfg.get("y_axis_label_size_px", 20)
    small_px = cfg.get("small_label_size_px", 14)
    desc_px = cfg.get("desc_size_px", 16)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5

    y_col_in = _px_in(y_col_px)
    x_row_in = _px_in(x_row_px)
    gap_in = _px_in(gap_px)

    grid_left = area_left + y_col_in
    grid_top = area_top
    grid_w = area_w - y_col_in
    grid_h = area_h - x_row_in
    cell_w = (grid_w - gap_in) / 2
    cell_h = (grid_h - gap_in) / 2

    border_color = _color(theme, border_token, "#1F1F1F")
    primary = _color(theme, "primary", "#C7000B")
    primary_soft = _color(theme, "primary_soft", "#FDECEE")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    ink_mute = _color(theme, "ink_mute", "#9A9A9A")
    paper_2 = _color(theme, "paper_2", "#F4F4F4")

    # Y 轴（左侧，垂直文字）
    y_axis = _extra(spec, "y_axis", default=None) or _extra(spec, "axis_y", default="")
    if y_axis:
        # L-4 fix: 旋转 textbox 视觉宽度 = textbox.height. 长文本时单行宽度 = chars * font_pt/72,
        # 超出可用空间时会侵入 grid (压象限). 字号自适应或截断 + footnote 双策略.
        y_text_str = str(y_axis)
        y_text_pt_default = _px_pt(y_label_px)
        cn_chars_y = sum(1 for ch in y_text_str if "一" <= ch <= "鿿")
        weighted_y = max(1, cn_chars_y + (len(y_text_str) - cn_chars_y) * 0.55)
        # 旋转中心 X = area_left + 0.45, grid 起点 = area_left + y_col_in
        # 视觉宽度 (即 textbox.height) 不能压 grid: 0.45 + height/2 ≤ y_col_in - 0.1
        max_text_in = max(0.8, (y_col_in - 0.55) * 2)
        text_width_in = weighted_y * y_text_pt_default / 72
        if text_width_in <= max_text_in:
            y_text_pt = y_text_pt_default
            y_text_main = y_text_str
            y_text_footnote = None
        else:
            # 先尝试缩字号到能容纳 (但 ≥ 9pt 保证可读)
            scaled_pt = max_text_in / weighted_y * 72
            if scaled_pt >= 9:
                y_text_pt = scaled_pt
                y_text_main = y_text_str
                y_text_footnote = None
            else:
                # 字号过小, 改 truncate + 把全文放底部 footnote
                # 算阈值 = max_chars 在 9pt 下能装的中文字符数
                max_chars = max(2, int(max_text_in * 72 / 9))
                y_text_main = y_text_str[: max_chars - 1] + "…" if len(y_text_str) > max_chars else y_text_str
                y_text_pt = 9
                y_text_footnote = y_text_str

        rot_w = 0.5
        rot_h = max(0.8, weighted_y * y_text_pt / 72 + 0.3)  # textbox 视觉宽度恰好容纳文字
        rbx = slide.shapes.add_textbox(
            Inches(area_left + 0.2), Inches(grid_top + grid_h / 2 - rot_h / 2),
            Inches(rot_w), Inches(rot_h),
        )
        rbx.rotation = 270.0
        rtf = rbx.text_frame
        rtf.word_wrap = False
        rp = rtf.paragraphs[0]
        rp.alignment = PP_ALIGN.CENTER
        rp.text = y_text_main
        rp.font.name = font
        rp.font.size = Pt(y_text_pt)
        rp.font.bold = True
        rp.font.color.rgb = hex_to_rgb(ink_soft)
        # footnote: 完整 y_axis 文字放在 axis 列底部, 横向小字
        if y_text_footnote:
            footnote_pt = max(8, int(y_text_pt * 0.85))
            add_textbox(
                slide, area_left, grid_top + grid_h - 0.4, y_col_in, 0.35,
                y_text_footnote, font, footnote_pt,
                ink_mute, alignment=PP_ALIGN.CENTER,
            )

    # X 轴（底部）
    x_axis = _extra(spec, "x_axis", default=None) or _extra(spec, "axis_x", default="")
    if x_axis:
        add_textbox(
            slide, grid_left, grid_top + grid_h + 0.15, grid_w, x_row_in - 0.2,
            str(x_axis), font, _px_pt(y_label_px),
            ink_soft, alignment=PP_ALIGN.CENTER, bold=True,
        )

    # 4 象限（quadrants）
    quadrants = _extra(spec, "quadrants", default=None) or []
    if not quadrants:
        pts = get_points(spec)
        for p in pts[:4]:
            quadrants.append({
                "heading": p.heading or "",
                "desc": p.body or "",
                "highlight": False,
            })
    # 顺序约定：左上 / 右上 / 左下 / 右下
    positions = [
        (grid_left, grid_top),
        (grid_left + cell_w + gap_in, grid_top),
        (grid_left, grid_top + cell_h + gap_in),
        (grid_left + cell_w + gap_in, grid_top + cell_h + gap_in),
    ]
    heading_pt = _px_pt(heading_px)
    desc_pt = _px_pt(desc_px)
    small_pt = _px_pt(small_px)

    for i, (x, y) in enumerate(positions):
        item = quadrants[i] if i < len(quadrants) else {}
        if isinstance(item, str):
            item = {"heading": "", "desc": item}
        highlight = bool(_dget(item, "highlight", default=False))
        fill = primary_soft if highlight else paper_2
        _add_rect(slide, x, y, cell_w, cell_h, fill,
                  line_hex=border_color, line_width_pt=max(border_px * 0.5, 0.5))

        # 角标（tag / small_label）
        tag = _dget(item, "tag", default="") or _dget(item, "corner_tag", default="")
        if tag:
            add_textbox(
                slide, x + 0.15, y + 0.1, cell_w - 0.3, small_pt / 72 * 1.5 + 0.1,
                str(tag), font, small_pt,
                primary if highlight else ink_mute, bold=True,
            )
        heading = _dget(item, "heading", default="") or _dget(item, "title", default="")
        add_textbox(
            slide, x + 0.2, y + 0.4, cell_w - 0.4, heading_pt / 72 * 1.4 + 0.2,
            str(heading), font, heading_pt,
            ink, bold=True,
        )
        desc = _dget(item, "desc", default="") or _dget(item, "description", default="") or _dget(item, "body", default="")
        if desc:
            add_textbox(
                slide, x + 0.2, y + 0.4 + heading_pt / 72 * 1.4 + 0.1,
                cell_w - 0.4, cell_h - 0.4 - heading_pt / 72 * 1.4 - 0.3,
                str(desc), font, desc_pt,
                ink_soft,
            )


@register_renderer("architecture-layered")
def render_architecture_layered(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """分层架构图：每层一行，左侧 header 列 + 右侧 N 个 cell；支持 highlight cell 红底白字。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "architecture_layered")
    header_col_px = cfg.get("header_column_px", 260)
    cell_columns = cfg.get("cell_columns", 6)
    layer_gap_px = cfg.get("layer_gap_px", 14)
    header_bg_token = cfg.get("header_background_token", "paper_3")
    header_name_px = cfg.get("header_name_size_px", 24)
    header_label_px = cfg.get("header_label_size_px", 14)
    cell_title_px = cfg.get("cell_title_size_px", 16)
    cell_desc_px = cfg.get("cell_desc_size_px", 14)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5

    layers = _extra(spec, "layers", default=None) or []
    if not layers:
        for p in get_points(spec):
            layers.append({
                "name": p.heading or "",
                "cells": [{"title": p.body or ""}],
            })
    n_layers = max(len(layers), 1)

    header_w = _px_in(header_col_px)
    gap_in = _px_in(layer_gap_px)
    layer_h = (area_h - gap_in * (n_layers - 1)) / n_layers
    cell_area_w = area_w - header_w - 0.15

    header_name_pt = _px_pt(header_name_px)
    header_label_pt = _px_pt(header_label_px)
    cell_title_pt = _px_pt(cell_title_px)
    cell_desc_pt = _px_pt(cell_desc_px)

    header_bg = _color(theme, header_bg_token, "#EAEAEA")
    paper_2 = _color(theme, "paper_2", "#F4F4F4")
    primary = _color(theme, "primary", "#C7000B")
    paper = _color(theme, "paper", "#FFFFFF")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    ink_mute = _color(theme, "ink_mute", "#9A9A9A")

    for li, layer in enumerate(layers):
        if isinstance(layer, str):
            layer = {"name": layer, "cells": []}
        y = area_top + li * (layer_h + gap_in)

        # Header 格
        _add_rect(slide, area_left, y, header_w, layer_h, header_bg)
        # schemas: layer.header.{code, name}; legacy: layer.{name, label}
        header_obj = _dget(layer, "header", default=None)
        if header_obj is not None:
            name_text = _dget(header_obj, "name", default="")
            label_text = _dget(header_obj, "code", default="") or f"L{li + 1}"
        else:
            name_text = _dget(layer, "name", default="")
            label_text = _dget(layer, "label", default="") or f"L{li + 1}"
        add_textbox(
            slide, area_left + 0.15, y + 0.15,
            header_w - 0.3, header_name_pt / 72 * 1.4 + 0.1,
            name_text, font, header_name_pt,
            ink, bold=True,
        )
        add_textbox(
            slide, area_left + 0.15,
            y + 0.15 + header_name_pt / 72 * 1.4 + 0.05,
            header_w - 0.3, header_label_pt / 72 * 1.4 + 0.05,
            label_text, font, header_label_pt,
            ink_mute,
        )

        # 右侧 cells (schemas 与 legacy 字段名一致)
        cells = _dget(layer, "cells", default=None) or []
        n_cells = len(cells) if cells else cell_columns
        cw = cell_area_w / max(n_cells, 1)
        for ci in range(n_cells):
            cell = cells[ci] if ci < len(cells) else {}
            if isinstance(cell, str):
                cell = {"title": cell}
            cx = area_left + header_w + 0.15 + ci * cw
            highlight = bool(_dget(cell, "highlight", default=False))
            fill = primary if highlight else paper_2
            _add_rect(slide, cx, y, cw - 0.05, layer_h, fill)

            ctitle = _dget(cell, "title", default="")
            cdesc = _dget(cell, "desc", "description", default="")
            title_color = paper if highlight else ink
            desc_color = paper if highlight else ink_soft
            add_textbox(
                slide, cx + 0.1, y + 0.15,
                cw - 0.2, cell_title_pt / 72 * 1.5 + 0.1,
                ctitle, font, cell_title_pt,
                title_color, bold=True,
            )
            if cdesc:
                add_textbox(
                    slide, cx + 0.1,
                    y + 0.15 + cell_title_pt / 72 * 1.5 + 0.05,
                    cw - 0.2, layer_h - 0.3 - cell_title_pt / 72 * 1.5,
                    cdesc, font, cell_desc_pt,
                    desc_color,
                )


@register_renderer("timeline-huawei")
def render_timeline_huawei(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """华为时间轴：顶部年份大字 48px + 阶段连接线 + 阶段标题 22px + 描述 16px。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "timeline_huawei")
    phase_col_px = cfg.get("phase_column_px", 180)
    year_px = cfg.get("year_size_px", 48)
    title_px = cfg.get("title_size_px", 22)
    desc_px = cfg.get("desc_size_px", 16)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5

    phases = _extra(spec, "phases", default=None) or _extra(spec, "timeline", default=None) or []
    if not phases:
        for p in get_points(spec):
            phases.append({
                "year": p.heading or "",
                "title": p.metric_label or "",
                "desc": p.body or "",
            })
    n = max(len(phases), 1)
    col_w = area_w / n

    year_pt = _px_pt(year_px)
    title_pt = _px_pt(title_px)
    desc_pt = _px_pt(desc_px)

    primary = _color(theme, "primary", "#C7000B")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    rule_c = _color(theme, "rule", "#D9D9D9")

    # 横贯连接线
    line_y = area_top + year_pt / 72 * 1.3 + 0.4
    _add_rect(
        slide, area_left, line_y, area_w, 0.03,
        rule_c,
    )

    for i, phase in enumerate(phases):
        if isinstance(phase, str):
            phase = {"year": "", "title": phase, "desc": ""}
        x = area_left + i * col_w
        cy = area_top

        year_text = str(_dget(phase, "year", default=""))
        add_textbox(
            slide, x, cy, col_w, year_pt / 72 * 1.3 + 0.2,
            year_text, font, year_pt,
            primary, bold=True,
        )

        # 节点圆点
        node_d = 0.25
        node_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + col_w / 2 - node_d / 2),
            Inches(line_y - node_d / 2 + 0.015),
            Inches(node_d), Inches(node_d),
        )
        node_shape.fill.solid()
        node_shape.fill.fore_color.rgb = hex_to_rgb(primary)
        node_shape.line.fill.background()

        body_top = line_y + 0.3
        # schemas: TimelinePhase.label/items[]; legacy: title/desc/description
        title_text = _dget(phase, "title", "label", default="")
        add_textbox(
            slide, x + 0.1, body_top, col_w - 0.2, title_pt / 72 * 1.4 + 0.15,
            title_text, font, title_pt,
            ink, bold=True,
        )
        # desc 优先 string; 若 schemas 给的是 items list, 折叠成 "• item1\n• item2"
        desc_text = _dget(phase, "desc", "description", default="")
        if not desc_text:
            items_list = _dget(phase, "items", default=None)
            if isinstance(items_list, (list, tuple)) and items_list:
                desc_text = "\n".join(f"• {it}" for it in items_list)
        if desc_text:
            add_textbox(
                slide, x + 0.1,
                body_top + title_pt / 72 * 1.4 + 0.1,
                col_w - 0.2,
                area_top + area_h - body_top - title_pt / 72 * 1.4 - 0.1,
                desc_text, font, desc_pt,
                ink_soft,
            )


@register_renderer("process-flow-huawei")
def render_process_flow_huawei(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """华为流程：N 个圆角 rect + 箭头连接器（步骤标题 24px / 描述 16px / 编号 14px）。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "process_flow_huawei")
    step_min_px = cfg.get("step_min_width_px", 240)
    arrow_gap_px = cfg.get("arrow_gap_px", 40)
    title_px = cfg.get("step_title_size_px", 24)
    desc_px = cfg.get("step_desc_size_px", 16)
    number_px = cfg.get("step_number_size_px", 14)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5

    steps = _extra(spec, "steps", default=None) or _extra(spec, "process", default=None) or []
    if not steps:
        for p in get_points(spec):
            steps.append({"title": p.heading or "", "desc": p.body or ""})
    n = max(len(steps), 1)

    arrow_gap_in = _px_in(arrow_gap_px)
    step_w = (area_w - arrow_gap_in * (n - 1)) / n
    step_h = min(area_h, 2.6)
    step_top = area_top + (area_h - step_h) / 2

    title_pt = _px_pt(title_px)
    desc_pt = _px_pt(desc_px)
    number_pt = _px_pt(number_px)

    primary = _color(theme, "primary", "#C7000B")
    primary_soft = _color(theme, "primary_soft", "#FDECEE")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    ink_mute = _color(theme, "ink_mute", "#9A9A9A")

    for i, step in enumerate(steps):
        if isinstance(step, str):
            step = {"title": step, "desc": ""}
        x = area_left + i * (step_w + arrow_gap_in)

        # 圆角卡片（浅红底）
        add_rounded_rect(slide, x, step_top, step_w, step_h, primary_soft, corner_radius=0.1)

        number_text = f"{i + 1:02d}"
        add_textbox(
            slide, x + 0.2, step_top + 0.15, step_w - 0.4,
            number_pt / 72 * 1.5 + 0.1,
            number_text, font, number_pt,
            primary, bold=True,
        )
        # schemas: ProcessStep.label/desc; legacy: title
        title_text = _dget(step, "title", "label", default="")
        add_textbox(
            slide, x + 0.2, step_top + 0.5, step_w - 0.4,
            title_pt / 72 * 1.4 + 0.15,
            title_text, font, title_pt,
            ink, bold=True,
        )
        desc_text = _dget(step, "desc", "description", default="")
        if desc_text:
            add_textbox(
                slide, x + 0.2,
                step_top + 0.5 + title_pt / 72 * 1.4 + 0.15,
                step_w - 0.4,
                step_h - 0.7 - title_pt / 72 * 1.4,
                desc_text, font, desc_pt,
                ink_soft,
            )

        # 箭头（非末步）
        if i < n - 1:
            ax = x + step_w + arrow_gap_in * 0.1
            ay = step_top + step_h / 2 - 0.15
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(ax), Inches(ay),
                Inches(arrow_gap_in * 0.8), Inches(0.3),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = hex_to_rgb(primary)
            arr.line.fill.background()


# ===========================================================================
# P1 五版式
# ===========================================================================

@register_renderer("swot")
def render_swot(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """SWOT 2x2：四象限 (S/W/O/T) + 中心红十字分隔线。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "swot")
    cross_w_px = cfg.get("cross_line_width_px", 2)
    cross_token = cfg.get("cross_line_token", "primary")

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5
    cell_w = area_w / 2
    cell_h = area_h / 2

    primary = _color(theme, "primary", "#C7000B")
    cross_color = _color(theme, cross_token, primary)
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    paper_2 = _color(theme, "paper_2", "#F4F4F4")

    quadrant_map = [
        ("strengths", "S", "优势"),
        ("weaknesses", "W", "劣势"),
        ("opportunities", "O", "机会"),
        ("threats", "T", "威胁"),
    ]
    body_pt = _px_pt(_tokens(theme).get("type_scale_px", {}).get("body", 24))
    heading_pt = _px_pt(_tokens(theme).get("type_scale_px", {}).get("h4", 28))

    for i, (key, letter, label) in enumerate(quadrant_map):
        row, col = divmod(i, 2)
        x = area_left + col * cell_w
        y = area_top + row * cell_h
        _add_rect(slide, x, y, cell_w, cell_h, paper_2)

        add_textbox(
            slide, x + 0.2, y + 0.15, 1.0, heading_pt / 72 * 1.3 + 0.1,
            letter, font, heading_pt * 1.5,
            primary, bold=True,
        )
        add_textbox(
            slide, x + 1.2, y + 0.3, cell_w - 1.4, heading_pt / 72 * 1.3 + 0.1,
            label, font, heading_pt,
            ink, bold=True,
        )
        items = _extra(spec, key, None) or []
        if isinstance(items, str):
            items = [items]
        text = "\n".join(f"• {str(it)}" for it in items[:5]) if items else ""
        if text:
            add_textbox(
                slide, x + 0.3, y + 1.1, cell_w - 0.5, cell_h - 1.3,
                text, font, body_pt,
                ink_soft,
            )

    # 中心红十字线
    line_pt = max(cross_w_px * 0.5, 1.0)
    # 横线
    hl = slide.shapes.add_connector(
        1,
        Inches(area_left), Inches(area_top + cell_h),
        Inches(area_left + area_w), Inches(area_top + cell_h),
    )
    hl.line.color.rgb = hex_to_rgb(cross_color)
    hl.line.width = Pt(line_pt)
    # 竖线
    vl = slide.shapes.add_connector(
        1,
        Inches(area_left + cell_w), Inches(area_top),
        Inches(area_left + cell_w), Inches(area_top + area_h),
    )
    vl.line.color.rgb = hex_to_rgb(cross_color)
    vl.line.width = Pt(line_pt)


@register_renderer("roadmap")
def render_roadmap(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """Roadmap 泳道图：左侧 lane 列 + 右侧多阶段横向条，按 emphasis 映射填充色。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "roadmap")
    lane_col_px = cfg.get("lane_column_px", 180)
    bar_h_px = cfg.get("phase_bar_height_px", 56)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5

    lane_w = _px_in(lane_col_px)
    bar_h = _px_in(bar_h_px)
    body_pt = _px_pt(_type_px(theme, "body", 24))
    small_pt = _px_pt(_type_px(theme, "small", 20))

    lanes = _extra(spec, "lanes", default=None) or []
    phases = _extra(spec, "phases", default=None) or ["Q1", "Q2", "Q3", "Q4"]
    n_phases = max(len(phases), 1)
    right_x = area_left + lane_w + 0.15
    right_w = area_w - lane_w - 0.15
    phase_col_w = right_w / n_phases

    primary = _color(theme, "primary", "#C7000B")
    accent = _color(theme, "accent", "#D97706")
    paper_2 = _color(theme, "paper_2", "#F4F4F4")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    rule = _color(theme, "rule_soft", "#ECECEC")

    # Phase 表头: title 主显 + summary 副显 (双层)
    title_h = 0.28
    summary_h = 0.32
    header_h = title_h + summary_h
    summary_pt = max(_px_pt(_type_px(theme, "small", 20)) - 4, 10)
    for pi, ph in enumerate(phases):
        if isinstance(ph, str):
            title_text = ph
            summary_text = ""
        else:
            title_text = str(_dget(ph, "title", default="") or "")
            summary_text = str(_dget(ph, "summary", default="") or "")
        add_textbox(
            slide, right_x + pi * phase_col_w, area_top,
            phase_col_w, title_h,
            title_text, font, small_pt,
            ink_soft, alignment=PP_ALIGN.CENTER, bold=True,
        )
        if summary_text:
            add_textbox(
                slide, right_x + pi * phase_col_w, area_top + title_h,
                phase_col_w, summary_h,
                summary_text, font, summary_pt,
                ink_soft, alignment=PP_ALIGN.CENTER, bold=False,
            )

    # 泳道
    if not lanes:
        for p in get_points(spec):
            lanes.append({"name": p.heading or "", "bars": []})
    n_lanes = max(len(lanes), 1)
    lane_h = (area_h - header_h - 0.15) / n_lanes

    emphasis_fill = {
        "primary": primary,
        "accent": accent,
        "bar": primary,
        "bar2": accent,
        "default": paper_2,
    }

    for li, lane in enumerate(lanes):
        if isinstance(lane, str):
            lane = {"name": lane, "bars": []}
        ly = area_top + header_h + 0.15 + li * lane_h

        # 泳道标签
        add_textbox(
            slide, area_left, ly + (lane_h - bar_h) / 2, lane_w, bar_h,
            str(_dget(lane, "name", default="")), font, body_pt,
            ink, bold=True,
        )
        # 泳道底色
        _add_rect(slide, right_x, ly, right_w, lane_h, rule)

        for bar in _dget(lane, "bars", default=[]) or []:
            start = int(_dget(bar, "start", default=0))
            span = int(_dget(bar, "span", default=1))
            emphasis = str(_dget(bar, "emphasis", default="default"))
            fill = emphasis_fill.get(emphasis, paper_2)
            bx = right_x + start * phase_col_w + 0.06
            bw = max(span * phase_col_w - 0.12, 0.3)
            by = ly + (lane_h - bar_h) / 2
            add_rounded_rect(slide, bx, by, bw, bar_h, fill, corner_radius=0.08)
            label = str(_dget(bar, "label", default=""))
            if label:
                text_color = "#FFFFFF" if emphasis in ("primary", "accent", "bar", "bar2") else ink
                add_textbox(
                    slide, bx + 0.15, by + 0.05, bw - 0.3, bar_h - 0.1,
                    label, font, small_pt,
                    text_color, bold=True,
                )


@register_renderer("pyramid")
def render_pyramid(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """金字塔：N 层 TRAPEZOID 自顶向下递增宽度。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "pyramid")
    levels_default = cfg.get("levels_default", 4)
    top_ratio = float(cfg.get("top_width_ratio", 0.35))

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5

    # L-7 fix Bug B 配套: 用 _extra_list_merge 同时取 path X (variant) 和 path Y (content) 的 levels,
    # 避免 LLM 双写时字段不一致导致数据丢失 (如 variant.levels 只有 name, content.levels 有 title+desc).
    levels = _extra_list_merge(spec, "levels")
    if not levels:
        for p in get_points(spec):
            levels.append({"title": p.heading or "", "desc": p.body or ""})
    n = len(levels) or levels_default

    # L-7 主体: 优先用独立 descriptions[] (schemas 字段, 比 levels.desc 更精细), 以 layer center 为 anchor 对齐.
    descriptions = _extra_list_merge(spec, "descriptions")

    # 金字塔整体：左侧 55% 宽，右侧留给描述
    pyr_w = area_w * 0.55
    pyr_cx = area_left + pyr_w / 2
    layer_h = area_h / n

    primary = _color(theme, "primary", "#C7000B")
    primary_dark = _color(theme, "primary_dark", "#8A0008")
    paper_2 = _color(theme, "paper_2", "#F4F4F4")
    paper = _color(theme, "paper", "#FFFFFF")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    body_pt = _px_pt(_type_px(theme, "body", 24))
    small_pt = _px_pt(_type_px(theme, "small", 20))
    h4_pt = _px_pt(_type_px(theme, "h4", 28))

    # 每层梯形：自顶向下宽度线性递增 top_ratio → 1.0
    for i in range(n):
        level = levels[i] if i < len(levels) else {}
        if isinstance(level, str):
            level = {"title": level, "desc": ""}
        frac_top = top_ratio + (1.0 - top_ratio) * (i / n)
        frac_bot = top_ratio + (1.0 - top_ratio) * ((i + 1) / n)
        top_w = pyr_w * frac_top
        bot_w = pyr_w * frac_bot
        # 用 TRAPEZOID 近似：取底宽为 bot_w，adj 控制上底比例
        y = area_top + i * layer_h
        shape = slide.shapes.add_shape(
            MSO_SHAPE.TRAPEZOID,
            Inches(pyr_cx - bot_w / 2), Inches(y),
            Inches(bot_w), Inches(layer_h - 0.03),
        )
        shape.fill.solid()
        # 顶层主色，越往下越浅
        if i == 0:
            shape.fill.fore_color.rgb = hex_to_rgb(primary_dark)
        elif i == 1:
            shape.fill.fore_color.rgb = hex_to_rgb(primary)
        else:
            shape.fill.fore_color.rgb = hex_to_rgb(paper_2)
        shape.line.fill.background()

        # TRAPEZOID 的 adj 设置（控制上底比例）
        try:
            sp = shape._element
            prstGeom = sp.find(qn("a:prstGeom"))
            if prstGeom is not None:
                avLst = prstGeom.find(qn("a:avLst"))
                if avLst is None:
                    avLst = etree.SubElement(prstGeom, qn("a:avLst"))
                # adj = (bot_w - top_w) / 2 / bot_w × 100000（留边百分比）
                adj_val = int(max(0.0, (bot_w - top_w) / 2 / max(bot_w, 0.01)) * 100000)
                gd = etree.SubElement(avLst, qn("a:gd"))
                gd.set("name", "adj")
                gd.set("fmla", f"val {adj_val}")
        except Exception:
            pass

        # 层内文字 (schemas: name; legacy: title)
        title_text = _dget(level, "title", "name", default="")
        text_color = paper if i < 2 else ink
        add_textbox(
            slide, pyr_cx - bot_w / 2, y + (layer_h - 0.03) / 2 - 0.2,
            bot_w, 0.4,
            title_text, font, body_pt,
            text_color, alignment=PP_ALIGN.CENTER, bold=True,
        )

        # 右侧描述: 优先用独立 descriptions[i], 以 layer center 为 anchor 链式定位
        # 兜底: legacy 用 level.desc/description/tag
        rx = area_left + pyr_w + 0.3
        rw = area_w - pyr_w - 0.3
        layer_center_y = y + (layer_h - 0.03) / 2
        if i < len(descriptions):
            d = descriptions[i]
            num = _dget(d, "number", default="")
            heading = _dget(d, "heading", default="")
            body = _dget(d, "body", default="")
            head_pt = small_pt
            head_h = head_pt / 72 * 1.4 + 0.1
            body_pt_local = small_pt - 1
            body_h = max(
                body_pt_local / 72 * 1.5 + 0.2,
                _estimate_text_height_in(body, rw, body_pt_local, line_spacing=1.4, padding_in=0.1),
            ) if body else 0.0
            block_h = head_h + body_h
            block_top = layer_center_y - block_h / 2
            head_text = f"{num}  {heading}".strip() if num else heading
            if head_text:
                add_textbox(
                    slide, rx, block_top, rw, head_h,
                    head_text, font, head_pt,
                    ink, bold=True,
                )
            if body:
                add_textbox(
                    slide, rx, block_top + head_h, rw, body_h,
                    body, font, body_pt_local,
                    ink_soft,
                )
        else:
            desc_text = _dget(level, "desc", "description", "tag", default="")
            if desc_text:
                add_textbox(
                    slide, rx, y + 0.1, rw, layer_h - 0.1,
                    f"• {desc_text}", font, small_pt,
                    ink_soft,
                )


@register_renderer("heatmap-matrix")
def render_heatmap_matrix(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """热力图矩阵：N+1 行 × (cols+2) 列表格。

    Unicode 方案 A：score_chars = 5 个 U+2588 Full Block，按 score 截取前 N 个显示为红色。
    末列为 total 合计。
    """
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "heatmap_matrix")
    row_label_col_px = cfg.get("row_label_column_px", 220)
    total_col_px = cfg.get("total_column_px", 220)
    # Unicode 方案 A：5 个 U+2588 Full Block。fallback 兜底防 layouts 缺字段或误填退格符。
    score_chars = cfg.get("score_chars") or "█████"

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5

    rows_data = _extra(spec, "rows", default=None) or []
    columns = _extra(spec, "columns", default=None) or []
    if not rows_data:
        for p in get_points(spec):
            rows_data.append({
                "label": p.heading or "",
                "scores": [0, 0, 0, 0],
                "total": p.body or "",
            })
    if not columns:
        columns = [f"C{i + 1}" for i in range(len(rows_data[0].get("scores", [])) if rows_data else 4)]

    n_rows = len(rows_data) + 1  # +1 header
    n_cols = len(columns) + 2    # label + N + total

    # 列宽
    row_label_w = _px_in(row_label_col_px)
    total_w = _px_in(total_col_px)
    middle_w = max(area_w - row_label_w - total_w, 1.0)
    mid_col_w = middle_w / max(len(columns), 1)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(area_left), Inches(area_top),
        Inches(area_w), Inches(min(area_h, n_rows * 0.55)),
    )
    tbl = table_shape.table
    # 设列宽
    tbl.columns[0].width = Emu(int(row_label_w * 914400))
    for ci in range(len(columns)):
        tbl.columns[ci + 1].width = Emu(int(mid_col_w * 914400))
    tbl.columns[n_cols - 1].width = Emu(int(total_w * 914400))

    primary = _color(theme, "primary", "#C7000B")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_mute = _color(theme, "ink_mute", "#9A9A9A")
    paper_3 = _color(theme, "paper_3", "#EAEAEA")
    body_pt = _px_pt(_type_px(theme, "body", 24))
    small_pt = _px_pt(_type_px(theme, "small", 20))

    def _set_cell(cell, text: str, color_hex: str, bold: bool = False, fill_hex: str | None = None,
                  size_pt: float = body_pt, align=PP_ALIGN.CENTER) -> None:
        if fill_hex:
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(fill_hex)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        p = tf.paragraphs[0]
        p.alignment = align
        if p.runs:
            run = p.runs[0]
            run.font.name = font
            run.font.size = Pt(size_pt)
            run.font.color.rgb = hex_to_rgb(color_hex)
            run.font.bold = bold

    # Header 行
    _set_cell(tbl.cell(0, 0), "", ink, fill_hex=paper_3, bold=True, size_pt=small_pt)
    for ci, col_name in enumerate(columns):
        _set_cell(tbl.cell(0, ci + 1), str(col_name), ink, bold=True,
                  fill_hex=paper_3, size_pt=small_pt)
    _set_cell(tbl.cell(0, n_cols - 1), "合计", ink, bold=True,
              fill_hex=paper_3, size_pt=small_pt)

    # 数据行
    for ri, row in enumerate(rows_data):
        if isinstance(row, str):
            row = {"label": row, "scores": [0] * len(columns), "total": ""}
        _set_cell(tbl.cell(ri + 1, 0), str(_dget(row, "label", default="")), ink,
                  bold=True, align=PP_ALIGN.LEFT, size_pt=small_pt)
        scores = _dget(row, "scores", default=[]) or []
        for ci in range(len(columns)):
            sc = int(scores[ci]) if ci < len(scores) else 0
            sc = max(0, min(sc, 5))
            filled = score_chars[:sc]
            empty = "░" * (5 - sc)  # U+2591 Light Shade
            cell = tbl.cell(ri + 1, ci + 1)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            if filled:
                r1 = p.runs[0] if p.runs else p.add_run()
                r1.text = filled
                r1.font.name = font
                r1.font.size = Pt(body_pt)
                r1.font.color.rgb = hex_to_rgb(primary)
                r1.font.bold = True
            if empty:
                r2 = p.add_run()
                r2.text = empty
                r2.font.name = font
                r2.font.size = Pt(body_pt)
                r2.font.color.rgb = hex_to_rgb(ink_mute)
        total_text = str(_dget(row, "total", default=""))
        _set_cell(tbl.cell(ri + 1, n_cols - 1), total_text, ink_mute,
                  align=PP_ALIGN.LEFT, size_pt=small_pt)


@register_renderer("thankyou")
def render_thankyou(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """感谢页：红底白字超大 "Thank you." 360px + 4 列联系方式。"""
    sw, sh = _slide_dims_in()
    cfg = _layout(theme, "thankyou")
    bg_token = cfg.get("background_token", "primary")
    title_px = cfg.get("title_size_px", 360)
    contact_cols = cfg.get("contact_columns", 4)

    set_slide_background(slide, _color(theme, bg_token, "#C7000B"))

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)
    paper = _color(theme, "paper", "#FFFFFF")

    # 主标题
    title_text = spec.content.title or "Thank you."
    title_pt = _px_pt(title_px)
    title_h = title_pt / 72 * 1.1 + 0.5
    tbx = slide.shapes.add_textbox(
        Inches(left_pad), Inches(top_pad + 0.8),
        Inches(sw - left_pad - right_pad), Inches(title_h),
    )
    tf = tbx.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = font
    p.font.size = Pt(title_pt)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(paper)

    # 副标题
    if spec.content.subtitle:
        sub_pt = _px_pt(_type_px(theme, "subtitle", 36))
        add_textbox(
            slide, left_pad, top_pad + 0.8 + title_h + 0.2,
            sw - left_pad - right_pad, sub_pt / 72 * 1.5 + 0.3,
            spec.content.subtitle, font, sub_pt,
            paper,
        )

    # 底部联系方式
    contacts = _extra(spec, "contacts", default=None) or []
    if contacts:
        small_pt = _px_pt(_type_px(theme, "small", 20))
        n_cols = min(len(contacts), contact_cols)
        col_w = (sw - left_pad - right_pad) / max(n_cols, 1)
        contact_y = sh - bottom_pad - 0.8
        for i, c in enumerate(contacts[:n_cols]):
            label_v = _dget(c, "label", default="")
            value_v = _dget(c, "value", default="")
            if label_v or value_v:
                text = f"{label_v}\n{value_v}".strip()
            else:
                text = str(c) if c else ""
            add_textbox(
                slide, left_pad + i * col_w, contact_y, col_w, 0.8,
                text, font, small_pt,
                paper,
            )


# ===========================================================================
# P2 四版式（+ V17 变体集合由 design.layout_variant 分派）
# ===========================================================================

@register_renderer("cards-6")
def render_cards_6(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """6 卡（3×2）：直接复用 render.py 的 _render_cards。"""
    _render_cards(slide, spec, theme, safe, 6, total_slides)


@register_renderer("rings")
def render_rings(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """同心环：左侧 3 层同心椭圆 + 右侧步骤列表。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "rings")
    ring_max = cfg.get("ring_count_max", 3)
    right_gap_px = cfg.get("right_list_gap_px", 28)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5

    rings = _extra(spec, "rings", default=None) or []
    steps = _extra(spec, "steps", default=None) or []
    if not rings:
        # 退化：用 points 前 ring_max 条
        for p in get_points(spec)[:ring_max]:
            rings.append({"label": p.heading or p.body})

    n = min(len(rings), ring_max) if rings else ring_max
    left_w = area_w * 0.45
    ring_cx = area_left + left_w / 2
    ring_cy = area_top + area_h / 2

    primary = _color(theme, "primary", "#C7000B")
    primary_dark = _color(theme, "primary_dark", "#8A0008")
    primary_soft = _color(theme, "primary_soft", "#FDECEE")
    paper = _color(theme, "paper", "#FFFFFF")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")

    # 最大直径
    max_d = min(left_w * 0.85, area_h * 0.85)
    step_d = max_d / n

    # 外→内绘制：外层浅 / 空心边框，内层实心
    for i in range(n):
        d = max_d - i * step_d
        x = ring_cx - d / 2
        y = ring_cy - d / 2
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d),
        )
        if i == n - 1:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(primary)
            shape.line.fill.background()
        else:
            shape.fill.background()
            shape.line.color.rgb = hex_to_rgb(primary_dark if i == 0 else primary)
            shape.line.width = Pt(2.5)

        # 环标签
        label = ""
        if i < len(rings):
            item = rings[i]
            label = str(_dget(item, "label", default="")) if isinstance(item, dict) else str(item)
        if label:
            label_w = d * 0.9
            label_x = ring_cx - label_w / 2
            # 外两层顶部标注，内层中心
            if i == n - 1:
                label_y = ring_cy - 0.25
                color = paper
            else:
                label_y = y + 0.1
                color = ink
            add_textbox(
                slide, label_x, label_y, label_w, 0.5,
                label, font, _px_pt(_type_px(theme, "small", 20)),
                color, alignment=PP_ALIGN.CENTER, bold=True,
            )

    # 右侧步骤列表
    if not steps:
        for p in get_points(spec):
            steps.append({"title": p.heading or "", "desc": p.body or ""})
    right_x = area_left + left_w + 0.3
    right_w = area_w - left_w - 0.3
    gap_in = _px_in(right_gap_px)
    body_pt = _px_pt(_type_px(theme, "body", 24))
    small_pt = _px_pt(_type_px(theme, "small", 20))
    cy = area_top + 0.1
    for i, step in enumerate(steps):
        if isinstance(step, str):
            step = {"title": step, "desc": ""}
        title_text = str(_dget(step, "title", default=""))
        desc_text = str(_dget(step, "desc", "description", default=""))

        num_text = f"{i + 1:02d}"
        add_textbox(
            slide, right_x, cy, 0.8, body_pt / 72 * 1.4 + 0.1,
            num_text, font, body_pt,
            primary, bold=True,
        )
        add_textbox(
            slide, right_x + 0.8, cy, right_w - 0.8, body_pt / 72 * 1.4 + 0.1,
            title_text, font, body_pt,
            ink, bold=True,
        )
        if desc_text:
            add_textbox(
                slide, right_x + 0.8,
                cy + body_pt / 72 * 1.4 + 0.05,
                right_w - 0.8,
                small_pt / 72 * 1.6 + 0.2,
                desc_text, font, small_pt,
                ink_soft,
            )
        cy += body_pt / 72 * 1.4 + small_pt / 72 * 1.6 + 0.25 + gap_in
        if cy > area_top + area_h - 0.3:
            break


@register_renderer("personas")
def render_personas(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """人物画像卡：2-4 列 role/name/attrs/quote。quote 左侧 3pt 红 accent bar。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "personas")
    cols_max = cfg.get("columns_max", 4)
    accent_bar_px = cfg.get("quote_accent_bar_px", 3)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5

    personas = _extra(spec, "personas", default=None) or []
    if not personas:
        for p in get_points(spec):
            personas.append({
                "role": p.heading or "",
                "name": p.metric_label or "",
                "attrs": [],
                "quote": p.body or "",
            })
    n = min(len(personas), cols_max) or 1
    gap = 0.25
    col_w = (area_w - gap * (n - 1)) / n

    primary = _color(theme, "primary", "#C7000B")
    paper_2 = _color(theme, "paper_2", "#F4F4F4")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    ink_mute = _color(theme, "ink_mute", "#9A9A9A")
    body_pt = _px_pt(_type_px(theme, "body", 24))
    small_pt = _px_pt(_type_px(theme, "small", 20))
    h4_pt = _px_pt(_type_px(theme, "h4", 28))
    accent_bar_in = _px_in(accent_bar_px)

    for i in range(n):
        item: Any = personas[i] if i < len(personas) else {}
        if isinstance(item, str):
            item = {"role": "", "name": "", "attrs": [], "quote": item}
        x = area_left + i * (col_w + gap)
        cy = area_top

        # role
        role_text = str(_dget(item, "role", default=""))
        if role_text:
            add_textbox(
                slide, x, cy, col_w, small_pt / 72 * 1.4 + 0.1,
                role_text, font, small_pt,
                primary, bold=True,
            )
            cy += small_pt / 72 * 1.4 + 0.05
        # name
        name_text = str(_dget(item, "name", default=""))
        if name_text:
            add_textbox(
                slide, x, cy, col_w, h4_pt / 72 * 1.4 + 0.1,
                name_text, font, h4_pt,
                ink, bold=True,
            )
            cy += h4_pt / 72 * 1.4 + 0.2

        # attrs（label/value 对）
        attrs = _dget(item, "attrs", default=[]) or []
        for attr in attrs[:6]:
            if isinstance(attr, dict):
                label = str(_dget(attr, "label", default=""))
                value = str(_dget(attr, "value", default=""))
                line = f"{label}  {value}"
            else:
                line = str(attr)
            add_textbox(
                slide, x, cy, col_w, small_pt / 72 * 1.5 + 0.1,
                line, font, small_pt,
                ink_soft,
            )
            cy += small_pt / 72 * 1.5 + 0.05

        # quote：左红条 + 灰底 rounded rect
        quote = str(_dget(item, "quote", default=""))
        if quote:
            cy += 0.1
            quote_h = min(area_top + area_h - cy - 0.1, 2.0)
            if quote_h > 0.4:
                # 背景
                add_rounded_rect(slide, x, cy, col_w, quote_h, paper_2, corner_radius=0.06)
                # 红条
                _add_rect(slide, x, cy, accent_bar_in, quote_h, primary)
                add_textbox(
                    slide, x + accent_bar_in + 0.15, cy + 0.15,
                    col_w - accent_bar_in - 0.3, quote_h - 0.3,
                    f"“{quote}”", font, small_pt,
                    ink_soft,
                )


@register_renderer("risk-list")
def render_risk_list(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """风险列表：2 列卡片，每卡 number/title/desc/mitigation + severity 标签。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "risk_list")
    cols_max = cfg.get("columns_max", 2)
    sev_bar_px = cfg.get("severity_bar_width_px", 4)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5

    risks = _extra(spec, "risks", default=None) or []
    if not risks:
        for p in get_points(spec):
            risks.append({
                "title": p.heading or "",
                "desc": p.body or "",
                "severity": "MED",
            })
    # L-6 fix: cols 自适应 ceil(sqrt(n)), card 尺寸按 (canvas - margins - gutters) / (cols, rows) 撑满
    # 旧实现 cols 固定 2 + card_h 上限 2.4", n=4/5 时画布大量空白.
    n = max(len(risks), 1)
    import math as _math
    cols = max(1, min(cols_max, int(_math.ceil(_math.sqrt(n)))))
    # 横向布局优先 (PPT 16:9): 列数 ≥ 行数, 让 cols ≥ rows
    rows = (n + cols - 1) // cols
    if rows > cols and cols < cols_max:
        # 升 cols 到不超过 cols_max, 让 rows 减
        cols = min(cols_max, rows)
        rows = (n + cols - 1) // cols
    gap = 0.25
    card_w = (area_w - gap * (cols - 1)) / cols
    card_h = (area_h - gap * (rows - 1)) / rows  # 不再设 2.4 上限, 数量少时单卡变大填满

    primary = _color(theme, "primary", "#C7000B")
    accent = _color(theme, "accent", "#D97706")
    ok_c = _color(theme, "ok", "#1E7B3A")
    paper_2 = _color(theme, "paper_2", "#F4F4F4")
    paper = _color(theme, "paper", "#FFFFFF")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    ink_mute = _color(theme, "ink_mute", "#9A9A9A")
    body_pt = _px_pt(_type_px(theme, "body", 24))
    small_pt = _px_pt(_type_px(theme, "small", 20))
    h4_pt = _px_pt(_type_px(theme, "h4", 28))
    sev_bar_in = _px_in(sev_bar_px)

    sev_color = {
        "HIGH": primary,
        "H": primary,
        "MED": accent,
        "M": accent,
        "LOW": ok_c,
        "L": ok_c,
    }

    for i, risk in enumerate(risks):
        if isinstance(risk, str):
            risk = {"title": risk, "desc": "", "severity": "MED"}
        r = i // cols
        c = i % cols
        x = area_left + c * (card_w + gap)
        y = area_top + r * (card_h + gap)
        if y + card_h > area_top + area_h:
            break

        # 背景
        add_rounded_rect(slide, x, y, card_w, card_h, paper_2, corner_radius=0.08)
        # 左侧 severity bar
        sev = str(_dget(risk, "severity", default="MED")).upper()
        bar_color = sev_color.get(sev, accent)
        _add_rect(slide, x, y, sev_bar_in, card_h, bar_color)

        # 序号
        num_text = f"R{i + 1:02d}"
        add_textbox(
            slide, x + sev_bar_in + 0.2, y + 0.1, 1.0, small_pt / 72 * 1.4 + 0.1,
            num_text, font, small_pt,
            bar_color, bold=True,
        )
        # severity 标签
        add_textbox(
            slide, x + card_w - 1.2, y + 0.1, 1.0, small_pt / 72 * 1.4 + 0.1,
            sev, font, small_pt,
            bar_color, alignment=PP_ALIGN.RIGHT, bold=True,
        )
        # 标题
        title_text = str(_dget(risk, "title", default=""))
        add_textbox(
            slide, x + sev_bar_in + 0.2, y + 0.5,
            card_w - sev_bar_in - 0.3, h4_pt / 72 * 1.4 + 0.15,
            title_text, font, h4_pt,
            ink, bold=True,
        )
        # 描述
        desc_text = str(_dget(risk, "desc", "description", default=""))
        if desc_text:
            add_textbox(
                slide, x + sev_bar_in + 0.2,
                y + 0.5 + h4_pt / 72 * 1.4 + 0.15,
                card_w - sev_bar_in - 0.3, body_pt / 72 * 1.8 + 0.2,
                desc_text, font, small_pt,
                ink_soft,
            )
        # 缓解措施
        mit_text = str(_dget(risk, "mitigation", default=""))
        if mit_text:
            add_textbox(
                slide, x + sev_bar_in + 0.2,
                y + card_h - small_pt / 72 * 1.5 - 0.3,
                card_w - sev_bar_in - 0.3, small_pt / 72 * 1.5 + 0.1,
                f"→ {mit_text}", font, small_pt,
                bar_color, bold=True,
            )


@register_renderer("governance")
def render_governance(slide, spec: SlideSpec, theme: dict, safe: SafeArea, total_slides: int) -> None:
    """治理结构：顶部决策委员会 rect + 下方 N 个工作组 rect + 连接线。"""
    sw, sh = _slide_dims_in()
    set_slide_background(slide, _color(theme, "paper", "#FFFFFF"))
    render_title_zone(slide, spec, theme, safe)
    render_footer(slide, spec, theme, safe, total_slides)

    cfg = _layout(theme, "governance")
    top_h_px = cfg.get("top_box_height_px", 120)
    unit_cols_max = cfg.get("unit_columns_max", 4)

    top_pad, right_pad, bottom_pad, left_pad = _canvas_padding_in(theme)
    font = _font_family(spec, theme)

    area_left = left_pad
    area_top = top_pad + 1.4
    area_w = sw - area_left - right_pad
    area_h = sh - area_top - bottom_pad - 0.5

    # schemas: spec.variant.top (GovTop {name, tag}); legacy: spec.content.top_box {title, desc}
    top_box = _extra(spec, "top_box", "top", default=None) or {}
    units = _extra(spec, "units", default=None) or []
    if not units:
        for p in get_points(spec):
            units.append({"title": p.heading or "", "desc": p.body or ""})

    primary = _color(theme, "primary", "#C7000B")
    paper = _color(theme, "paper", "#FFFFFF")
    paper_2 = _color(theme, "paper_2", "#F4F4F4")
    ink = _color(theme, "ink", "#1F1F1F")
    ink_soft = _color(theme, "ink_soft", "#4B4B4B")
    rule = _color(theme, "rule", "#D9D9D9")
    body_pt = _px_pt(_type_px(theme, "body", 24))
    small_pt = _px_pt(_type_px(theme, "small", 20))
    h4_pt = _px_pt(_type_px(theme, "h4", 28))

    # 顶部框（决策委员会）
    top_h = _px_in(top_h_px)
    top_w = area_w * 0.6
    top_x = area_left + (area_w - top_w) / 2
    top_y = area_top
    add_rounded_rect(slide, top_x, top_y, top_w, top_h, primary, corner_radius=0.1)

    # schemas: GovTop.{name, tag}; legacy: {title, desc}
    if isinstance(top_box, (dict,)) or hasattr(top_box, "__class__") and not isinstance(top_box, str):
        top_title = _dget(top_box, "title", "name", default="")
        top_desc = _dget(top_box, "desc", "tag", default="")
    else:
        top_title = str(top_box) if top_box else ""
        top_desc = ""
    if top_title:
        add_textbox(
            slide, top_x, top_y + 0.15, top_w, h4_pt / 72 * 1.4 + 0.15,
            top_title, font, h4_pt,
            paper, alignment=PP_ALIGN.CENTER, bold=True,
        )
    if top_desc:
        add_textbox(
            slide, top_x, top_y + 0.15 + h4_pt / 72 * 1.4 + 0.1,
            top_w, small_pt / 72 * 1.5 + 0.15,
            top_desc, font, small_pt,
            paper, alignment=PP_ALIGN.CENTER,
        )

    # 底部 N 个工作组
    n = min(len(units), unit_cols_max) or 1
    gap = 0.25
    unit_w = (area_w - gap * (n - 1)) / n
    unit_top = top_y + top_h + 0.6
    unit_h = min(area_h - top_h - 0.6 - 0.1, 3.0)

    # 连接线：顶框底边中点 → 水平横杆 → 下垂到每个 unit 顶部
    hub_y = top_y + top_h
    trunk_y = hub_y + 0.25
    hub_x_center = top_x + top_w / 2
    # 垂直干线
    slide.shapes.add_connector(
        1, Inches(hub_x_center), Inches(hub_y),
        Inches(hub_x_center), Inches(trunk_y),
    ).line.color.rgb = hex_to_rgb(rule)
    # 水平横杆（覆盖所有 unit 中心）
    if n > 1:
        first_cx = area_left + unit_w / 2
        last_cx = area_left + (n - 1) * (unit_w + gap) + unit_w / 2
        slide.shapes.add_connector(
            1, Inches(first_cx), Inches(trunk_y),
            Inches(last_cx), Inches(trunk_y),
        ).line.color.rgb = hex_to_rgb(rule)

    for i in range(n):
        item = units[i] if i < len(units) else {}
        if isinstance(item, str):
            item = {"title": item, "desc": ""}
        ux = area_left + i * (unit_w + gap)
        # 垂直支线
        cx = ux + unit_w / 2
        slide.shapes.add_connector(
            1, Inches(cx), Inches(trunk_y),
            Inches(cx), Inches(unit_top),
        ).line.color.rgb = hex_to_rgb(rule)

        add_rounded_rect(slide, ux, unit_top, unit_w, unit_h, paper_2, corner_radius=0.08)
        # schemas: GovUnit.{code, name, items[]}; legacy: {title, desc, description}
        ut = _dget(item, "title", "name", default="")
        ud = _dget(item, "desc", "description", default="")
        if not ud:
            items_list = _dget(item, "items", default=None)
            if isinstance(items_list, (list, tuple)) and items_list:
                ud = "\n".join(f"• {it}" for it in items_list)
        # 同时显示 code (如 "L1") 作为 ut 前缀, 若 schemas 给了 code
        unit_code = _dget(item, "code", default="")
        if unit_code and unit_code not in ut:
            ut = f"{unit_code} {ut}".strip()
        add_textbox(
            slide, ux + 0.15, unit_top + 0.15, unit_w - 0.3,
            body_pt / 72 * 1.4 + 0.15,
            ut, font, body_pt,
            ink, bold=True,
        )
        if ud:
            add_textbox(
                slide, ux + 0.15,
                unit_top + 0.15 + body_pt / 72 * 1.4 + 0.1,
                unit_w - 0.3,
                unit_h - 0.4 - body_pt / 72 * 1.4,
                ud, font, small_pt,
                ink_soft,
            )
