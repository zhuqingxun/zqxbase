# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=1.0.0"]
# ///
"""density_check.py: 渲染后密度预警 — 分析 PPTX 每页 shape 面积占比。

仅预警，不阻断管线。shape 数量 <= EXEMPT_MAX_SHAPES (默认 2) 的标题型页面豁免。

计算方式: 所有 shape 的包围盒面积**累加和** / slide 面积. 不做重叠去重 — 因为 huawei renderer 大量
用"卡片背景 + 文本框叠加"模式 (cards-N 中 textbox 与 add_rounded_rect 重叠), 累加和往往 > 100%.
**这是设计限制不是 bug** (2026-05-19 audit 更新 docstring 与实现一致): 真正去重叠需要 shapely
依赖, 项目阶段不引入. 该预警实际只对极简页面 (shape 数 ≤ 2 已被豁免, 故几乎不触发) 有意义,
留作"未来引入 shapely 时再激活" 的占位机制.

阈值: 内容面积 < 20% 时 WARN.

Usage:
    uv run --script engine/density_check.py <output.pptx>
    uv run --script engine/density_check.py <output.pptx> --json
    uv run --script engine/density_check.py <output.pptx> --threshold 25
"""

from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Emu


# ---------------------------------------------------------------------------
# 配置
# ---------------------------------------------------------------------------

DEFAULT_THRESHOLD_PCT = 20  # 内容面积 < 此值时 WARN

# 豁免：slide 标题含这些关键词时跳过（匹配 visual type 无法从 PPTX 反推，用标题启发式）
EXEMPT_KEYWORDS: set[str] = set()  # 不用关键词，用 shape 数量启发式

# shape 数量 <= 此值认为是标题型页面，豁免
EXEMPT_MAX_SHAPES = 2


# ---------------------------------------------------------------------------
# 数据结构
# ---------------------------------------------------------------------------

@dataclass
class SlideMetrics:
    """单页密度指标。"""
    index: int
    title: str
    shape_count: int
    content_area_pct: float
    status: str  # "OK" | "WARN" | "SKIP"

    def to_dict(self) -> dict:
        return {
            "index": self.index,
            "title": self.title,
            "shape_count": self.shape_count,
            "content_area_pct": round(self.content_area_pct, 1),
            "status": self.status,
        }


@dataclass
class DensityReport:
    """整份 PPTX 的密度报告。"""
    pptx_path: str
    threshold_pct: float
    slides: list[SlideMetrics] = field(default_factory=list)

    @property
    def warned(self) -> int:
        return sum(1 for s in self.slides if s.status == "WARN")

    @property
    def ok(self) -> int:
        return sum(1 for s in self.slides if s.status == "OK")

    @property
    def skipped(self) -> int:
        return sum(1 for s in self.slides if s.status == "SKIP")

    def to_dict(self) -> dict:
        return {
            "pptx_path": self.pptx_path,
            "threshold_pct": self.threshold_pct,
            "total": len(self.slides),
            "ok": self.ok,
            "warned": self.warned,
            "skipped": self.skipped,
            "slides": [s.to_dict() for s in self.slides],
        }


# ---------------------------------------------------------------------------
# 分析逻辑
# ---------------------------------------------------------------------------

def _analyze_slide_shapes(
    slide, slide_width_emu: int, slide_height_emu: int
) -> tuple[str, int, float]:
    """单次遍历 slide.shapes 同时取 title + shape_count + 累计面积占比.

    2026-05-19 audit fix: 原版 _get_slide_title 和 _compute_content_area_pct 各遍历一次
    (50 页 × 25 shapes = 5000 → 10000 shape 访问), 合并为单次.
    """
    title = "(no title)"
    title_taken = False
    total_shape_area = 0
    shape_count = 0
    for shape in slide.shapes:
        shape_count += 1
        if not title_taken and shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                title = txt[:60]
                title_taken = True
        w = shape.width if shape.width else 0
        h = shape.height if shape.height else 0
        total_shape_area += w * h
    slide_area = slide_width_emu * slide_height_emu
    pct = (total_shape_area / slide_area * 100) if slide_area else 0.0
    return title, shape_count, pct


def analyze_density(pptx_path: str, threshold_pct: float = DEFAULT_THRESHOLD_PCT) -> DensityReport:
    """分析 PPTX 每页密度。"""
    prs = Presentation(pptx_path)
    report = DensityReport(pptx_path=pptx_path, threshold_pct=threshold_pct)

    for i, slide in enumerate(prs.slides):
        title, shape_count, pct = _analyze_slide_shapes(
            slide, prs.slide_width, prs.slide_height
        )

        # 标题型页面豁免（shape 数量 <= 2）
        if shape_count <= EXEMPT_MAX_SHAPES:
            report.slides.append(SlideMetrics(
                index=i + 1,
                title=title,
                shape_count=shape_count,
                content_area_pct=0.0,
                status="SKIP",
            ))
            continue

        status = "WARN" if pct < threshold_pct else "OK"
        report.slides.append(SlideMetrics(
            index=i + 1,
            title=title,
            shape_count=shape_count,
            content_area_pct=pct,
            status=status,
        ))

    return report


# ---------------------------------------------------------------------------
# 输出格式化
# ---------------------------------------------------------------------------

def format_text_report(report: DensityReport) -> str:
    """生成人类可读的密度报告。"""
    lines: list[str] = []
    lines.append(f"Density check: {report.pptx_path} (threshold: {report.threshold_pct}%)\n")

    for s in report.slides:
        tag = s.status
        if s.status == "SKIP":
            lines.append(f"  SKIP  Slide {s.index:2d}: \"{s.title}\" ({s.shape_count} shapes, exempt)")
        elif s.status == "WARN":
            lines.append(f"  WARN  Slide {s.index:2d}: \"{s.title}\" — {s.content_area_pct:.1f}% content")
        else:
            lines.append(f"  OK    Slide {s.index:2d}: \"{s.title}\" — {s.content_area_pct:.1f}% content")

    lines.append(f"\nSummary: {report.ok} OK / {report.warned} WARN / {report.skipped} SKIP")
    if report.warned > 0:
        lines.append(f"NOTE: {report.warned} slides below {report.threshold_pct}% content density.")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> None:
    parser = argparse.ArgumentParser(description="PPTX 密度预警")
    parser.add_argument("pptx", help=".pptx 文件路径")
    parser.add_argument("--threshold", type=float, default=DEFAULT_THRESHOLD_PCT,
                        help=f"内容面积阈值百分比 (default: {DEFAULT_THRESHOLD_PCT})")
    parser.add_argument("--json", action="store_true", help="输出 JSON 格式")
    args = parser.parse_args()

    report = analyze_density(args.pptx, args.threshold)

    if args.json:
        print(json.dumps(report.to_dict(), ensure_ascii=False, indent=2))
    else:
        print(format_text_report(report))

    # 预警脚本不影响退出码，始终 exit 0
    sys.exit(0)


if __name__ == "__main__":
    main()
